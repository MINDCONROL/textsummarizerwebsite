# Import necessary libraries
import streamlit as st  # For creating the web application
import os  # For operating system dependent functionalities (like file paths)
import tempfile  # For creating temporary files and directories
import random  # For generating random numbers (used in analogies, study questions)
from typing import Optional, List, Dict, Tuple, Any  # For type hinting, improving code readability

# NLP and Machine Learning Libraries
from transformers import T5Tokenizer, T5ForConditionalGeneration, pipeline  # For T5 summarization and specialized pipelines
import nltk  # Natural Language Toolkit for text processing
# nltk.download('punkt') # Downloaded via function call later
# nltk.download('stopwords') # Downloaded via function call later
import spacy  # For advanced NLP tasks like NER, POS tagging

# Text Processing and Data Handling
import re  # For regular expression operations
from collections import defaultdict  # For dictionaries with default values (not explicitly used anymore but good to know)
from sklearn.feature_extraction.text import TfidfVectorizer  # For TF-IDF calculation
import numpy as np  # For numerical operations (used with TF-IDF)

# File Extraction Libraries
from PyPDF2 import PdfReader  # For reading text from PDF files
from docx import Document as DocxDocument  # For reading text from DOCX files (aliased to avoid conflict)
from pptx import Presentation  # For reading text from PPTX files

# Visualization Libraries
import networkx as nx  # For creating and manipulating complex networks
import matplotlib.pyplot as plt  # For creating static, animated, and interactive visualizations
from wordcloud import WordCloud  # For generating word cloud images
import plotly.graph_objects as go  # For creating interactive Plotly graphs

# Web Content and Audio/Video Processing
import requests  # For making HTTP requests to fetch web content
from bs4 import BeautifulSoup  # For parsing HTML and XML documents
from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled  # For fetching YouTube transcripts
from moviepy.editor import VideoFileClip # For processing video and audio files
import whisper # For audio transcription using OpenAI's Whisper model

# Concurrency
import concurrent.futures  # For running operations in parallel (like multiple summaries)

# --- Global Variables & Configuration ---
# Attempt to download NLTK data, handle potential errors gracefully
try:
    nltk.download('punkt', quiet=True)  # Tokenizer models
    nltk.download('stopwords', quiet=True)  # Stopwords list
except Exception as e:
    st.warning(f"Could not download NLTK data automatically: {e}. Some features might be limited.")

# --- Streamlit Page Configuration ---
# Set up the page layout, title, and icon for the Streamlit app
# This should be the first Streamlit command in your script, and it can only be set once.
st.set_page_config(
    layout="wide",  # Use wide layout for more space
    page_title="AI Text & Media Analyzer Suite",  # Title appearing in the browser tab
    page_icon="✨"  # Icon for the browser tab (can be an emoji or a URL)
)

# --- Custom CSS Styling ---
# Apply custom CSS for a more appealing and consistent user interface
st.markdown("""
<style>
    /* Main app background */
    .stApp {
        color: black; /* Default text color to black for overall readability */
    }

    /* Sidebar styling */
    .css-1d391kg { /* Sidebar main class */
        background-color: #004d40; /* Dark teal background for sidebar */
        padding: 10px;
        color: white; /* Ensure text in the sidebar is white for contrast */
    }
    .css-1d391kg .stSelectbox label { /* Sidebar selectbox label */
        color: white; /* Ensure label text is white */
    }
    .st-emotion-cache-17x3wpl { /* Sidebar elements general */
        color: white; /* Ensure other text elements are white */
    }


    /* Main content area titles */
    h1, h2, h3 {
        color: white; /* Set header text to white */
    }

    /* Buttons styling */
    .stButton>button {
        background-color: #00796b; /* Teal background */
        color: white; /* White text */
        border-radius: 5px;
        padding: 10px 15px;
        border: none;
    }
    .stButton>button:hover {
        background-color: #004d40; /* Darker teal on hover */
        color: white; /* Maintain white text on hover */
    }

    /* File uploader styling */
    .stFileUploader label {
        color: black; /* Set file uploader label to black */
    }
    .stFileUploader div[data-baseweb="input"] > div { /* Style the input text color */
        color: black !important;
    }

    /* Text input/area styling */
    .stTextInput input, .stTextArea textarea {
        border: 1px solid #00796b; /* Teal border */
        border-radius: 5px;
        color: black; /* Set input text color to black */
    }
    .stTextInput label, .stTextArea label {
        color: black; /* Set label text color to black */
    }

    /* Summary boxes and content cards */
    .summary-box, .content-card {
        border-left: 5px solid #00796b; /* Teal left border */
        padding: 1rem;
        margin: 1rem 0;
        background-color: #e0f2f1; /* Light teal background */
        border-radius: 5px;
        box-shadow: 2px 2px 5px rgba(0,0,0,0.1);
        color: black; /* Set text color within these boxes to black */
    }

    .child-mode {
        background-color: #fff8e1; /* Light yellow for child mode */
        padding: 1rem;
        border-radius: 10px;
        border-left: 5px solid #ffc107; /* Amber border */
        color: black; /* Set text color in child mode to black */
    }

    .flashcard {
        background-color: #f1f8e9; /* Light green background */
        border-radius: 8px;
        padding: 15px;
        margin: 10px 0;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        border-left: 4px solid #8bc34a; /* Green border */
        color: black; /* Set text color in flashcards to black */
    }

    /* Tabs styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 24px; /* Space between tabs */
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        white-space: pre-wrap;
        background-color: #e0f2f1; /* Light teal for inactive tabs */
        border-radius: 4px 4px 0px 0px;
        padding: 10px;
        color: black; /* Set text color for inactive tabs to black */
    }
    .stTabs [aria-selected="true"] {
        background-color: #00796b; /* Teal for active tab */
        color: white; /* Set text color for active tab to white for contrast */
    }

    /* Metrics styling */
    .stMetric {
        background-color: #e8eaf6; /* Indigo light background */
        padding: 10px;
        border-radius: 5px;
        border: 1px solid #3f51b5; /* Indigo border */
        color: black; /* Set text color in metrics to black */
    }
    .stMetricLabel {
        color: black; /* Ensure metric labels are black */
    }
    .stMetricValue {
        color: black; /* Ensure metric values are black */
    }
</style>
""", unsafe_allow_html=True) # Allow HTML for styling purposes

# --- CORE MODEL LOADING ---
@st.cache_resource # Cache the loaded models to avoid reloading on each interaction
def load_models() -> Dict[str, Any]:
    """Load and cache all required NLP models."""
    try:
        # Main T5 model for general summarization
        t5_tokenizer = T5Tokenizer.from_pretrained("t5-small")
        t5_model = T5ForConditionalGeneration.from_pretrained("t5-small")

        # Specialized summarizer for child-friendly summaries (BART is good for this)
        child_summarizer = pipeline(
            "summarization",
            model="facebook/bart-large-cnn", # BART model fine-tuned on CNN/DailyMail
            tokenizer="facebook/bart-large-cnn"
        )

        # spaCy NLP processor for tasks like NER, POS tagging, noun chunks
        nlp = spacy.load("en_core_web_sm") # Small English model

        # Whisper model for audio transcription
        # It will download the model on first use if not already present
        # Using "base" model for a balance of speed and accuracy in a web app context
        # Other options: "tiny", "small", "medium", "large"
        transcription_model = whisper.load_model("base")

        return {
            't5_tokenizer': t5_tokenizer,
            't5_model': t5_model,
            'child_summarizer': child_summarizer,
            'nlp': nlp,
            'transcription_model': transcription_model
        }
    except Exception as e: # Catch any exception during model loading
        st.error(f"Fatal Error: Model loading failed: {str(e)}. The application cannot continue.")
        st.stop() # Stop the Streamlit application if models can't load

# --- FILE AND CONTENT PROCESSING UTILITIES ---
def extract_text_from_file(file_path: str, file_type: str) -> Optional[str]:
    """Universal text extraction for multiple file formats."""
    try:
        # PDF file processing
        if file_type == "application/pdf":
            reader = PdfReader(file_path) # Initialize PDF reader
            # Join text from all pages, checking if text extraction yields content
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

        # DOCX file processing
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = DocxDocument(file_path) # Initialize DOCX document reader
            # Join text from all paragraphs, checking if paragraph text exists
            return "\n".join(para.text for para in doc.paragraphs if para.text)

        # PPTX file processing
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file_path) # Initialize PPTX presentation reader
            text_runs = [] # List to hold text from shapes
            for slide in prs.slides: # Iterate through each slide
                for shape in slide.shapes: # Iterate through each shape on the slide
                    if hasattr(shape, "text") and shape.text.strip(): # Check if shape has text and it's not empty
                        text_runs.append(shape.text)
            return "\n\n".join(text_runs) # Join text from shapes with double newline

        # TXT file processing
        elif file_type == "text/plain":
            with open(file_path, 'r', encoding='utf-8') as f: # Open in read mode with UTF-8 encoding
                return f.read() # Read the entire file content

        return None # Return None if file type is not supported or no text found
    except Exception as e: # Catch any exception during file extraction
        st.error(f"File Extraction Error: {str(e)}")
        return None

def preprocess_text(text: str, mode: str = "default") -> str:
    """Text cleaning with mode-specific rules."""
    if not text: return "" # Return empty string if input text is None or empty
    text = re.sub(r'\s+', ' ', text).strip() # Replace multiple whitespaces with single, strip leading/trailing

    # Mode-specific preprocessing
    if mode == "child":
        # Keep simple punctuation, remove complex characters for child-friendliness
        text = re.sub(r'[^\w\s.,!?\'"-]', '', text) # Allow alphanumeric, whitespace, and basic punctuation
        text = re.sub(r'\b\w{15,}\b', '', text)  # Remove very long words (often jargon)
    else: # Default mode for other audiences
        text = re.sub(r'\[.*?\]|\(.*?\)', '', text)  # Remove bracketed/parenthesized citations or notes

    return text # Return the preprocessed text

def transcribe_audio_with_whisper(audio_file_path: str, model: Any) -> Optional[str]:
    """Transcribes an audio file using OpenAI's Whisper model."""
    try:
        with st.spinner("Transcribing audio... This may take a few moments."):
            result = model.transcribe(audio_file_path) # Perform transcription
        return result["text"] # Return the transcribed text
    except Exception as e: # Catch transcription errors
        st.error(f"Audio Transcription Error: {str(e)}")
        return None

def extract_audio_from_video(video_file_path: str, audio_output_path: str) -> bool:
    """Extracts audio from a video file and saves it."""
    try:
        with st.spinner("Extracting audio from video..."):
            video_clip = VideoFileClip(video_file_path) # Load video clip
            audio_clip = video_clip.audio # Get audio track
            if audio_clip:
                audio_clip.write_audiofile(audio_output_path, codec='pcm_s16le') # Write audio to file
                audio_clip.close() # Close audio clip
                video_clip.close() # Close video clip
                return True # Indicate success
            else:
                st.warning("No audio track found in the video.")
                video_clip.close()
                return False # Indicate no audio track
    except Exception as e: # Catch audio extraction errors
        st.error(f"Audio Extraction Error: {str(e)}")
        return False

def fetch_and_parse_url(url: str) -> Optional[str]:
    """Fetches content from a URL (webpage or YouTube) and extracts text."""
    try:
        with st.spinner(f"Fetching content from {url}..."):
            # Check if it's a YouTube URL
            if "youtube.com/watch?v=" in url or "youtu.be/" in url:
                video_id = None
                if "watch?v=" in url:
                    video_id = url.split("watch?v=")[1].split("&")[0]
                elif "youtu.be/" in url:
                    video_id = url.split("youtu.be/")[1].split("?")[0]

                if video_id:
                    try:
                        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
                        # Try to get a manually created transcript first, then an auto-generated one
                        transcript = transcript_list.find_manually_created_transcript(['en']) \
                                     or transcript_list.find_generated_transcript(['en'])
                        transcript_text = " ".join([item['text'] for item in transcript.fetch()])
                        return transcript_text
                    except NoTranscriptFound:
                        st.error(f"No English transcript found for YouTube video: {url}")
                        return None
                    except TranscriptsDisabled:
                        st.error(f"Transcripts are disabled for YouTube video: {url}")
                        return None
                else:
                    st.error("Could not parse YouTube video ID.")
                    return None
            else: # Regular webpage
                headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
                response = requests.get(url, timeout=10, headers=headers) # Make HTTP GET request
                response.raise_for_status() # Raise an exception for bad status codes (4XX or 5XX)
                soup = BeautifulSoup(response.content, 'html.parser') # Parse HTML content

                # Remove script and style elements
                for script_or_style in soup(["script", "style"]):
                    script_or_style.decompose() # Remove the element from the tree

                # Get text, join paragraphs, and clean up
                text_elements = soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'article'])
                text = "\n".join(element.get_text(separator=" ", strip=True) for element in text_elements)
                return text if text else None
    except requests.exceptions.RequestException as e: # Catch network/HTTP errors
        st.error(f"URL Fetch Error: {str(e)}")
        return None
    except Exception as e: # Catch other parsing errors
        st.error(f"Content Parsing Error: {str(e)}")
        return None

# --- SUMMARY GENERATION ---
def generate_summary(models: dict, text: str, audience: str, max_input_length: int = 1024) -> str:
    """Generate audience-specific summaries using appropriate models."""
    if not text: return "Cannot generate summary from empty text." # Handle empty input
    try:
        # Child-friendly summary using BART model
        if audience == "child":
            # Truncate text for child summarizer if it's too long (BART has token limits)
            processed_text_for_child = preprocess_text(text, "child")[:max_input_length * 3] # Give it a bit more context than T5
            if not processed_text_for_child.strip(): return "Not enough content for a child-friendly summary after preprocessing."

            summary_output = models['child_summarizer'](
                processed_text_for_child,
                max_length=150,  # Shorter summary for children
                min_length=30,
                do_sample=True,  # Use sampling for more creative output
                temperature=0.8, # Control randomness
                top_p=0.95       # Nucleus sampling
            )
            summary = summary_output[0]['summary_text']

            # Simple analogies to make it more relatable
            analogies = [
                "like building with colorful blocks",
                "similar to how a tiny seed grows into a big plant",
                "just like sharing your favorite toys with a friend",
                "comparable to all the different colors in a rainbow",
                "like different animals in a zoo, each special in its own way"
            ]
            # Format the child summary with an analogy
            return f"Imagine this story is about {summary.lower()}. It's a bit {random.choice(analogies)}."

        # Summaries for other audiences (student, researcher, expert) using T5 model
        else:
            # Define prompts for T5 based on the target audience
            prompts = {
                "student": "summarize for a high school student, focusing on main ideas and concepts: ",
                "researcher": "summarize for a researcher, highlighting key findings, methodology, and implications: ",
                "expert": "provide a detailed technical summary for an expert in the field, including nuances: "
            }
            # Process text in chunks if it's too long for T5's input limit
            # T5 typically has a 512 token limit, but we chunk by character length for simplicity.
            # Average token is ~4 chars. Chunk size 3000 chars ~ 750 tokens.
            # Model will truncate to its max_length (e.g., 512 tokens for t5-small) during encoding.
            chunk_size = 3000 # Characters per chunk
            chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
            summaries = [] # List to store summaries of each chunk

            # Limit to processing a few chunks to manage time and resources
            num_chunks_to_process = min(len(chunks), 3) # Process up to 3 chunks

            for chunk_idx in range(num_chunks_to_process):
                chunk = chunks[chunk_idx]
                # Encode the prompt + chunk for T5
                inputs = models['t5_tokenizer'].encode(
                    prompts[audience] + chunk,
                    return_tensors="pt", # Return PyTorch tensors
                    max_length=max_input_length, # Max tokens for the model input
                    truncation=True # Truncate if longer than max_length
                )

                # Configure generation parameters based on audience
                gen_kwargs = {
                    "max_length": 250 if audience == "student" else 400, # Max length of the generated summary
                    "min_length": 75 if audience == "student" else 150,  # Min length of the summary
                    "num_beams": 4,       # Beam search for better quality
                    "early_stopping": True, # Stop when beams converge
                    "temperature": 0.7 if audience == "student" else 0.6 # Control randomness
                }
                # Specific parameters for expert summaries
                if audience == "expert":
                    gen_kwargs.update({
                        "length_penalty": 2.0, # Encourage longer, more detailed summaries
                        "no_repeat_ngram_size": 3 # Avoid repeating 3-grams
                    })

                # Generate summary using T5 model
                outputs = models['t5_model'].generate(inputs, **gen_kwargs)
                # Decode the generated summary
                summaries.append(models['t5_tokenizer'].decode(outputs[0], skip_special_tokens=True))

            full_summary = "\n\n[CHUNK BREAK]\n\n".join(summaries) # Join chunk summaries

            # Add key findings (top TF-IDF terms) for researcher summaries
            if audience == "researcher" and text:
                try:
                    vectorizer = TfidfVectorizer(stop_words='english', max_features=1000)
                    tfidf = vectorizer.fit_transform([text]) # Calculate TF-IDF
                    features = vectorizer.get_feature_names_out() # Get feature names (terms)
                    # Get top 5 terms based on TF-IDF scores
                    top_terms_indices = np.argsort(tfidf.toarray().sum(axis=0))[-5:][::-1]
                    top_terms = [features[i] for i in top_terms_indices]
                    full_summary += "\n\n**KEY TERMS IDENTIFIED:**\n- " + "\n- ".join(top_terms)
                except Exception as tfidf_e:
                    st.warning(f"Could not extract key terms for researcher summary: {tfidf_e}")

            return full_summary # Return the complete summary

    except Exception as e: # Catch any summary generation error
        st.error(f"Summary Generation Failed: {str(e)}")
        return "Could not generate summary due to an error."

# --- STUDY MATERIAL GENERATION ---
def generate_study_materials(text: str, nlp_model: Any) -> Optional[Dict[str, List[str]]]:
    """Create comprehensive study materials from text using spaCy."""
    if not text or not text.strip(): return None # Return None if text is empty

    doc = nlp_model(text[:50000]) # Process with spaCy (limit length for performance)
    sentences = list(doc.sents) # Get sentences
    if not sentences: return None # Return None if no sentences found

    # Extract key phrases (noun chunks, prioritizing those with adjectives/numbers)
    key_phrases = []
    for chunk in doc.noun_chunks:
        # Prioritize phrases that are descriptive or quantitative
        if any(token.pos_ in ('ADJ', 'NUM') for token in chunk) or len(chunk.text.split()) > 1:
             # Basic filter for too short or too common phrases (customize as needed)
            if len(chunk.text) > 3 and chunk.text.lower() not in ["this study", "the paper", "the author"]:
                key_phrases.append(chunk.text.strip())
    key_phrases = list(dict.fromkeys(key_phrases)) # Remove duplicates while preserving order
    key_phrases = [phrase for phrase in key_phrases if len(phrase.split()) <= 5][:15] # Limit length and count

    if not key_phrases: return {"flashcards": ["No key phrases found to generate flashcards."], "questions": [], "diagrams": [], "plan": [], "key_phrases": []}


    # Generate Flashcards with context
    flashcards = []
    for phrase in key_phrases[:10]: # Limit to 10 flashcards
        context_sentences = [s.text for s in sentences if phrase.lower() in s.text.lower()]
        context = random.choice(context_sentences) if context_sentences else f"An important concept related to {phrase}."
        # Ensure context is not overly long for a flashcard
        context = (context[:200] + '...') if len(context) > 200 else context
        flashcards.append(f"**Q: What is '{phrase}'?**\nA: *Context:* {context}")

    # Generate Practice Questions
    questions = []
    if len(key_phrases) > 1:
        for i in range(min(5, len(key_phrases) -1 )): # Limit to 5 questions
            # Create varied question types
            q_type = random.choice(["relate", "significance", "example"])
            if q_type == "relate" and i + 1 < len(key_phrases):
                questions.append(
                    f"**Q: How might '{key_phrases[i]}' be related to '{key_phrases[i+1]}'?**\n"
                    f"A: *Hint:* Consider their roles, definitions, or how one might influence the other."
                )
            elif q_type == "significance":
                questions.append(
                    f"**Q: What is the potential significance of '{key_phrases[i]}'?**\n"
                    f"A: *Hint:* Think about its impact or importance in the broader context."
                )
            else: # example
                 questions.append(
                    f"**Q: Can you think of an example or application of '{key_phrases[i]}'?**\n"
                    f"A: *Hint:* Try to connect it to a real-world scenario or a concept you already know."
                )
    else:
        questions.append(f"**Q: Elaborate on '{key_phrases[0]}'.**\nA: *Hint:* Provide more details or explain its meaning.")


    # Suggest Diagrams based on key phrases
    diagrams = []
    if len(key_phrases) >= 3:
        diagrams.append(f"**Concept Map:** Visually connect '{key_phrases[0]}', '{key_phrases[1]}', and '{key_phrases[2]}' to show their relationships.")
    if len(key_phrases) >= 1:
        diagrams.append(f"**Timeline (if applicable):** If '{key_phrases[0]}' involves a process or history, create a timeline.")
    if len(key_phrases) >= 2:
        diagrams.append(f"**Venn Diagram:** Compare and contrast '{key_phrases[0]}' and '{key_phrases[1]}'.")
    diagrams.append("**Flowchart:** If the text describes a process, try to map it out.")


    # Generate a simple Study Plan
    study_plan = []
    for i, phrase in enumerate(key_phrases[:min(3, len(key_phrases))]): # Plan for first 3 key phrases
        study_plan.append(f"**Day {i+1}:** Focus on understanding '{phrase}'. Review flashcards and related text sections.")
    if len(key_phrases) > 3:
        study_plan.append(f"**Day {min(4, len(key_phrases))}:** Review all key phrases and attempt practice questions.")

    return {
        "flashcards": flashcards,
        "questions": questions,
        "diagrams": diagrams,
        "plan": study_plan,
        "key_phrases": key_phrases
    }

# --- VISUALIZATION UTILITIES ---
def create_wordcloud(text: str) -> Optional[plt.Figure]:
    """Generate a word cloud visualization."""
    if not text or not text.strip(): return None # Return None if no text
    try:
        wordcloud_generator = WordCloud(
            width=800,
            height=400,
            background_color='white',
            colormap='viridis', # A visually appealing colormap
            stopwords=nltk.corpus.stopwords.words('english'), # Use NLTK's English stopwords
            collocations=False # Avoid showing collocations (pairs of words)
        ).generate(text)

        fig, ax = plt.subplots(figsize=(12, 6)) # Create matplotlib figure and axes
        ax.imshow(wordcloud_generator, interpolation='bilinear') # Display word cloud
        ax.axis("off") # Hide axes
        plt.tight_layout(pad=0) # Adjust layout
        return fig # Return the figure object
    except Exception as e:
        st.warning(f"Word cloud generation failed: {e}")
        return None

def create_entity_network(entities: List[Tuple[str, str]], nlp_doc: Any) -> Optional[go.Figure]:
    """Create an interactive entity relationship network using Plotly and NetworkX."""
    if not entities: return None # Return None if no entities

    G = nx.Graph() # Initialize a NetworkX graph

    # Add nodes with type grouping and size based on entity length
    unique_entities = list(dict.fromkeys(entities))[:50] # Limit for performance, ensure uniqueness
    for entity_text, label in unique_entities:
        G.add_node(entity_text, type=label, size=8 + len(entity_text) / 2.5, label=label)

    # Add edges between entities that co-occur in the same sentence
    # Create a dictionary for faster lookup of entity types
    entity_map = {text: label for text, label in unique_entities}
    sentences = list(nlp_doc.sents)

    for sentence in sentences:
        sent_entities = [ent.text for ent in sentence.ents if ent.text in entity_map]
        # Add edges between co-occurring entities in the sentence
        for i in range(len(sent_entities)):
            for j in range(i + 1, len(sent_entities)):
                if G.has_node(sent_entities[i]) and G.has_node(sent_entities[j]):
                    if G.has_edge(sent_entities[i], sent_entities[j]):
                        G[sent_entities[i]][sent_entities[j]]['weight'] += 0.2 # Increase weight for co-occurrence
                    else:
                        G.add_edge(sent_entities[i], sent_entities[j], weight=0.5)


    if not G.nodes: return None # If no nodes were added (e.g. due to filtering)

    pos = nx.spring_layout(G, k=0.8, iterations=50) # Position nodes using spring layout algorithm

    # Create Plotly figure components for edges
    edge_x, edge_y = [], []
    for edge in G.edges(data=True):
        x0, y0 = pos[edge[0]]
        x1, y1 = pos[edge[1]]
        edge_x.extend([x0, x1, None]) # X coordinates for edges
        edge_y.extend([y0, y1, None]) # Y coordinates for edges

    edge_trace = go.Scatter(
        x=edge_x, y=edge_y,
        line=dict(width=max(0.5, G[edge[0]][edge[1]]['weight'] if 'weight' in G[edge[0]][edge[1]] else 0.5), color='#888'), # Edge line style
        hoverinfo='none', # No hover info for edges
        mode='lines' # Draw lines for edges
    )

    # Create Plotly figure components for nodes
    node_x, node_y, node_text, node_color, node_size_viz = [], [], [], [], []
    entity_types = list(set(nx.get_node_attributes(G, 'type').values()))
    color_map_types = plt.cm.get_cmap('Accent', len(entity_types) if entity_types else 1) # Using 'Accent'

    for node_idx, node in enumerate(G.nodes()):
        x, y = pos[node]
        node_x.append(x) # X coordinate for node
        node_y.append(y) # Y coordinate for node
        node_info = G.nodes[node]
        type_info = node_info.get('type', 'Unknown')
        node_text.append(f"{node}<br>Type: {type_info}") # Text for hover info
        # Assign color based on entity type
        try:
            type_index = entity_types.index(type_info)
            rgb_color = color_map_types(type_index)
            node_color.append(f'rgb({int(rgb_color[0]*255)},{int(rgb_color[1]*255)},{int(rgb_color[2]*255)})')
        except ValueError: # Fallback color
             node_color.append('rgb(128,128,128)') # Grey for unknown types

        node_size_viz.append(node_info.get('size', 10)) # Visual size of node

    node_trace = go.Scatter(
        x=node_x, y=node_y,
        mode='markers+text', # Show markers and text for nodes
        text=[node for node in G.nodes()], # Display node name as text
        textposition='top center',
        hoverinfo='text', # Show text from node_text on hover
        hovertext=node_text,
        marker=dict(
            showscale=False, # No color scale legend
            colorscale='Rainbow', # Colorscale for markers (can also use fixed colors)
            color=node_color, # Assign colors based on type
            size=node_size_viz, # Node sizes
            line_width=1.5,
            line_color='black'
        )
    )

    # Create the final Plotly figure
    fig = go.Figure(
        data=[edge_trace, node_trace], # Combine edge and node traces
        layout=go.Layout(
            title='Interactive Entity Relationship Network', # Figure title
            titlefont_size=16,
            showlegend=False, # Hide legend
            hovermode='closest', # Hover mode
            margin=dict(b=10, l=5, r=5, t=30), # Figure margins
            height=600, # Figure height
            xaxis=dict(showgrid=False, zeroline=False, showticklabels=False), # X-axis style
            yaxis=dict(showgrid=False, zeroline=False, showticklabels=False)  # Y-axis style
        )
    )
    return fig # Return the Plotly figure object

# --- MAIN APPLICATION LOGIC ---

# Load models once at the start
# This is a global call, and the result is used throughout the app.
# It's crucial this completes successfully.
MODELS = load_models()

def display_analysis_results(text_to_analyze: str, source_name: str = "Uploaded Content"):
    """
    A centralized function to display summaries, visualizations, and study materials
    for the provided text.
    """
    st.subheader(f"Analysis Results for: {source_name}")

    # Preprocess text for different audiences
    child_text_processed = preprocess_text(text_to_analyze, "child")
    standard_text_processed = preprocess_text(text_to_analyze)

    # Tabbed interface for different analysis outputs
    analysis_tabs = st.tabs(["📚 Summaries", "📊 Visualizations", "📝 Study Aids"])

    with analysis_tabs[0]: # Summaries Tab
        st.markdown("#### Audience-Specific Summaries")
        st.info("Summaries are generated for different understanding levels. Select an audience.")

        # Audience selection for summary
        summary_audience = st.selectbox(
            "Choose summary audience:",
            ["child", "student", "researcher", "expert"],
            key=f"summary_audience_{source_name.replace(' ','_')}" # Unique key for selectbox
        )

        if st.button(f"Generate {summary_audience.capitalize()} Summary", key=f"gen_summary_btn_{source_name.replace(' ','_')}"):
            with st.spinner(f"Generating {summary_audience} summary... Please wait."):
                text_for_summary = child_text_processed if summary_audience == "child" else standard_text_processed
                if not text_for_summary.strip():
                    st.warning("The text is too short or empty after preprocessing for this summary type.")
                else:
                    summary = generate_summary(MODELS, text_for_summary, summary_audience)
                    summary_container_class = "child-mode" if summary_audience == "child" else "summary-box"
                    st.markdown(f'<div class="{summary_container_class}">{summary}</div>', unsafe_allow_html=True)

    with analysis_tabs[1]: # Visualizations Tab
        st.markdown("#### Text Visualizations")
        st.write("Visual tools to explore the text's key themes and entities.")

        # Word Cloud Generation
        st.markdown("##### Word Cloud")
        if st.button("Generate Word Cloud", key=f"wc_btn_{source_name.replace(' ','_')}"):
            with st.spinner("Creating word cloud..."):
                wordcloud_fig = create_wordcloud(standard_text_processed)
                if wordcloud_fig:
                    st.pyplot(wordcloud_fig)
                    st.caption("Most frequent words in the text (excluding common stop words).")
                else:
                    st.info("Not enough text to generate a word cloud or an error occurred.")

        # Entity Network Visualization
        st.markdown("##### Entity Network")
        if st.button("Generate Entity Network", key=f"en_btn_{source_name.replace(' ','_')}"):
            with st.spinner("Analyzing entities and building network..."):
                doc = MODELS['nlp'](standard_text_processed[:50000]) # Limit for performance
                entities = [(ent.text, ent.label_) for ent in doc.ents]
                if entities:
                    network_fig = create_entity_network(entities, doc)
                    if network_fig:
                        st.plotly_chart(network_fig, use_container_width=True)
                        st.caption("Interactive network of named entities (people, organizations, locations, etc.) and their co-occurrences.")
                    else:
                        st.info("Could not generate entity network (e.g., too few entities).")
                else:
                    st.info("No named entities found to create a network.")

        # Text Statistics (moved to visualization for quicker insights)
        st.markdown("##### Text Statistics")
        if standard_text_processed:
            doc_stats = MODELS['nlp'](standard_text_processed[:100000]) # Limit for performance
            num_sentences = len(list(doc_stats.sents))
            num_unique_tokens = len(set(token.lemma_.lower() for token in doc_stats if not token.is_stop and not token.is_punct))
            num_words = len([token for token in doc_stats if not token.is_punct])
            # Simple readability: Flesch Reading Ease approximation (very basic)
            # score = 206.835 - 1.015 * (total_words / total_sentences) - 84.6 * (total_syllables / total_words)
            # For simplicity here, a proxy:
            avg_sent_len = num_words / num_sentences if num_sentences > 0 else 0
            readability_proxy = max(0, 100 - avg_sent_len) # Higher is better (simpler)

            col1, col2, col3 = st.columns(3)
            with col1: st.metric("Sentences", f"{num_sentences:,}")
            with col2: st.metric("Unique Terms", f"{num_unique_tokens:,}")
            with col3: st.metric("Avg. Sentence Length", f"{avg_sent_len:.1f} words")
        else:
            st.info("Not enough text for statistics.")


    with analysis_tabs[2]: # Study Aids Tab
        st.markdown("#### Study Aids Generator")
        st.write("Create flashcards, practice questions, and more from the text.")
        if st.button("Generate Study Materials", key=f"study_aid_btn_{source_name.replace(' ','_')}"):
            with st.spinner("Crafting study materials..."):
                materials = generate_study_materials(standard_text_processed, MODELS['nlp'])
                if materials and materials["key_phrases"]:
                    st.markdown("##### 🔑 Key Phrases Identified")
                    st.info(", ".join(materials["key_phrases"]))

                    st.markdown("##### 📇 Flashcards (Q&A)")
                    for i, card in enumerate(materials["flashcards"]):
                        question, answer = card.split("\nA: ", 1)
                        with st.expander(f"Flashcard {i+1}: {question.replace('**Q:','').replace('**','').strip()}", expanded=False):
                            st.markdown(f"A: {answer}")

                    st.markdown("##### ❓ Practice Questions")
                    for q_text in materials["questions"]:
                        st.markdown(f'<div class="flashcard">{q_text}</div>', unsafe_allow_html=True)

                    st.markdown("##### 📊 Recommended Diagram Types")
                    for diag_suggestion in materials["diagrams"]:
                        st.markdown(f"- {diag_suggestion}")

                    st.markdown("##### 📅 Suggested Study Plan")
                    for plan_item in materials["plan"]:
                        st.markdown(f"- {plan_item}")

                    # Download study materials
                    study_content_str = f"KEY PHRASES:\n{', '.join(materials['key_phrases'])}\n\n"
                    study_content_str += "FLASHCARDS:\n" + "\n\n".join(materials["flashcards"]) + "\n\n"
                    study_content_str += "PRACTICE QUESTIONS:\n" + "\n\n".join(materials["questions"]) + "\n\n"
                    study_content_str += "DIAGRAM SUGGESTIONS:\n" + "\n".join(materials["diagrams"]) + "\n\n"
                    study_content_str += "STUDY PLAN:\n" + "\n".join(materials["plan"])

                    st.download_button(
                        label="📥 Download All Study Materials (TXT)",
                        data=study_content_str,
                        file_name=f"study_materials_{source_name.replace(' ','_')}.txt",
                        mime="text/plain"
                    )
                elif materials and "flashcards" in materials and materials["flashcards"] and "No key phrases" in materials["flashcards"][0]:
                     st.warning("No key phrases were found in the text to generate detailed study materials.")
                else:
                    st.warning("Could not generate study materials. The text might be too short or lack distinct concepts.")


def main_application():
    """Main function to run the Streamlit application."""

    # --- Sidebar Navigation ---
    from PIL import Image
    st.sidebar.image(Image.open('QR3XoLs.jpeg'), width=100)  # Placeholder logo, width=100
    st.sidebar.caption("pokeman")
    

   
    st.sidebar.title("✨ AI Analyzer Suite ✨")
    st.sidebar.markdown("---") # Divider
    # Define application modes for navigation
    app_mode = st.sidebar.radio(
        "Choose a Tool:",
        ["🏠 Home", "💬 Text Input Analysis", "📄 Document File Analysis", "🎤 Media File Analysis (Audio/Video)", "🔗 URL Content Analysis", "💡 About"]
    )
    st.sidebar.markdown("---") # Divider
    st.sidebar.info("Upload content or paste text, then explore summaries, visualizations, and study aids.")


    # --- Page Content Based on Navigation ---
    if app_mode == "🏠 Home":
        st.title("Welcome to the AI Text & Media Analyzer Suite! 👋")
        st.markdown("""
            <div class="content-card">
                <p>Unlock insights from your text and media with our powerful AI-driven tools. This suite is designed for students, researchers, educators, and anyone looking to understand content deeply and efficiently.</p>
                <h4>🚀 **What You Can Do:**</h4>
                <ul>
                    <li>Paste text directly for quick analysis (<b>Text Input Analysis</b>).</li>
                    <li>Upload documents (PDF, DOCX, PPTX, TXT) to extract and analyze text (<b>Document File Analysis</b>).</li>
                    <li>Analyze spoken content from audio (MP3, WAV) or video (MP4) files (<b>Media File Analysis</b>).</li>
                    <li>Fetch and analyze text from web pages or YouTube transcripts (<b>URL Content Analysis</b>).</li>
                </ul>
                <h4>🎯 **For Each Input, Get:**</h4>
                <ul>
                    <li>👶 **Child-friendly** explanations.</li>
                    <li>🎓 **Student-focused** summaries.</li>
                    <li>🔬 **Researcher-oriented** analyses with key terms.</li>
                    <li>🧠 **Expert-level** detailed summaries.</li>
                    <li>📊 **Visualizations** like word clouds and entity networks.</li>
                    <li>📝 **Study Aids** including flashcards, practice questions, and study plans.</li>
                </ul>
                <p>Select an analysis tool from the sidebar to begin!</p>
            </div>
        """, unsafe_allow_html=True)
        st.balloons()


    elif app_mode == "💬 Text Input Analysis":
        st.header("💬 Direct Text Input Analysis")
        st.markdown("<div class='content-card'>Enter or paste your text below for a comprehensive analysis.</div>", unsafe_allow_html=True)
        # Text area for user input
        text_input = st.text_area("Paste your text here:", height=250, placeholder="Type or paste your content...", key="text_input_main")
        if text_input: # If text is entered
            display_analysis_results(text_input, "Pasted Text")
        else:
            st.info("Please paste some text in the area above to enable analysis.")


    elif app_mode == "📄 Document File Analysis":
        st.header("📄 Document File Analysis")
        st.markdown("<div class='content-card'>Upload your documents (PDF, DOCX, PPTX, TXT) to extract text and generate insights.</div>", unsafe_allow_html=True)
        # File uploader for document files
        uploaded_file = st.file_uploader(
            "Upload your document:",
            type=["pdf", "docx", "pptx", "txt"], # Supported file types
            key="doc_uploader"
        )
        if uploaded_file: # If a file is uploaded
            with st.spinner(f"Processing '{uploaded_file.name}'... Please wait."):
                # Save to a temporary file to pass its path to extraction functions
                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
                    tmp_file.write(uploaded_file.getvalue()) # Write uploaded file content to temp file
                    tmp_file_path = tmp_file.name # Get path of the temp file

                raw_text = extract_text_from_file(tmp_file_path, uploaded_file.type) # Extract text
                os.unlink(tmp_file_path) # Delete the temporary file after use

                if raw_text and raw_text.strip(): # If text extraction is successful
                    st.success(f"Successfully extracted text from '{uploaded_file.name}'.")
                    with st.expander("View Extracted Text", expanded=False):
                        st.text_area("Extracted Text:", value=raw_text, height=200, disabled=True)
                    display_analysis_results(raw_text, uploaded_file.name)
                else:
                    st.error(f"Could not extract text from '{uploaded_file.name}' or the document is empty.")


    elif app_mode == "🎤 Media File Analysis (Audio/Video)":
        st.header("🎤 Audio & Video Transcription and Analysis")
        st.markdown("<div class='content-card'>Upload audio (MP3, WAV, M4A) or video (MP4, MOV, AVI) files. The audio will be transcribed to text for analysis.</div>", unsafe_allow_html=True)
        # File uploader for audio/video files
        uploaded_media_file = st.file_uploader(
            "Upload your audio or video file:",
            type=["mp3", "wav", "m4a", "mp4", "mov", "avi", "ogg", "flac"], # Supported media types
            key="media_uploader"
        )
        if uploaded_media_file: # If a media file is uploaded
            with st.spinner(f"Processing media file '{uploaded_media_file.name}'..."):
                # Save to a temporary file
                file_extension = os.path.splitext(uploaded_media_file.name)[1]
                with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_media_file:
                    tmp_media_file.write(uploaded_media_file.getvalue())
                    tmp_media_file_path = tmp_media_file.name

                transcribed_text = None
                # Determine if it's video or audio
                if file_extension.lower() in ['.mp4', '.mov', '.avi']: # Video file
                    audio_output_path = os.path.join(tempfile.gettempdir(), "extracted_audio.wav")
                    if extract_audio_from_video(tmp_media_file_path, audio_output_path):
                        transcribed_text = transcribe_audio_with_whisper(audio_output_path, MODELS['transcription_model'])
                        if os.path.exists(audio_output_path): os.unlink(audio_output_path) # Clean up extracted audio
                    else:
                        st.error("Failed to extract audio from video.")
                else: # Audio file
                    transcribed_text = transcribe_audio_with_whisper(tmp_media_file_path, MODELS['transcription_model'])

                if os.path.exists(tmp_media_file_path): os.unlink(tmp_media_file_path) # Clean up uploaded media temp file

                if transcribed_text and transcribed_text.strip():
                    st.success(f"Successfully transcribed audio from '{uploaded_media_file.name}'.")
                    with st.expander("View Transcription", expanded=False):
                        st.text_area("Transcribed Text:", value=transcribed_text, height=200, disabled=True)
                    display_analysis_results(transcribed_text, f"Transcription of {uploaded_media_file.name}")
                elif transcribed_text is None: # Explicit check for None for transcription errors
                    st.error(f"Transcription failed for '{uploaded_media_file.name}'.")
                else: # Empty transcription
                    st.warning(f"Transcription resulted in empty text for '{uploaded_media_file.name}'. The media might be silent or unclear.")


    elif app_mode == "🔗 URL Content Analysis":
        st.header("🔗 Web Content & YouTube Transcript Analysis")
        st.markdown("<div class='content-card'>Enter a URL (webpage or YouTube video) to fetch its text content for analysis.</div>", unsafe_allow_html=True)
        # Text input for URL
        url_input = st.text_input("Enter URL (e.g., https://example.com or YouTube video link):", key="url_input_main")
        if st.button("Fetch and Analyze URL", key="fetch_url_btn"):
            if url_input:
                fetched_text = fetch_and_parse_url(url_input) # Fetch and parse URL
                if fetched_text and fetched_text.strip():
                    st.success("Successfully fetched content from URL.")
                    with st.expander("View Fetched Text", expanded=False):
                        st.text_area("Fetched Text:", value=fetched_text, height=200, disabled=True)
                    display_analysis_results(fetched_text, f"Content from {url_input.split('//')[-1].split('/')[0]}") # Use domain as source name
                elif fetched_text is None: # Error handled by fetch_and_parse_url
                    pass # Error message already shown by the function
                else:
                    st.error("Could not fetch significant text content from the URL or the content is empty.")
            else:
                st.warning("Please enter a URL.")


    elif app_mode == "💡 About":
        st.title("💡 About the AI Analyzer Suite")
        st.markdown(
            """
            <style>
                <div class="content-card">
                    <h4><strong>AI Text & Media Analyzer Suite - Version 2.5</strong></h4>
                    <p>This application leverages state-of-the-art Natural Language Processing (NLP) models to provide comprehensive text analysis and understanding tools.
                    It's built with Python using the Streamlit framework for the user interface.

                    <h5><strong>Core Technologies:</strong></h5>
                    <ul>
                        <li><strong>Summarization:</strong>
                            <ul>
                                <li>General Purpose: Google's T5 (Text-To-Text Transfer Transformer) - <code>t5-small</code> model.</li>
                                <li>Child-Friendly: Facebook's BART (Bidirectional Auto-Regressive Transformer) - <code>facebook/bart-large-cnn</code> model.</li>
                            </ul>
                        </li>
                        <li><strong>NLP Tasks (Entity Recognition, POS Tagging, etc.):</strong> SpaCy - <code>en_core_web_sm</code> model.</li>
                        <li><strong>Audio Transcription:</strong> OpenAI's Whisper - <code>base</code> model.</li>
                        <li><strong>Text Extraction:</strong> Libraries like <code>PyPDF2</code> (PDFs), <code>python-docx</code> (Word), <code>python-pptx</code> (PowerPoint).</li>
                        <li><strong>Web Content:</strong> <code>requests</code> and <code>BeautifulSoup4</code> for webpages, <code>youtube_transcript_api</code> for YouTube.</li>
                        <li><strong>Visualizations:</strong> <code>WordCloud</code>, <code>Matplotlib</code>, <code>Plotly</code>, and <code>NetworkX</code>.</li>
                    </ul>

                    <h5><strong>Key Features:</strong></h5>
                    <ul>
                        <li>Multi-format input: Direct text, various document types, audio/video files, and web URLs.</li>
                        <li>Audience-specific summarization tailored for children, students, researchers, and experts.</li>
                        <li>Automated generation of study aids: key phrases, flashcards, practice questions, diagram suggestions, and study plans.</li>
                        <li>Interactive visualizations: Word clouds and entity relationship networks.</li>
                        <li>Detailed text statistics.</li>
                    </ul>

                    <h5><strong>Developed For:</strong></h5>
                    <p>Educational purposes, research assistance, content understanding, and general productivity.
                    This tool aims to make complex information more accessible and actionable.</p>

                    <h5><strong>Disclaimer:</strong></h5>
                    <p>AI-generated content can sometimes be inaccurate or incomplete. Always critically evaluate the results.
                    Processing very large files or long audio/video may take time and consume resources.</p>

                    <hr>
                    <p style="text-align:center;"><em>Happy Analyzing!</em></p>
                </div>
            """,
            unsafe_allow_html=True,
        )

# Entry point of the script
if __name__ == "__main__":
    main_application()

