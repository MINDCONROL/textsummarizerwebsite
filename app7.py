# Import all original libraries plus additional ones needed for new features
import streamlit as st
import os
import tempfile
import random
from typing import Optional, List, Dict, Tuple, Any
from datetime import datetime
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER
import textwrap
import pandas as pd
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import plotly.graph_objects as go
import networkx as nx
import numpy as np
import torch
from collections import defaultdict
import re
import concurrent.futures
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from transformers import T5Tokenizer, T5ForConditionalGeneration, pipeline, BertTokenizer, BertModel
import spacy
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from youtube_transcript_api import NoTranscriptFound, TranscriptsDisabled
# Add the import where it's actually used for audio extraction from video
from moviepy.editor import VideoFileClip # Keep this import for video audio extraction
# Remove the redundant import if VideoFileClip is used directly: import moviepy.editor as mp
import whisper
from PIL import Image
from spacy import displacy
# Initialize NLTK with enhanced error handling
try:
    nltk.download('punkt', quiet=True)
    nltk.download('stopwords', quiet=True)
    nltk.download('wordnet', quiet=True)  # Added for better lemmatization
except Exception as e:
    st.warning(f"Could not download NLTK data automatically: {e}. Some features might be limited.")

# Enhanced Streamlit page configuration
st.set_page_config(
    layout="wide",
    page_title="Advanced AI Text & Media Analyzer Suite",
    page_icon=" 🧠 ",
    menu_items={
        'Get Help': 'https://example.com/help',
        'Report a bug': "https://example.com/bug",
        'About': "# Advanced AI Content Analysis Suite"
    }
)

# Enhanced CSS styling with additional UI improvements
st.markdown("""
<style>
/* General Streamlit App Styling */
.stApp {
    background-color: #004d40; /* Dark teal background for main page */
    color: white; /* Default text color to white for overall readability on dark background */
    line-height: 1.6; /* Increased line spacing for better readability from first block */
    font-family: 'Segoe UI', Roboto, Arial, sans-serif; /* Modern, readable font stack from first block */
}

/* Sidebar Styling (White) */
.css-1d391kg { /* Target the main sidebar container - Class name may vary with Streamlit versions */
    background-color: #FFFFFF; /* White background for sidebar from second block */
    padding: 20px; /* Increased padding from first block */
    box-shadow: 2px 2px 10px rgba(0,0,0,0.05); /* Softer shadow from first block */
    border-radius: 10px; /* Rounded corners from first block */
    color: black; /* Ensure default text color in sidebar is black for contrast */
}

/* Ensure text elements within the sidebar are black */
/* This might require targeting specific Streamlit component classes within the sidebar if default color inheritance is an issue */
.css-1d391kg label, /* Labels in sidebar (e.g., selectbox, radio, checkbox) */
.css-1d391kg .st-bv, /* Markdown text in sidebar */
.css-1d391kg .st-cc, /* Text elements potentially */
.css-1d391kg .st-cq, /* More text elements potentially */
.css-1d391kg div { /* General div within sidebar - be cautious with this */
    color: black; /* Set text color to black for these elements */
}
/* Note: The exact classes for text within the sidebar can be tricky and might need
   inspection of the rendered HTML if the above general rules aren't sufficient. */


/* Main Content Area Headings */
h1, h2, h3, h4, h5, h6 {
    color: white; /* Header text remains white for visibility on dark teal background */
}

/* Buttons Styling */
.stButton>button { /* Target Streamlit buttons */
    background-color: #00796b; /* Teal background from second block */
    color: white; /* White text for contrast */
    border-radius: 8px; /* Rounded corners from first block */
    padding: 12px 24px; /* Increased padding from first block */
    font-weight: 500; /* Medium font weight from first block */
    transition: all 0.3s ease; /* Smooth transition for hover effects from first block */
    cursor: pointer; /* Indicate interactivity from first block */
    border: none; /* Remove default border from first block */
}

.stButton>button:hover { /* Style for button on hover */
    background-color: #004d40; /* Darker shade on hover from second block */
    transform: translateY(-2px); /* Slight lift effect from first block */
    box-shadow: 0 4px 12px rgba(0,0,0,0.15); /* More pronounced shadow on hover from first block */
}

/* File Uploader Styling */
.stFileUploader label { /* Target file uploader label */
    color: white; /* Set file uploader label to white for readability on dark background */
}

.stFileUploader>div>section { /* Dropzone area */
    border: 2px dashed #00796b; /* Dotted border with teal accent color (from second block) */
    border-radius: 8px; /* Rounded corners from first block */
    background-color: rgba(0, 121, 107, 0.1); /* Very light tint of teal */
    padding: 20px; /* Padding inside the uploader area from first block */
    text-align: center; /* Center align text from first block */
    color: white; /* Make dropzone text white */
}

.stFileUploader section button[data-testid="baseButton-secondary"] { /* Browse files button */
    color: black; /* Text color for "Browse files" button */
    background-color: #e0f2f1; /* Light background for button from second block */
    border: 1px solid #00796b; /* Border for button from second block */
    border-radius: 8px; /* Match button border-radius */
    padding: 8px 15px; /* Adjust padding for this button */
}
.stFileUploader section button[data-testid="baseButton-secondary"]:hover {
    background-color: #c0e0df; /* Slightly darker on hover */
}


/* Text Input and Text Area Styling */
.stTextInput>div>div>input, .stTextArea>div>textarea {
    border: 1px solid #00796b; /* Teal border */
    border-radius: 8px; /* Rounded corners from first block */
    padding: 10px; /* Padding inside input fields from first block */
    width: 100%; /* Full width from first block */
    box-sizing: border-box; /* Include padding and border in size from first block */
    color: white; /* Set input text color to white */
    background-color: #005548; /* Slightly lighter background than main page */
}

.stTextInput label, .stTextArea label { /* Target labels for text input and text area */
    color: white; /* Set label text color to white */
}

.stTextInput>div>div>input:focus, .stTextArea>div>textarea:focus {
    border-color: #00796b; /* Highlight color on focus (teal) */
    outline: none; /* Remove default outline */
    box-shadow: 0 0 0 0.2rem rgba(0, 121, 107, 0.25); /* Glow effect on focus (teal) */
}

/* Content Cards and Similar Display Boxes */
.content-card {
    background-color: #e0f2f1; /* Light teal background from second block */
    border-radius: 10px; /* Rounded corners from first block */
    padding: 25px; /* Increased padding inside the card from first block */
    margin-bottom: 20px; /* Space below the card from first block */
    box-shadow: 0 2px 5px rgba(0,0,0,0.08); /* Subtle shadow from first block */
    border-left: 5px solid #00796b; /* Accent color left border (teal) from second block */
    line-height: 1.7; /* Ensure good line spacing within cards from first block */
    color: black; /* Set text color within these boxes to black from second block */
}

.summary-box { /* Styling for summary boxes from second block */
    border-left: 5px solid #00796b; /* Teal left border for emphasis */
    padding: 1rem; /* Padding inside the box */
    margin: 1rem 0; /* Margin around the box */
    background-color: #e0f2f1; /* Light teal background */
    border-radius: 5px; /* Rounded corners */
    box-shadow: 2px 2px 5px rgba(0,0,0,0.1); /* Subtle shadow for depth */
    color: black; /* Set text color within these boxes to black */
}

.child-mode { /* Specific styling for child-mode content from second block */
    background-color: #fff8e1; /* Light yellow for child mode */
    padding: 1rem; /* Padding */
    border-radius: 10px; /* Rounded corners */
    border-left: 5px solid #ffc107; /* Amber border */
    color: black; /* Set text color in child mode to black */
}

.flashcard { /* Styling for flashcard elements from second block */
    background-color: #f1f8e9; /* Light green background */
    border-radius: 8px; /* Rounded corners */
    padding: 15px; /* Padding */
    margin: 10px 0; /* Margin */
    box-shadow: 0 2px 4px rgba(0,0,0,0.1); /* Subtle shadow */
    border-left: 4px solid #8bc34a; /* Green border */
    color: black; /* Set text color in flashcards to black */
}


/* Tabs Styling */
.stTabs [data-baseweb="tab-list"] { /* Target the list of tabs */
    gap: 24px; /* Space between tabs from second block */
}
.stTabs [data-baseweb="tab"] { /* Target individual tab elements */
    height: 50px; /* Height of tabs from second block */
    white-space: pre-wrap; /* Allow text wrapping in tabs from second block */
    background-color: #e0f2f1; /* Light teal for inactive tabs from second block */
    border-radius: 8px 8px 0px 0px; /* Rounded top corners from first block */
    padding: 10px 15px; /* Padding for tabs (combined) */
    color: black; /* Set text color for inactive tabs to black from second block */
    transition: background-color 0.3s ease, color 0.3s ease; /* Smooth transition */
}
.stTabs [aria-selected="true"] { /* Target the currently selected (active) tab */
    background-color: #00796b !important; /* Teal for active tab from second block */
    color: white !important; /* Set text color for active tab to white from second block */
    font-weight: bold; /* Bold text for active tab from first block */
    border-radius: 8px 8px 0 0; /* Ensure rounded top corners */
    padding: 10px 15px; /* Ensure consistent padding */
}

/* Expander Styling */
.stExpander {
    background-color: #f8f9fa; /* Light background for expanders (from first block - readable on dark main) */
    border-radius: 8px; /* Rounded corners from first block */
    margin-bottom: 15px; /* Space below expanders from first block */
    border: 1px solid #e0e0e0; /* Subtle border from first block */
    color: black; /* Ensure text inside expander is black */
}

.stExpander div[role="button"] { /* Style for the expander header */
    padding: 15px; /* Padding in the header from first block */
    color: #1a374d; /* Slightly darker color for header text */
}

.stExpander div[data-testid="stExpanderBody"] { /* Style for expander content */
    color: black; /* Ensure text in the body is black */
}


/* Metrics Styling */
.stMetric { /* Target metric display elements */
    background-color: #e0f2f1; /* Light teal background (using color from second block's palette) */
    padding: 15px; /* Padding inside metrics from first block */
    border-radius: 8px; /* Rounded corners from first block */
    border-left: 4px solid #00796b; /* Accent color left border (teal) */
    margin-bottom: 15px; /* Space below metrics from first block */
    color: black; /* Set text color in metrics to black from second block */
}
.stMetric label { /* Target metric labels */
    color: black; /* Ensure metric labels are black */
}
.stMetric div[data-testid="stMetricValue"] { /* Target metric values */
    color: black; /* Ensure metric values are black */
}
.stMetric div[data-testid="stMetricDelta"] { /* Target metric delta */
    color: inherit; /* Inherit color, usually handled by Streamlit for +/- */
}


/* Responsive Layout Improvements for Smaller Screens */
@media (max-width: 768px) {
    .stTabs [data-baseweb="tab-list"] {
        flex-direction: column; /* Stack tabs vertically on small screens from first block */
    }
    .stButton>button {
        width: 100%; /* Full width buttons on small screens from first block */
        text-align: center;
    }
}

/* Additional minor improvements */
.stMarkdown h3 {
    margin-top: 1.5em; /* Space above h3 in markdown from first block */
}

.stMarkdown ul, .stMarkdown ol {
    margin-bottom: 1em; /* Space below lists in markdown from first block */
}

/* Ensure general markdown text is readable on dark background */
.stMarkdown {
    color: white;
}
.stMarkdown strong {
    color: white; /* Ensure bold text is white */
}
.stMarkdown em {
     color: #b2dfdb; /* A slightly lighter teal for emphasis */
}


/* Streamlit default text input/area labels - sometimes these need explicit color */
.st-emotion-cache-xyz label { /* Replace xyz with the actual class if needed */
   color: white; /* Ensure labels are white if not covered above */
}


</style>
""", unsafe_allow_html=True)

# Enhanced model loading with progress tracking and fallback options
@st.cache_resource
def load_models() -> Dict[str, Any]:
    """Load and cache all required NLP models with enhanced error handling and progress tracking."""
    try:
        with st.spinner("Loading AI models... This may take a few minutes"):
            # Load T5 model with progress indicator
            # st.write("Loading T5 model...")
            t5_tokenizer = T5Tokenizer.from_pretrained("t5-base")
            t5_model = T5ForConditionalGeneration.from_pretrained("t5-base")

            # Load BERT model with progress indicator
            # st.write("Loading BERT model...")
            bert_tokenizer = BertTokenizer.from_pretrained("bert-base-uncased")
            bert_model = BertModel.from_pretrained("bert-base-uncased")

            # Load specialized models
            # st.write("Loading child-friendly summarizer...")
            child_summarizer = pipeline(
                "summarization",
                model="facebook/bart-large-cnn",
                tokenizer="facebook/bart-large-cnn"
            )

            # Load spaCy with additional pipeline components
            # st.write("Loading spaCy NLP processor...")
            nlp = spacy.load("en_core_web_sm")
            nlp.add_pipe('sentencizer')

            # Load Whisper for transcription
            # st.write("Loading Whisper transcription model...")
            transcription_model = whisper.load_model("base")

            st.success("All models loaded successfully!")
            return {
                't5_tokenizer': t5_tokenizer,
                't5_model': t5_model,
                'bert_tokenizer': bert_tokenizer,
                'bert_model': bert_model,
                'child_summarizer': child_summarizer,
                'nlp': nlp,
                'transcription_model': transcription_model
            }
    except Exception as e:
        st.error(f"Model loading failed: {str(e)}")
        st.stop()

# Enhanced text extraction with better error handling and file type support
def extract_text_from_file(file_path: str, file_type: str) -> Optional[str]:
    """Universal text extraction with improved error handling and support for more file types."""
    try:
        if file_type == "application/pdf":
            reader = PdfReader(file_path)
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
            return text if text.strip() else None

        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = DocxDocument(file_path)
            return "\n".join(para.text for para in doc.paragraphs if para.text)

        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text)
            return "\n\n".join(text_runs) if text_runs else None

        elif file_type == "text/plain":
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        return None
    except Exception as e:
        st.error(f"File extraction error: {str(e)}")
        return None
    finally:
        # Ensure temporary file is deleted
        if 'file_path' in locals() and os.path.exists(file_path):
            try:
                os.unlink(file_path)
            except Exception as e:
                st.warning(f"Could not delete temporary file {file_path}: {e}")


# Advanced text preprocessing with more cleaning options
def preprocess_text(text: str, mode: str = "default") -> str:
    """Enhanced text cleaning with mode-specific rules and better handling of special cases."""
    if not text:
        return ""

    # Basic cleaning
    text = re.sub(r'\s+', ' ', text).strip()
    text = re.sub(r'\[.*?\]|\(.*?\)', '', text)  # Remove citations

    # Mode-specific processing
    if mode == "child":
        text = re.sub(r'[^\w\s.,!?\'\"\\]', '', text) # Added escaping for backslash
        text = re.sub(r'\b\w{15,}\b', '', text)
        text = text.lower()  # Make it more child-friendly
    elif mode == "technical":
        # Keep technical terms but clean other noise
        text = re.sub(r'[^\w\s.,;:\-+\'"()\[\]{}]', '', text)

    return text

# Utility function to get BERT embeddings for a given text
def get_bert_embeddings(text: str, tokenizer, model):
    """Get the mean-pooled BERT embedding for the input text."""
    if tokenizer is None or model is None:
        st.warning("BERT tokenizer or model not loaded.")
        return None
    try:
        # Consider making max_length configurable
        inputs = tokenizer(text, return_tensors="pt", truncation=True, max_length=128)
        with torch.no_grad():
            outputs = model(**inputs)
            # Use the last hidden state and mean pooling
            last_hidden_state = outputs.last_hidden_state # (batch_size, seq_len, hidden_size)
            attention_mask = inputs['attention_mask']
            mask_expanded = attention_mask.unsqueeze(-1).expand(last_hidden_state.size()).float()
            sum_embeddings = torch.sum(last_hidden_state * mask_expanded, 1)
            sum_mask = torch.clamp(mask_expanded.sum(1), min=1e-9)
            mean_pooled = sum_embeddings / sum_mask
            return mean_pooled
    except Exception as e:
        # More specific exception handling could be added here
        st.warning(f"BERT embedding generation failed: {e}")
        return None

# Enhanced BERT keyword extraction with semantic grouping and meaning generation
def extract_keywords_with_bert(text: str, models: dict) -> Dict[str, Dict[str, Any]]:
    """Advanced keyword extraction with semantic grouping, meaning generation, and no fixed limit."""
    if not text or not text.strip():
        return {}

    # Add checks for required models
    if 'bert_tokenizer' not in models or models['bert_tokenizer'] is None or \
       'bert_model' not in models or models['bert_model'] is None or \
       'nlp' not in models or models['nlp'] is None:
        st.warning("Required models for keyword extraction not loaded.")
        return {}

    try:
        # Get document embedding
        doc_embedding = get_bert_embeddings(text, models['bert_tokenizer'], models['bert_model'])
        if doc_embedding is None:
            return {}

        # Process text with spaCy (consider increasing the processing limit for long texts)
        doc_spacy = models['nlp'](text[:50000])

        # Extract candidate phrases (noun chunks and significant terms)
        candidates = list(set(
            [chunk.text.lower() for chunk in doc_spacy.noun_chunks if isinstance(chunk, spacy.tokens.span.Span)
             and len(chunk.text.split()) <= 5 and len(chunk.text) > 3] # Added type check
        ))

        candidates.extend(list(set(
            [token.lemma_.lower() for token in doc_spacy if isinstance(token, spacy.tokens.token.Token) # Added type check
             and token.pos_ in ['NOUN', 'PROPN', 'ADJ', 'VERB'] and not token.is_stop
             and not token.is_punct and len(token.lemma_) > 3]
        )))

        candidates = list(set(candidates))
        if not candidates:
            return {}

        # Get embeddings for candidates
        phrase_embeddings = []
        valid_phrases = []
        for phrase in candidates:
            phrase_emb = get_bert_embeddings(phrase, models['bert_tokenizer'], models['bert_model'])
            if phrase_emb is not None:
                phrase_embeddings.append(phrase_emb)
                valid_phrases.append(phrase)

        if not phrase_embeddings:
            return {}

        # Calculate similarities
        phrase_embeddings_tensor = torch.cat(phrase_embeddings, dim=0)
        similarities = cosine_similarity(
            doc_embedding.cpu().numpy(),
            phrase_embeddings_tensor.cpu().numpy()
        )[0]

        # Get all keywords with meanings and scores
        keywords_with_meanings = {}
        sorted_indices = np.argsort(similarities)[::-1]

        for i in sorted_indices:
            phrase = valid_phrases[i]
            score = float(similarities[i])

            # Generate contextual meaning using surrounding text
            context_sentences = []
            if hasattr(doc_spacy, 'sents') and doc_spacy.sents is not None:
                 try:
                     context_sentences = [sent.text for sent in doc_spacy.sents if isinstance(sent, spacy.tokens.span.Span) and phrase in sent.text.lower()]
                 except Exception as e:
                     st.warning(f"Error processing sentences for keyword context: {e}")


            context = context_sentences[0] if context_sentences else ""

            # Determine category based on POS tags
            pos_tags = []
            if models.get('nlp'):
                try:
                    doc_phrase = models['nlp'](phrase)
                    pos_tags = [token.pos_ for token in doc_phrase if hasattr(token, 'pos_')] # Check for pos_ attribute
                except Exception as e:
                    st.warning(f"Error processing phrase for POS tags: {e}")


            category = "Concept"
            if 'VERB' in pos_tags:
                category = "Action"
            elif 'ADJ' in pos_tags:
                category = "Attribute"

            keywords_with_meanings[phrase] = {
                'score': score,
                'meaning': f"{category} related to {context[:100]}..." if context else f"Important {category}",
                'category': category,
                'context': context[:200] + "..." if context else ""
            }

        return keywords_with_meanings

    except Exception as e:
        # More specific exception handling could be added here
        st.warning(f"Advanced keyword extraction failed: {e}")
        return {}

# Enhanced summary generation with multiple note-taking methods
def generate_advanced_summary(models: dict, text: str, audience: str) -> Dict[str, Any]:
    """Generate comprehensive summaries with multiple note-taking methods and structured formats."""
    if not text or not text.strip():
        return {"error": "Cannot generate summary from empty text."}

    try:
        # Preprocess text based on audience
        processed_text = preprocess_text(text, "child" if audience == "child" else "default")

        # --- Base Summary Generation ---
        base_summary = "Summary could not be generated." # Default value in case of failure
        if processed_text: # Only attempt summary if there's processed text
            if audience == "child":
                # Ensure models['child_summarizer'] is available and callable
                if models.get('child_summarizer') and callable(models['child_summarizer']):
                    try:
                        summary_output = models['child_summarizer'](
                            processed_text[:1024*4], # Consider making this limit configurable
                            max_length=150,
                            min_length=30,
                            do_sample=True,
                            temperature=0.8,
                            top_p=0.95
                        )
                        # **CORRECTION:** Check if summary_output is a list and not empty before subscripting
                        if isinstance(summary_output, list) and len(summary_output) > 0 and isinstance(summary_output[0], dict) and 'summary_text' in summary_output[0]:
                            base_summary = summary_output[0]['summary_text']
                        else:
                             st.warning("Child summarizer returned unexpected output format or is empty.")
                             base_summary = "Child summary generation failed due to unexpected output."
                    except Exception as e:
                        st.warning(f"Child summary generation failed: {e}")
                        base_summary = "Child summary generation failed."
                else:
                    st.warning("Child summarizer model not loaded or not callable.")
                    base_summary = "Child summarizer not available."
            else: # student, researcher, expert
                 # Ensure T5 model and tokenizer are available
                if models.get('t5_tokenizer') and models.get('t5_model'):
                    try:
                        prompts = {
                            "student": "summarize for a high school student in bullet points: ",
                            "researcher": "create a detailed research summary with key findings: ",
                            "expert": "provide a technical summary with analysis and implications: "
                        }
                        prompt_text = prompts.get(audience, "summarize: ") + processed_text[:3000] # Default prompt
                        inputs = models['t5_tokenizer'].encode(
                            prompt_text,
                            return_tensors="pt",
                            max_length=1024, # Consider making this limit configurable
                            truncation=True
                        )
                        outputs = models['t5_model'].generate(
                            inputs,
                            max_length=500,
                            min_length=150,
                            num_beams=4,
                            early_stopping=True,
                            temperature=0.7 if audience == "student" else 0.6
                        )
                         # **CORRECTION:** Check if outputs is a tensor and not empty before subscripting and decoding
                        if isinstance(outputs, torch.Tensor) and outputs.shape[0] > 0:
                            base_summary = models['t5_tokenizer'].decode(outputs[0], skip_special_tokens=True)
                        else:
                            st.warning("T5 model generated unexpected output format or is empty.")
                            base_summary = "Advanced summary generation failed due to unexpected output."

                    except Exception as e:
                         st.warning(f"T5 summary generation failed: {e}")
                         base_summary = "Advanced summary generation failed."
                else:
                    st.warning("T5 model or tokenizer not loaded.")
                    base_summary = "Advanced summarization model not available."
        else:
            st.warning("Processed text is empty, skipping summary generation.")
            base_summary = "No text to summarize."


        # --- spaCy Processing and Note Generation ---
        doc = None
        sentences = []
        entities = []
        try:
            # Process text with spaCy (consider increasing the processing limit)
            # Add a check for the nlp model
            if processed_text and models.get('nlp'):
                 doc = models['nlp'](processed_text[:50000])
                 # Safely extract sentences, ensuring they are spaCy spans
                 sentences = [sent for sent in doc.sents if isinstance(sent, spacy.tokens.span.Span) and len(sent.text.split()) > 5]
                 # Safely extract entities, ensuring they are spaCy spans
                 entities = [(ent.text, ent.label_) for ent in doc.ents if isinstance(ent, spacy.tokens.span.Span)]
            elif processed_text and not models.get('nlp'):
                 st.warning("spaCy model not loaded for detailed analysis.")
            else:
                 st.info("No processed text for detailed analysis (notes, etc.).")

        except Exception as e:
            st.warning(f"spaCy processing failed: {e}")
            # Keep sentences and entities as empty lists if spaCy fails


        # Generate Cornell Method Notes
        cornell_notes = {"Cue": [], "Notes": [], "Summary": []}
        if sentences:
             for i, sent in enumerate(sentences[:10]): # Limit to first 10 sentences for Cornell notes
                 cornell_notes["Cue"].append(f"Key Point {i+1}")
                 cornell_notes["Notes"].append(sent.text) # Use .text here
                 # Generate a simple summary based on the first few words of the sentence
                 summary_snippet = ' '.join(sent.text.split()[:5]) + "..." if sent.text else "..."
                 cornell_notes["Summary"].append(f"Important concept about {summary_snippet}")


        # Generate Mind Map Structure
        mind_map = {
            "Central Idea": base_summary.split('.')[0] if base_summary and '.' in base_summary else base_summary[:min(50, len(base_summary))] + "..." if base_summary else "Document Analysis", # Safer slicing
            "Main Branches": [],
            "Sub-branches": defaultdict(list)
        }
        if entities:
             for i, (ent_text, label) in enumerate(entities[:5]): # Limit branches for mind map
                 mind_map["Main Branches"].append(f"{label}: {ent_text}")
                 # Add related sentences as sub-branches
                 if sentences:
                     related_sents = [sent.text for sent in sentences if ent_text.lower() in sent.text.lower()][:2] # Limit sub-branches
                     for sent_text in related_sents:
                          mind_map["Sub-branches"][f"Branch {i+1}"].append(sent_text)


        # Generate Linear Notes
        linear_notes = []
        if sentences:
            for i, sent in enumerate(sentences[:15]): # Limit to first 15 sentences for linear notes
                linear_notes.append(f"{i+1}. {sent.text}") # Use .text here

        # Generate Concept Map suggestions
        concept_map_suggestions = []
        # Get concepts from noun chunks, ensuring they are spans before accessing .text
        concepts = [chunk.text for chunk in (doc.noun_chunks if doc and hasattr(doc, 'noun_chunks') else []) if isinstance(chunk, spacy.tokens.span.Span) and len(chunk.text.split()) > 1][:5] # Add check for noun_chunks existence and type check

        if len(concepts) >= 2:
             concept_map_suggestions.append(f"Consider a Concept Map showing relationships between: {', '.join(concepts)}")
             # Add more specific suggestions if possible
             if len(concepts) >= 3:
                 concept_map_suggestions.append(f"Explore how '{concepts[0]}' influences '{concepts[1]}' and '{concepts[2]}'.")


        return {
            "base_summary": base_summary,
            "cornell_notes": cornell_notes,
            "mind_map": mind_map,
            "linear_notes": linear_notes,
            "concept_map_suggestions": concept_map_suggestions,
            "key_phrases": [chunk.text for chunk in (doc.noun_chunks if doc and hasattr(doc, 'noun_chunks') else []) if isinstance(chunk, spacy.tokens.span.Span)][:15], # Add check for noun_chunks existence and type check
            "entities": entities # entities already filtered and formatted
        }

    except Exception as e:
        # Catch any other unexpected errors during the process
        st.error(f"An unexpected error occurred during advanced summary generation: {str(e)}")
        return {"error": f"Summary generation failed: {str(e)}"}


# Enhanced study materials generation with abstractive content
def generate_advanced_study_materials(text: str, nlp_model: Any) -> Dict[str, Any]:
    """Generate comprehensive study materials with abstractive content and contextual understanding."""
    if not text or not text.strip():
        return {}

    # Add a check for the nlp model
    if nlp_model is None:
        st.warning("spaCy model not loaded for study materials generation.")
        return {}

    try:
        # Process text with spaCy (consider increasing the processing limit)
        doc = nlp_model(text[:50000])
        sentences = [sent for sent in doc.sents if isinstance(sent, spacy.tokens.span.Span)] # Ensure sentences are spans

        # Abstractive Flashcards with contextual understanding
        flashcards = []
        # Ensure noun_chunks attribute exists and chunks are spans
        important_concepts = [chunk.text for chunk in (doc.noun_chunks if hasattr(doc, 'noun_chunks') else []) if isinstance(chunk, spacy.tokens.span.Span) and
                              len(chunk.text.split()) <= 3 and
                              any(t.pos_ in ['NOUN', 'PROPN'] for t in chunk)][:10] # Limit flashcards


        for concept in important_concepts:
            # Find sentences containing the concept
            context_sents = [sent.text for sent in sentences if concept.lower() in sent.text.lower()]
            context = random.choice(context_sents) if context_sents else ""

            # Generate question and answer based on context (consider using a question generation model for more abstractive questions)
            question = f"What is the significance of '{concept}' in this context?"
            answer = f"The concept '{concept}' relates to {context[:150]}..." if context else f"'{concept}' is an important concept mentioned."

            flashcards.append(f"Q: {question}\nA: {answer}")

        # Conceptual Questions based on relationships
        questions = []
        # Ensure noun_chunks attribute exists and chunks are spans
        concepts = [chunk.text for chunk in (doc.noun_chunks if hasattr(doc, 'noun_chunks') else []) if isinstance(chunk, spacy.tokens.span.Span) and len(chunk.text.split()) > 1][:10] # Limit concepts for questions

        if len(concepts) >= 2:
            for i in range(min(5, len(concepts) - 1)): # Limit questions
                questions.append(
                    f"How does '{concepts[i]}' relate to '{concepts[i+1]}' in this context?"
                )

        # Diagram Suggestions with specific guidance
        diagrams = []
        if len(concepts) >= 3: # Require at least 3 concepts for some suggestions
            diagrams.append(f"Concept Map: Connect '{concepts[0]}', '{concepts[1]}', and '{concepts[2]}' showing their relationships")

        diagrams.append("Timeline: Create a chronological sequence of events mentioned (if applicable)")
        diagrams.append("Venn Diagram: Compare and contrast two main concepts (if applicable)")
        diagrams.append("Flowchart: Map out any processes or sequences described (if applicable)")


        # Enhanced Study Plan with spaced repetition principles
        study_plan = [
            f"Day 1: Familiarize yourself with key concepts like '{important_concepts[0] if important_concepts else 'main topics'}'",
            "Day 2: Review the summary and linear notes",
            "Day 3: Test your understanding using the flashcards",
            "Day 4: Visualize relationships by sketching a concept map or flowchart",
            "Day 7: Review the material again using a different method (e.g., re-summarize a section)",
            "Weekly: Briefly review previous topics to reinforce learning"
        ]


        return {
            "flashcards": flashcards,
            "questions": questions,
            "diagrams": diagrams,
            "plan": study_plan,
            "key_concepts": important_concepts, # Use important_concepts which are more focused
            "summary": " ".join(sent.text for sent in sentences[:5]) if sentences else "" # Provide a brief summary
        }

    except Exception as e:
        # More specific exception handling could be added here
        st.warning(f"Study materials generation failed: {e}")
        return {}


# PDF generation function for comprehensive reports
def create_pdf_report(content: Dict[str, Any], title: str) -> BytesIO:
    """Generate a professional PDF report from analysis content."""
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)

    # Get the default stylesheet
    styles = getSampleStyleSheet()

    # Modify existing styles from the sample stylesheet
    # Check if styles exist before modifying
    if 'Heading1' in styles:
        styles['Heading1'].fontSize = 16
        styles['Heading1'].leading = 20
        styles['Heading1'].textColor = colors.darkblue
    if 'Heading2' in styles:
        styles['Heading2'].fontSize = 14
        styles['Heading2'].leading = 18
        styles['Heading2'].textColor = colors.darkblue
    if 'Bullet' in styles:
        styles['Bullet'].fontSize = 10
        styles['Bullet'].leading = 14
        styles['Bullet'].leftIndent = 20
    # **CORRECTION:** Modify the existing 'BodyText' style instead of trying to add a new one
    if 'BodyText' in styles:
         styles['BodyText'].fontSize = 10
         styles['BodyText'].leading = 14
         # If you wanted to base it on 'Normal' explicitly, you would do that here
         # For example: styles['BodyText'] = ParagraphStyle(name='BodyText', parent=styles['Normal'], ...)
         # but simply modifying the default BodyText is usually sufficient.


    # Add the custom 'Justify' style if it doesn't exist (add handles this, but checking is robust)
    if 'Justify' not in styles:
         styles.add(ParagraphStyle(name='Justify', alignment=TA_JUSTIFY, fontSize=10, leading=14))


    elements = []

    # Add title
    # Use .get() with a fallback to handle missing styles gracefully
    title_style = styles.get('Heading1', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
    elements.append(Paragraph(title, title_style))
    elements.append(Spacer(1, 12))

    # Add creation date
    # Use .get() with a fallback for BodyText style
    body_text_style = styles.get('BodyText', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
    elements.append(Paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M')}", body_text_style))
    elements.append(Spacer(1, 24))


    # Add summary section
    if 'base_summary' in content and content['base_summary']:
        summary_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Comprehensive Summary", summary_heading_style))
        # Ensure 'Justify' style exists before using it
        summary_content_style = styles.get('Justify', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback to Normal
        elements.append(Paragraph(content['base_summary'], summary_content_style))
        elements.append(Spacer(1, 18)) # Increased spacing

    # Add Cornell Notes section
    if 'cornell_notes' in content and content['cornell_notes'].get('Cue'):
        cornell_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Cornell Method Notes", cornell_heading_style))
        # Create a simple table-like structure using paragraphs for Cornell notes
        # Use .get() with fallback for BodyText style
        body_text_style = styles.get('BodyText', styles.get('Normal', getSampleStyleSheet()['Normal']))
        elements.append(Paragraph("<b>Cue | Notes</b>", body_text_style))
        elements.append(Spacer(1, 6))
        # Use .get() with fallback for justify style
        justify_style = styles.get('Justify', styles.get('Normal', getSampleStyleSheet()['Normal']))
        # Ensure cue and note lists are of the same length and contain strings
        for cue, note in zip(content['cornell_notes'].get('Cue', []), content['cornell_notes'].get('Notes', [])):
             if isinstance(cue, str) and isinstance(note, str):
                 elements.append(Paragraph(f"<b>{cue}:</b> {note}", justify_style)) # Use justify/summary style
                 elements.append(Spacer(1, 5))
        elements.append(Spacer(1, 18)) # Increased spacing


    # Add Key Concepts section
    if 'key_concepts' in content and content['key_concepts']:
        key_concepts_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Key Concepts", key_concepts_heading_style))
        # Ensure Bullet style exists before using it
        bullet_style = styles.get('Bullet', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        # Ensure concepts are strings
        for concept in content['key_concepts']:
            if isinstance(concept, str):
                 elements.append(Paragraph(f"• {concept}", bullet_style))
        elements.append(Spacer(1, 18)) # Increased spacing

    # Add Flashcards section
    if 'flashcards' in content and content['flashcards']:
        flashcards_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Flashcards", flashcards_heading_style))
        # Use .get() with fallback for BodyText style
        body_text_style = styles.get('BodyText', styles.get('Normal', getSampleStyleSheet()['Normal']))
        # Use .get() with fallback for justify style
        justify_style = styles.get('Justify', styles.get('Normal', getSampleStyleSheet()['Normal']))
        for card in content['flashcards']:
            if isinstance(card, str) and "\nA: " in card:
                q, a = card.split("\nA: ")
                elements.append(Paragraph(f"<b>Q:</b> {q}", body_text_style))
                elements.append(Paragraph(f"<b>A:</b> {a}", justify_style)) # Use justify/summary style
                elements.append(Spacer(1, 12)) # Increased spacing between flashcards
        elements.append(Spacer(1, 18)) # Increased spacing

    # Add Mind Map Structure suggestion
    if 'mind_map' in content and content['mind_map'].get('Central Idea'):
         mind_map_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
         elements.append(Paragraph("Mind Map Structure Suggestion", mind_map_heading_style))
         # Use .get() with fallback for BodyText and Bullet styles
         body_text_style = styles.get('BodyText', styles.get('Normal', getSampleStyleSheet()['Normal']))
         bullet_style = styles.get('Bullet', styles.get('Normal', getSampleStyleSheet()['Normal']))
         if isinstance(content['mind_map']['Central Idea'], str):
             elements.append(Paragraph(f"Central Idea: {content['mind_map']['Central Idea']}", body_text_style))
         if content['mind_map'].get('Main Branches'):
             elements.append(Paragraph("Main Branches:", body_text_style))
             for branch in content['mind_map']['Main Branches']:
                  if isinstance(branch, str):
                     elements.append(Paragraph(f"- {branch}", bullet_style))
             # Ensure sub-branches exist and are lists of strings
             for i, branch in enumerate(content['mind_map'].get('Main Branches', [])):
                  if content['mind_map']['Sub-branches'].get(f"Branch {i+1}") and isinstance(content['mind_map']['Sub-branches'][f"Branch {i+1}"], list):
                      elements.append(Paragraph(f"  Sub-branches for {branch}:", body_text_style))
                      for sub_branch in content['mind_map']['Sub-branches'][f"Branch {i+1}"][:]: # Iterate over a copy
                           if isinstance(sub_branch, str):
                              elements.append(Paragraph(f"    • {sub_branch}", bullet_style))

         elements.append(Spacer(1, 18)) # Increased spacing


    # Add Linear Notes section
    if 'linear_notes' in content and content['linear_notes']:
        linear_notes_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Linear Notes", linear_notes_heading_style))
        # Use .get() with fallback for BodyText style
        body_text_style = styles.get('BodyText', styles.get('Normal', getSampleStyleSheet()['Normal']))
        # Ensure notes are strings
        for note in content['linear_notes']:
            if isinstance(note, str):
                elements.append(Paragraph(note, body_text_style))
                elements.append(Spacer(1, 5)) # Space between notes
        elements.append(Spacer(1, 18)) # Increased spacing

    # Add Concept Map Suggestions section
    if 'concept_map_suggestions' in content and content['concept_map_suggestions']:
        concept_map_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Concept Map Suggestions", concept_map_heading_style))
        # Use .get() with fallback for Bullet style
        bullet_style = styles.get('Bullet', styles.get('Normal', getSampleStyleSheet()['Normal']))
        # Ensure suggestions are strings
        for suggestion in content['concept_map_suggestions']:
             if isinstance(suggestion, str):
                 elements.append(Paragraph(f"- {suggestion}", bullet_style))
        elements.append(Spacer(1, 18)) # Increased spacing

    # Add Practice Questions section
    if 'questions' in content and content['questions']:
        questions_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Practice Questions", questions_heading_style))
        # Use .get() with fallback for BodyText style
        body_text_style = styles.get('BodyText', styles.get('Normal', getSampleStyleSheet()['Normal']))
        # Ensure questions are strings
        for i, question in enumerate(content['questions']):
            if isinstance(question, str):
                 elements.append(Paragraph(f"{i+1}. {question}", body_text_style))
                 elements.append(Spacer(1, 10)) # Space between questions
        elements.append(Spacer(1, 18)) # Increased spacing

    # Add Diagram Suggestions section
    if 'diagrams' in content and content['diagrams']:
        diagrams_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Diagram Suggestions", diagrams_heading_style))
        # Use .get() with fallback for Bullet style
        bullet_style = styles.get('Bullet', styles.get('Normal', getSampleStyleSheet()['Normal']))
        # Ensure suggestions are strings
        for suggestion in content['diagrams']:
            if isinstance(suggestion, str):
                elements.append(Paragraph(f"- {suggestion}", bullet_style))
        elements.append(Spacer(1, 18)) # Increased spacing

     # Add Study Plan section
    if 'plan' in content and content['plan']:
        study_plan_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Study Plan", study_plan_heading_style))
        # Use .get() with fallback for BodyText style
        body_text_style = styles.get('BodyText', styles.get('Normal', getSampleStyleSheet()['Normal']))
        # Ensure plan items are strings
        for i, step in enumerate(content['plan']):
             if isinstance(step, str):
                 elements.append(Paragraph(f"{i+1}. {step}", body_text_style))
                 elements.append(Spacer(1, 5)) # Space between steps
        elements.append(Spacer(1, 18)) # Increased spacing


    # Add Entities section
    if 'entities' in content and content['entities']:
        entities_heading_style = styles.get('Heading2', styles.get('Normal', getSampleStyleSheet()['Normal'])) # Fallback
        elements.append(Paragraph("Named Entities", entities_heading_style))
        # Use .get() with fallback for BodyText style
        body_text_style = styles.get('BodyText', styles.get('Normal', getSampleStyleSheet()['Normal']))
        # Ensure entities are in the expected tuple format with string elements
        for ent_tuple in content['entities']:
            if isinstance(ent_tuple, tuple) and len(ent_tuple) == 2 and isinstance(ent_tuple[0], str) and isinstance(ent_tuple[1], str):
                ent, label = ent_tuple
                elements.append(Paragraph(f"<b>{ent}:</b> {label}", body_text_style))
                elements.append(Spacer(1, 5)) # Space between entities
        elements.append(Spacer(1, 18)) # Increased spacing


    doc.build(elements)
    buffer.seek(0)
    return buffer

# Enhanced display function with all new features integrated
def display_advanced_analysis_results(text_to_analyze: str, models: dict, source_name: str = "Uploaded Content"):
    """Display comprehensive analysis results with all new features."""
    st.subheader(f"Advanced Analysis Results for: {source_name}")

    # Preprocess text
    processed_text = preprocess_text(text_to_analyze)

    # Create tabs for different analysis sections
    tab1, tab2, tab3, tab4 = st.tabs([" 📚  Summaries", " 🔑  Keywords", " 📝  Study Aids", " 📊  Visualizations"])

    with tab1:
        st.markdown("#### Advanced Summarization")
        audience = st.selectbox("Select audience:", ["student", "researcher", "expert", "child"])

        if st.button(f"Generate {audience.capitalize()} Summary"):
            if processed_text:
                with st.spinner(f"Generating advanced {audience} summary..."):
                    analysis = generate_advanced_summary(models, text_to_analyze, audience)

                    if 'error' in analysis:
                        st.error(analysis['error'])
                    else:
                        # Display base summary
                        st.markdown("### Comprehensive Summary")
                        st.markdown(f'<div class="content-card">{analysis.get("base_summary", "Summary not available.")}</div>', unsafe_allow_html=True)

                        # Display Cornell Notes
                        if analysis.get('cornell_notes') and analysis['cornell_notes'].get('Cue'):
                             with st.expander("Cornell Method Notes"):
                                 # Simple table-like display
                                 st.markdown("---")
                                 st.markdown("**Cue | Notes**")
                                 st.markdown("---")
                                 # Ensure cue and note lists are of the same length and contain strings
                                 for cue, note in zip(analysis['cornell_notes'].get('Cue', []), analysis['cornell_notes'].get('Notes', [])):
                                     if isinstance(cue, str) and isinstance(note, str):
                                         st.markdown(f"**{cue}:** {note}")
                                         st.markdown("---")


                        # Display Mind Map structure
                        if analysis.get('mind_map') and analysis['mind_map'].get('Central Idea'):
                            with st.expander("Mind Map Structure"):
                                st.write("Central Idea:", analysis['mind_map']['Central Idea'])
                                if analysis['mind_map'].get('Main Branches'):
                                    st.write("Main Branches:")
                                    for branch in analysis['mind_map']['Main Branches']:
                                        if isinstance(branch, str):
                                            st.write(f"- {branch}")
                                    # Ensure sub-branches exist and are lists of strings
                                    for i, branch in enumerate(analysis['mind_map'].get('Main Branches', [])):
                                         if analysis['mind_map']['Sub-branches'].get(f"Branch {i+1}") and isinstance(analysis['mind_map']['Sub-branches'][f"Branch {i+1}"], list):
                                             st.write(f"  Sub-branches for {branch}:")
                                             for sub_branch in analysis['mind_map']['Sub-branches'][f"Branch {i+1}"][:]: # Iterate over a copy
                                                 if isinstance(sub_branch, str):
                                                     st.write(f"  • {sub_branch}")


                        # Add PDF download button
                        pdf_buffer = create_pdf_report(analysis, f"Analysis Report - {source_name}")
                        st.download_button(
                            label=" 📥  Download Full Report (PDF)",
                            data=pdf_buffer,
                            file_name=f"analysis_report_{source_name.replace(' ','_')}.pdf",
                            mime="application/pdf"
                        )
            else:
                 st.warning("No text to summarize.")


    with tab2:
        st.markdown("#### Semantic Keywords")
        if st.button("Extract Keywords"):
            if processed_text:
                with st.spinner("Analyzing text for semantic keywords..."):
                    keywords = extract_keywords_with_bert(processed_text, models)

                    if keywords:
                        # Display keywords in a structured way
                        st.markdown("### Key Concepts with Contextual Meaning")
                        for kw, data in keywords.items():
                            if isinstance(kw, str) and isinstance(data, dict):
                                with st.expander(f"{kw} ({data.get('category', 'N/A')}) - Score: {data.get('score', 0.0):.2f}"):
                                    st.markdown(f"**Meaning:** {data.get('meaning', 'N/A')}")
                                    if data.get('context'):
                                        st.markdown(f"**Context:** {data['context']}")

                        # Prepare keywords data for download
                        if keywords: # Ensure keywords is not empty before creating DataFrame
                            keywords_df = pd.DataFrame.from_dict(keywords, orient='index')
                            csv = keywords_df.to_csv(index=True)

                            st.download_button(
                                label=" 📥  Download Keywords (CSV)",
                                data=csv,
                                file_name=f"keywords_{source_name.replace(' ','_')}.csv",
                                mime="text/csv"
                            )
                        else:
                             st.info("No keywords to download.")
                    else:
                        st.warning("No significant keywords could be extracted.")
            else:
                 st.warning("No text to extract keywords from.")

    with tab3:
        st.markdown("#### Study Aids Generator")
        if st.button("Generate Study Materials"):
            if processed_text:
                with st.spinner("Creating advanced study materials..."):
                    materials = generate_advanced_study_materials(processed_text, models.get('nlp')) # Pass nlp model safely

                    if materials:
                        # Display flashcards
                        st.markdown("### Flashcards")
                        if materials.get('flashcards'):
                             for card in materials['flashcards']:
                                 # Split Q and A and display with expander, ensure card is a string
                                 if isinstance(card, str) and "\nA: " in card:
                                     q, a = card.split("\nA: ")
                                     with st.expander(q):
                                         st.write(a)
                        else:
                            st.info("No flashcards generated.")


                        # Display practice questions
                        st.markdown("### Practice Questions")
                        if materials.get('questions'):
                             # Ensure questions are strings
                             for q in materials['questions']:
                                 if isinstance(q, str):
                                     st.markdown(f'<div class="flashcard">{q}</div>', unsafe_allow_html=True) # Use the flashcard CSS class
                        else:
                             st.info("No practice questions generated.")


                        # Display diagram suggestions
                        st.markdown("### Diagram Suggestions")
                        if materials.get('diagrams'):
                             # Ensure suggestions are strings
                             for diagram in materials['diagrams']:
                                 if isinstance(diagram, str):
                                     st.markdown(f"- {diagram}")
                        else:
                             st.info("No diagram suggestions generated.")


                        # Display study plan
                        st.markdown("### Study Plan")
                        if materials.get('plan'):
                            # Ensure plan items are strings
                            for item in materials['plan']:
                                if isinstance(item, str):
                                     st.markdown(f"- {item}")
                        else:
                            st.info("No study plan generated.")


                        # Add PDF download button
                        pdf_buffer = create_pdf_report(materials, f"Study Materials - {source_name}")
                        st.download_button(
                            label=" 📥  Download Study Materials (PDF)",
                            data=pdf_buffer,
                            file_name=f"study_materials_{source_name.replace(' ','_')}.pdf",
                            mime="application/pdf"
                        )
                    else:
                        st.warning("Could not generate study materials from this content.")
            else:
                 st.warning("No text to generate study materials from.")


    with tab4:
        st.markdown("#### Text Visualizations")

        # Word Cloud
        if st.button("Generate Word Cloud"):
            if processed_text:
                with st.spinner("Creating word cloud..."):
                    fig = create_wordcloud(processed_text)
                    if fig:
                        st.pyplot(fig)
                    else:
                        st.warning("Could not generate word cloud.")
            else:
                 st.warning("No text to generate word cloud from.")


        # Entity Network
        if st.button("Generate Entity Network"):
            if processed_text:
                with st.spinner("Analyzing entities..."):
                    # Add check for nlp model before calling it
                    if models.get('nlp'):
                        doc = models['nlp'](processed_text[:50000]) # Consider increasing the processing limit
                        entities = [(ent.text, ent.label_) for ent in doc.ents if isinstance(ent, spacy.tokens.span.Span)] # Ensure entities are spans
                        if entities:
                            fig = create_entity_network(entities, doc)
                            if fig:
                                st.plotly_chart(fig, use_container_width=True)
                            else:
                                st.warning("Could not generate entity network.")
                        else:
                            st.info("No named entities found to visualize.")
                    else:
                         st.warning("spaCy model not loaded for entity network visualization.")
            else:
                 st.warning("No text to generate entity network from.")


# Utility function to create a word cloud visualization using matplotlib
def create_wordcloud(text: str):
    """Generate a word cloud from the input text and return a matplotlib figure."""
    if not text or not text.strip():
        return None
    try:
        stop_words = set(stopwords.words('english'))
        wordcloud = WordCloud(
            width=800,
            height=400,
            background_color='white',
            stopwords=stop_words,
            colormap='viridis'
        ).generate(text)
        fig, ax = plt.subplots(figsize=(10, 5))
        ax.imshow(wordcloud, interpolation='bilinear')
        ax.axis('off')
        plt.tight_layout(pad=0)
        return fig
    except Exception as e:
        st.warning(f"Could not generate word cloud: {e}")
        return None

# Utility function to create an entity network visualization using Plotly and NetworkX
def create_entity_network(entities, doc):
    """Create a Plotly network graph of named entities and their co-occurrences."""
    if not entities or doc is None:
        return None

    try:
        G = nx.Graph()

        # Add nodes for each unique entity
        unique_entities = list(set(entities))
        if not unique_entities:
            return None # Return None if no unique entities

        for ent, label in unique_entities:
            G.add_node(ent, label=label)

        # Add edges based on co-occurrence in the same sentence
        # Add check to ensure doc.sents is iterable and contains spans
        if hasattr(doc, 'sents') and doc.sents is not None:
            try:
                for sent in doc.sents:
                     if isinstance(sent, spacy.tokens.span.Span): # Ensure sent is a span
                        sent_ents = [ent.text for ent in sent.ents if isinstance(ent, spacy.tokens.span.Span) and ent.text in G.nodes()] # Only consider entities that are nodes and are spans
                        for i in range(len(sent_ents)):
                            for j in range(i + 1, len(sent_ents)):
                                entity1 = sent_ents[i]
                                entity2 = sent_ents[j]
                                # Ensure both entities are in the graph before adding an edge
                                if G.has_node(entity1) and G.has_node(entity2):
                                    if G.has_edge(entity1, entity2):
                                        G[entity1][entity2]['weight'] += 1
                                    else:
                                        G.add_edge(entity1, entity2, weight=1)
            except Exception as e:
                 st.warning(f"Error processing sentences for entity network: {e}")


        if len(G.nodes) == 0:
            return None

        # Use a layout that is less likely to overlap for better visualization
        pos = nx.spring_layout(G, k=0.7, iterations=50, seed=42) # Adjusted k and iterations


        edge_x = []
        edge_y = []
        for edge in G.edges():
            x0, y0 = pos[edge[0]]
            x1, y1 = pos[edge[1]]
            edge_x += [x0, x1, None]
            edge_y += [y0, y1, None]

        edge_trace = go.Scatter(
            x=edge_x, y=edge_y,
            line=dict(width=0.5, color='#888'),
            hoverinfo='none',
            mode='lines'
        )

        node_x = []
        node_y = []
        node_text = []
        node_colors = [] # To color nodes by entity type

        # Define colors for entity types
        label_colors = {
            'PERSON': '#FF5733', # Orange-red
            'ORG': '#33FF57',    # Green
            'GPE': '#337AFF',    # Blue
            'LOC': '#FF33A1',    # Pink
            'DATE': '#FFDA33',   # Yellow
            'CARDINAL': '#8A33FF', # Purple
            'NORP': '#33FFCE',   # Cyan
            'FAC': '#FF8833',    # Orange
            'PRODUCT': '#8AFF33', # Lime
            'EVENT': '#3349FF',  # Indigo
            'LAW': '#FF3333',    # Red
            'LANGUAGE': '#33B5FF', # Light Blue
            'OTHER': '#808080', # Grey for unknown types
            # Add more as needed
        }


        for node in G.nodes():
            x, y = pos[node]
            node_x.append(x)
            node_y.append(y)
            label = G.nodes[node].get('label', 'OTHER') # Default to 'OTHER'
            node_text.append(f"{node} ({label})")
            node_colors.append(label_colors.get(label, '#808080')) # Use defined colors or grey for unknown types


        node_trace = go.Scatter(
            x=node_x, y=node_y,
            mode='markers+text',
            textposition="top center",
            hoverinfo='text',
            marker=dict(
                showscale=False,
                color=node_colors, # Use generated colors
                size=18,
                line_width=2
            )
        )

        fig = go.Figure(data=[edge_trace, node_trace],
                        layout=go.Layout(
                            showlegend=False,
                            hovermode='closest',
                            margin=dict(b=20, l=5, r=5, t=40),
                            xaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
                            yaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
                            title="Entity Network",
                            # Add hover information for edges (optional, requires more complex trace)
                        ))
        return fig
    except Exception as e:
        # More specific exception handling could be added here
        st.warning(f"Could not generate entity network: {e}")
        return None


# Utility function to extract audio from video files
def extract_audio_from_video(video_path: str, audio_output_path: str) -> bool:
    """Extract audio from a video file and save it as a WAV file."""
    try:
        clip = VideoFileClip(video_path)
        clip.audio.write_audiofile(audio_output_path, codec='pcm_s16le')
        clip.close()
        return True
    except Exception as e:
        st.warning(f"Audio extraction failed: {e}")
        return False


# Utility function to transcribe audio using Whisper
def transcribe_audio_with_whisper(audio_path: str, transcription_model: Any) -> str:
    """Transcribe audio file using Whisper model."""
    if transcription_model is None:
        st.warning("Transcription model not loaded.")
        return ""
    try:
        result = transcription_model.transcribe(audio_path)
        return result["text"]
    except Exception as e:
        st.warning(f"Transcription failed: {e}")
        return ""
    finally:
        # Clean up the temporary audio file
        if 'audio_path' in locals() and os.path.exists(audio_path):
            try:
                os.unlink(audio_path)
            except Exception as e:
                st.warning(f"Could not delete temporary audio file {audio_path}: {e}")


# Utility function to fetch and parse text content from a URL
def fetch_and_parse_url(url: str) -> str:
    """Fetches the content from a URL and extracts visible text."""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (compatible; AdvancedAIAnalyzer/1.0)" # Use a descriptive User-Agent
        }
        # Add a timeout to prevent hanging
        response = requests.get(url, headers=headers, timeout=15) # Increased timeout slightly
        response.raise_for_status() # Raise an HTTPError for bad responses (4xx or 5xx)
        soup = BeautifulSoup(response.content, "html.parser")
        # Remove scripts, styles, and irrelevant tags
        for tag in soup(["script", "style", "header", "footer", "nav", "aside", "form", "noscript", "img", "meta", "link"]): # Added more tags to remove
            tag.decompose()
        # Extract visible text
        text = soup.get_text(separator=" ", strip=True)
        # Clean up excessive whitespace
        text = re.sub(r'\s+', ' ', text).strip() # Added strip after regex
        return text
    except requests.exceptions.RequestException as e:
        st.warning(f"Error fetching URL: {e}")
        return ""
    except Exception as e:
        st.warning(f"Failed to fetch or parse URL: {e}")
        return ""


# Main application function with all original functionality plus new features
def main_application():
    """Main application function with enhanced features integrated."""
    MODELS = load_models()
    from PIL import Image # Import PIL for image handling, consider moving to top if used more widely
    
    # --- Sidebar Navigation ---
    try: # Try to load an image for the sidebar
        # Replace 'QR3XoLs.jpeg' with a valid path to your logo or remove if not needed.
        # For robustness, check if file exists or use a default.
        # For this example, assuming it might not be present, we'll wrap it.
        # sidebar_logo_path = 'QR3XoLs.jpeg' # Path to logo
        # if os.path.exists(sidebar_logo_path): # Check if logo exists
        # st.sidebar.image(Image.open(sidebar_logo_path), width=100) # Display logo
        # st.sidebar.caption("Your App Name/Logo Caption") # Caption for logo
        pass # Commenting out the logo part as the file isn't provided with the script
    except FileNotFoundError: # Handle if image not found
        st.sidebar.warning("Logo image not found.") # Warning message
    except Exception as img_e: # Handle other image errors
        st.sidebar.warning(f"Could not load logo: {img_e}") # Warning message

        # --- Sidebar Navigation ---
    st.sidebar.image(Image.open('QR3XoLs.jpeg'), width=100)  # Placeholder logo, width=100
    st.sidebar.caption("pokeman")

    # Sidebar navigation - enhanced with icons and better organization
    # Ensure 'QR3XoLs.jpeg' exists or replace with a valid path or remove
    # st.sidebar.image(Image.open('QR3XoLs.jpeg'), width=100)
    st.sidebar.title(" 🧠  Advanced AI Analyzer")
    st.sidebar.markdown("---")

    app_mode = st.sidebar.radio(
        "Navigation:",
        [" 🏠  Home", " 💬  Text Analysis", " 📄  Document Analysis", " 🎤  Media Analysis", " 🌐  Web Analysis", " ℹ ️ About"],
        help="Select the type of content you want to analyze"
    )
    st.sidebar.markdown("---")
    st.sidebar.info("Upload content or paste text to generate advanced analysis, summaries, and study materials.")

    if app_mode == " 🏠  Home":
        st.title("Advanced AI Content Analysis Suite")
        st.markdown("""
        <div class="content-card">
        <h3>Comprehensive Content Understanding Platform</h3>
        <p>This advanced tool provides deep analysis of text, documents, media, and web content using state-of-the-art NLP techniques.</p>

        <h4>New Advanced Features:</h4>
        <ul>
            <li>Semantic keyword extraction with contextual meanings</li>
            <li>Multiple note-taking formats (Cornell, Mind Mapping, Linear)</li>
            <li>Abstractive study materials generation</li>
            <li>Professional PDF report generation</li>
            <li>Enhanced visualizations and entity analysis</li>
        </ul>

        <h4>How to Use:</h4>
        <ol>
            <li>Select analysis mode from the sidebar</li>
            <li>Upload or input your content</li>
            <li>Explore different analysis tabs</li>
            <li>Download comprehensive reports</li>
        </ol>
        </div>
        """, unsafe_allow_html=True)

    elif app_mode == " 💬  Text Analysis":
        st.header("Text Content Analysis")
        text_input = st.text_area("Paste your text here:", height=300,
                                 placeholder="Enter text to analyze...",
                                 help="Paste any text content for comprehensive analysis")

        if st.button("Analyze Text"):
            if text_input and text_input.strip():
                display_advanced_analysis_results(text_input, MODELS, "Pasted Text")
            else:
                st.warning("Please enter some text to analyze.")

    elif app_mode == " 📄  Document Analysis":
        st.header("Document File Analysis")
        uploaded_file = st.file_uploader(
            "Upload document (PDF, DOCX, PPTX, TXT):",
            type=["pdf", "docx", "pptx", "txt"],
            help="Supported formats: PDF, Word, PowerPoint, Text"
        )

        if uploaded_file:
            with st.spinner(f"Processing {uploaded_file.name}..."):
                # Use tempfile.NamedTemporaryFile within a with statement for automatic cleanup
                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
                     tmp_file.write(uploaded_file.getvalue())
                     tmp_file_path = tmp_file.name # Keep path to pass to extraction function


                raw_text = extract_text_from_file(tmp_file_path, uploaded_file.type)

                # The cleanup for tmp_file_path is now handled in extract_text_from_file's finally block


                if raw_text and raw_text.strip():
                    with st.expander("View Extracted Text"):
                        st.text_area("Extracted Text", value=raw_text, height=300, disabled=True)

                    display_advanced_analysis_results(raw_text, MODELS, uploaded_file.name)
                else:
                    st.error("Could not extract meaningful text from the document.")


    elif app_mode == " 🎤  Media Analysis":
        st.header("Audio/Video Content Analysis")
        uploaded_media = st.file_uploader(
            "Upload media file (MP3, WAV, MP4, etc.):",
            type=["mp3", "wav", "mp4", "m4a", "ogg", "flac", "mov", "avi"], # Added mov and avi
            help="Supported formats: MP3, WAV, MP4, M4A, OGG, FLAC, MOV, AVI"
        )

        if uploaded_media:
            with st.spinner(f"Processing {uploaded_media.name}..."):
                 # Use tempfile.NamedTemporaryFile within a with statement for automatic cleanup
                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_media.name)[1]) as tmp_media:
                    tmp_media.write(uploaded_media.getvalue())
                    tmp_media_path = tmp_media.name # Keep path to pass to functions

                transcribed_text = None
                audio_path = tmp_media_path # Assume input is audio by default

                # Handle video files by extracting audio
                if uploaded_media.name.lower().endswith(('.mp4', '.mov', '.avi')):
                    audio_output = os.path.join(tempfile.gettempdir(), f"extracted_audio_{os.path.basename(tmp_media_path)}.wav")
                    if extract_audio_from_video(tmp_media_path, audio_output):
                        audio_path = audio_output
                    # Clean up the temporary video file immediately after extraction attempt
                    try:
                        os.unlink(tmp_media_path)
                    except Exception as e:
                        st.warning(f"Could not delete temporary video file {tmp_media_path}: {e}")
                    tmp_media_path = None # Set to None to indicate it's deleted


                # Transcribe audio if audio_path is valid
                if audio_path and os.path.exists(audio_path): # Check if audio_path exists
                     transcribed_text = transcribe_audio_with_whisper(audio_path, MODELS.get('transcription_model')) # Pass model safely
                elif tmp_media_path and os.path.exists(tmp_media_path): # If it was an audio file originally and wasn't processed as video
                     transcribed_text = transcribe_audio_with_whisper(tmp_media_path, MODELS.get('transcription_model')) # Pass model safely
                else:
                     st.error("Could not find a valid audio file to transcribe.")


                # The cleanup for audio_path (if created) and tmp_media_path (if not video) is handled in transcribe_audio_with_whisper's finally block

                if transcribed_text and transcribed_text.strip():
                    with st.expander("View Transcription"):
                        st.text_area("Transcribed Text", value=transcribed_text, height=300, disabled=True)

                    display_advanced_analysis_results(transcribed_text, MODELS, f"Transcript of {uploaded_media.name}")
                else:
                    st.error("Could not transcribe meaningful text from the media file.")


    elif app_mode == " 🌐  Web Analysis":
        st.header("Web Content Analysis")
        url_input = st.text_input("Enter URL:", placeholder="https://example.com")

        if st.button("Analyze URL Content"):
            if url_input and url_input.strip():
                with st.spinner("Fetching and analyzing URL content..."):
                    fetched_text = fetch_and_parse_url(url_input)

                    if fetched_text and fetched_text.strip():
                        # Improved source name extraction
                        try:
                            source_name = requests.utils.urlparse(url_input).netloc or "Web Content"
                        except Exception:
                             source_name = "Web Content"

                        with st.expander("View Fetched Content"):
                            st.text_area("Content", value=fetched_text, height=300, disabled=True)

                        display_advanced_analysis_results(fetched_text, MODELS, source_name)
                    else:
                        st.error("Could not fetch meaningful content from the URL. Please ensure the URL is valid and accessible.")
            else:
                st.warning("Please enter a valid URL.")

    elif app_mode == " ℹ ️ About":
        st.title("About the Advanced AI Analyzer")
        st.markdown("""
        <div class="content-card">
        <h3>Advanced AI Content Analysis Suite</h3>
        <p>Version 4.0 with enhanced features for comprehensive content understanding.</p>

        <h4>Core Technologies:</h4>
        <ul>
            <li>Transformers (T5, BERT) for summarization and semantic analysis</li>
            <li>spaCy for advanced NLP processing</li>
            <li>Whisper for audio transcription</li>
            <li>ReportLab for professional PDF generation</li>
            <li>Plotly/Matplotlib for interactive visualizations</li>
        </ul>

        <h4>Key Features:</h4>
        <ul>
            <li>Multi-format content analysis (text, documents, media, web)</li>
            <li>Audience-specific summarization techniques</li>
            <li>Contextual keyword extraction with semantic meanings</li>
            <li>Multiple note-taking formats (Cornell, Mind Maps, Linear)</li>
            <li>Abstractive study materials generation</li>
            <li>Interactive visualizations and entity networks</li>
            <li>Professional report generation in PDF format</li>
        </ul>

        <p>This tool is designed for students, researchers, educators, and professionals who need to deeply understand and work with various content formats.</p>
        </div>
        """, unsafe_allow_html=True)


if __name__ == "__main__":
    main_application()