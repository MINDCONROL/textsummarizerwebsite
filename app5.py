# Import necessary libraries
import streamlit as st  # For creating the web application
import os  # For operating system dependent functionalities (like file paths)
import tempfile  # For creating temporary files and directories
import random  # For generating random numbers (used in analogies, study questions)
from typing import Optional, List, Dict, Tuple, Any  # For type hinting, improving code readability

# NLP and Machine Learning Libraries
from transformers import T5Tokenizer, T5ForConditionalGeneration, pipeline, BertTokenizer, BertModel # For T5, BERT, and specialized pipelines
import nltk  # Natural Language Toolkit for text processing
import spacy  # For advanced NLP tasks like NER, POS tagging
import torch # PyTorch for tensor operations, especially with BERT

# Text Processing and Data Handling
import re  # For regular expression operations
from collections import defaultdict  # For dictionaries with default values
from sklearn.feature_extraction.text import TfidfVectorizer  # For TF-IDF calculation
from sklearn.metrics.pairwise import cosine_similarity # For comparing BERT embeddings
import numpy as np  # For numerical operations

# File Extraction Libraries
from PyPDF2 import PdfReader  # For reading text from PDF files
from docx import Document as DocxDocument  # For reading text from DOCX files (aliased to avoid conflict)
from pptx import Presentation  # For reading text from PPTX files

# Visualization Libraries
import networkx as nx  # For creating and manipulating complex networks
import matplotlib.pyplot as plt  # For creating static, animated, and interactive visualizations
from wordcloud import WordCloud  # For generating word cloud images
import plotly.graph_objects as go  # For creating interactive Plotly graphs

# Web Content and Audio/Video Processing
import requests  # For making HTTP requests to fetch web content
from bs4 import BeautifulSoup  # For parsing HTML and XML documents
from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled  # For fetching YouTube transcripts
from moviepy.editor import VideoFileClip # For processing video and audio files
import whisper # For audio transcription using OpenAI's Whisper model

# Concurrency
import concurrent.futures  # For running operations in parallel (like multiple summaries)

# --- Global Variables & Configuration ---
# Attempt to download NLTK data, handle potential errors gracefully
try:
    nltk.download('punkt', quiet=True)  # Tokenizer models for sentence splitting
    nltk.download('stopwords', quiet=True)  # Stopwords list for text cleaning
except Exception as e: # Catch any exception during NLTK data download
    st.warning(f"Could not download NLTK data automatically: {e}. Some features might be limited.") # Show a warning to the user

# --- Streamlit Page Configuration ---
# Set up the page layout, title, and icon for the Streamlit app
# This should be the first Streamlit command in your script, and it can only be set once.
st.set_page_config(
    layout="wide",  # Use wide layout for more space
    page_title="AI Text & Media Analyzer Suite (Enhanced)",  # Title appearing in the browser tab
    page_icon="🌟"  # Icon for the browser tab (can be an emoji or a URL)
)

# --- Custom CSS Styling ---
# Apply custom CSS for a more appealing and consistent user interface
st.markdown("""
<style>
    /* Main app background - NOW DARK TEAL */
    .stApp { /* Target the main Streamlit app container */
        background-color: #004d40; /* Dark teal background for main page */
        color: white; /* Default text color to white for overall readability on dark background */
    }

    /* Sidebar styling - NOW WHITE */
    .css-1d391kg { /* Sidebar main class (this class might change with Streamlit versions) */
        background-color: #FFFFFF; /* White background for sidebar */
        padding: 10px; /* Add some padding within the sidebar */
        color: black; /* Ensure text in the sidebar is black for contrast on light background */
    }

    .css-1d391kg .stSelectbox label { /* Sidebar selectbox label */
        color: black; /* Ensure label text is black for readability on light sidebar */
    }

    /* This class might be generic, ensure it targets sidebar specific elements if needed, or adjust more globally */
    /* Assuming .st-emotion-cache-17x3wpl is for sidebar text elements based on original context */
    .css-1d391kg .st-emotion-cache-17x3wpl, /* Targeting within sidebar */
    .st-emotion-cache-17x3wpl { /* General rule if it was meant for sidebar before, might need refinement */
        color: black; /* Ensure other text elements in sidebar are black */
    }
    /* If .st-emotion-cache-17x3wpl was a very general class, the above might make other parts black too.
       It's safer to target it within .css-1d391kg if it's specific to sidebar.
       If it's a general Streamlit component class, its color should be context-dependent.
       For now, making it black for the sidebar context.
    */


    /* Main content area titles */
    h1, h2, h3 { /* Target header elements */
        color: white; /* Header text remains white, which is good for the new dark teal main background */
    }

    /* Buttons styling - unchanged, should still look good */
    .stButton>button { /* Target Streamlit buttons */
        background-color: #00796b; /* Teal background */
        color: white; /* White text */
        border-radius: 5px; /* Rounded corners for buttons */
        padding: 10px 15px; /* Padding inside buttons */
        border: none; /* No border for a cleaner look */
    }
    .stButton>button:hover { /* Style for button on hover */
        background-color: #004d40; /* Darker teal on hover */
        color: white; /* Maintain white text on hover */
    }

    /* File uploader styling - text color adjusted for dark main background */
    .stFileUploader label { /* Target file uploader label */
        color: white; /* Set file uploader label to white for readability on dark background */
    }
    .stFileUploader div[data-baseweb="input"] > div { /* Style the input text color in file uploader */
        color: white !important; /* Ensure text inside file uploader input is white on dark background */
        background-color: #006050; /* Slightly lighter background for the input field itself for better visibility */
        border: 1px solid #00796b; /* Optional: border for the input part */
    }
    .stFileUploader section div[data-testid="stFileUploaderDropzone"] { /* Dropzone text */
        color: white; /* Make dropzone text white */
    }
     .stFileUploader section button[data-testid="baseButton-secondary"] { /* Browse files button */
        color: black; /* Text color for "Browse files" button */
        background-color: #e0f2f1; /* Light background for button */
        border: 1px solid #00796b; /* Border for button */
    }
    .stFileUploader section button[data-testid="baseButton-secondary"]:hover {
        background-color: #c0e0df; /* Slightly darker on hover */
    }


    /* Text input/area styling - text color adjusted for dark main background */
    .stTextInput input, .stTextArea textarea { /* Target text input and text area elements */
        border: 1px solid #00796b; /* Teal border (consider a lighter border if too dark) */
        border-radius: 5px; /* Rounded corners */
        color: white; /* Set input text color to white */
        background-color: #005548; /* Slightly lighter background than main page for text fields */
    }
    .stTextInput label, .stTextArea label { /* Target labels for text input and text area */
        color: white; /* Set label text color to white */
    }

    /* Summary boxes and content cards - unchanged, these have their own light backgrounds */
    .summary-box, .content-card { /* Class for styling summary and content display boxes */
        border-left: 5px solid #00796b; /* Teal left border for emphasis */
        padding: 1rem; /* Padding inside the box */
        margin: 1rem 0; /* Margin around the box */
        background-color: #e0f2f1; /* Light teal background */
        border-radius: 5px; /* Rounded corners */
        box-shadow: 2px 2px 5px rgba(0,0,0,0.1); /* Subtle shadow for depth */
        color: black; /* Set text color within these boxes to black */
    }
    .child-mode { /* Specific styling for child-mode content */
        background-color: #fff8e1; /* Light yellow for child mode */
        padding: 1rem; /* Padding */
        border-radius: 10px; /* Rounded corners */
        border-left: 5px solid #ffc107; /* Amber border */
        color: black; /* Set text color in child mode to black */
    }
    .flashcard { /* Styling for flashcard elements */
        background-color: #f1f8e9; /* Light green background */
        border-radius: 8px; /* Rounded corners */
        padding: 15px; /* Padding */
        margin: 10px 0; /* Margin */
        box-shadow: 0 2px 4px rgba(0,0,0,0.1); /* Subtle shadow */
        border-left: 4px solid #8bc34a; /* Green border */
        color: black; /* Set text color in flashcards to black */
    }

    /* Tabs styling - inactive/active tabs should still contrast well with their parent backgrounds */
    .stTabs [data-baseweb="tab-list"] { /* Target the list of tabs */
        gap: 24px; /* Space between tabs */
    }
    .stTabs [data-baseweb="tab"] { /* Target individual tab elements */
        height: 50px; /* Height of tabs */
        white-space: pre-wrap; /* Allow text wrapping in tabs */
        background-color: #e0f2f1; /* Light teal for inactive tabs (contrasts with dark main background) */
        border-radius: 4px 4px 0px 0px; /* Rounded top corners */
        padding: 10px; /* Padding */
        color: black; /* Set text color for inactive tabs to black */
    }
    .stTabs [aria-selected="true"] { /* Target the currently selected (active) tab */
        background-color: #00796b; /* Teal for active tab */
        color: white; /* Set text color for active tab to white for contrast */
    }

    /* Metrics styling - unchanged, these have their own light backgrounds */
    .stMetric { /* Target metric display elements */
        background-color: #e8eaf6; /* Indigo light background */
        padding: 10px; /* Padding */
        border-radius: 5px; /* Rounded corners */
        border: 1px solid #3f51b5; /* Indigo border */
        color: black; /* Set text color in metrics to black */
    }
    .stMetricLabel { /* Target metric labels */
        color: black; /* Ensure metric labels are black */
    }
    .stMetricValue { /* Target metric values */
        color: black; /* Ensure metric values are black */
    }
</style>
""", unsafe_allow_html=True)

# --- CORE MODEL LOADING ---
@st.cache_resource # Cache the loaded models to avoid reloading on each interaction, improving performance
def load_models() -> Dict[str, Any]:
    """Load and cache all required NLP models."""
    try: # Start try block for error handling during model loading
        # Main T5 model for general summarization (changed to t5-base)
        t5_tokenizer = T5Tokenizer.from_pretrained("t5-base") # Load T5-base tokenizer
        t5_model = T5ForConditionalGeneration.from_pretrained("t5-base") # Load T5-base model

        # BERT model for keyword extraction or other enhancements
        bert_tokenizer = BertTokenizer.from_pretrained("bert-base-uncased") # Load BERT tokenizer
        bert_model = BertModel.from_pretrained("bert-base-uncased") # Load BERT model

        # Specialized summarizer for child-friendly summaries (BART is good for this)
        child_summarizer = pipeline( # Create a summarization pipeline
            "summarization", # Specify the task
            model="facebook/bart-large-cnn",  # BART model fine-tuned on CNN/DailyMail
            tokenizer="facebook/bart-large-cnn" # Corresponding tokenizer for BART
        )
        # spaCy NLP processor for tasks like NER, POS tagging, noun chunks
        nlp = spacy.load("en_core_web_sm") # Load small English model from spaCy

        # Whisper model for audio transcription
        # It will download the model on first use if not already present
        # Using "base" model for a balance of speed and accuracy in a web app context
        transcription_model = whisper.load_model("base") # Load base Whisper model

        return { # Return a dictionary of loaded models
            't5_tokenizer': t5_tokenizer, # T5 tokenizer
            't5_model': t5_model, # T5 model
            'bert_tokenizer': bert_tokenizer, # BERT tokenizer
            'bert_model': bert_model, # BERT model
            'child_summarizer': child_summarizer, # Child-friendly summarizer pipeline
            'nlp': nlp, # spaCy NLP model
            'transcription_model': transcription_model # Whisper transcription model
        }
    except Exception as e: # Catch any exception during model loading
        st.error(f"Fatal Error: Model loading failed: {str(e)}. The application cannot continue.") # Display error message
        st.stop() # Stop the Streamlit application if models can't load

# --- FILE AND CONTENT PROCESSING UTILITIES ---
def extract_text_from_file(file_path: str, file_type: str) -> Optional[str]:
    """Universal text extraction for multiple file formats."""
    try: # Start try block for error handling during file processing
        # PDF file processing
        if file_type == "application/pdf": # Check if the file type is PDF
            reader = PdfReader(file_path) # Initialize PDF reader with the file path
            # Join text from all pages, checking if text extraction yields content
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        # DOCX file processing
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document": # Check for DOCX
            doc = DocxDocument(file_path) # Initialize DOCX document reader
            # Join text from all paragraphs, checking if paragraph text exists
            return "\n".join(para.text for para in doc.paragraphs if para.text)
        # PPTX file processing
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation": # Check for PPTX
            prs = Presentation(file_path) # Initialize PPTX presentation reader
            text_runs = [] # List to hold text extracted from shapes
            for slide in prs.slides: # Iterate through each slide in the presentation
                for shape in slide.shapes: # Iterate through each shape on the slide
                    if hasattr(shape, "text") and shape.text.strip(): # Check if shape has text and it's not empty
                        text_runs.append(shape.text) # Add the shape's text to the list
            return "\n\n".join(text_runs) # Join text from shapes with double newline for separation
        # TXT file processing
        elif file_type == "text/plain": # Check for plain text file
            with open(file_path, 'r', encoding='utf-8') as f: # Open in read mode with UTF-8 encoding
                return f.read() # Read the entire file content
        return None # Return None if file type is not supported or no text found
    except Exception as e: # Catch any exception during file extraction
        st.error(f"File Extraction Error: {str(e)}") # Display error message
        return None # Return None on error

def preprocess_text(text: str, mode: str = "default") -> str:
    """Text cleaning with mode-specific rules."""
    if not text: return "" # Return empty string if input text is None or empty
    text = re.sub(r'\s+', ' ', text).strip() # Replace multiple whitespaces with single, strip leading/trailing spaces
    # Mode-specific preprocessing
    if mode == "child": # If mode is for children
        # Keep simple punctuation, remove complex characters for child-friendliness
        text = re.sub(r'[^\w\s.,!?\'"-]', '', text) # Allow alphanumeric, whitespace, and basic punctuation
        text = re.sub(r'\b\w{15,}\b', '', text)  # Remove very long words (often jargon)
    else: # Default mode for other audiences
        text = re.sub(r'\[.*?\]|\(.*?\)', '', text)  # Remove bracketed/parenthesized citations or notes
    return text # Return the preprocessed text

def transcribe_audio_with_whisper(audio_file_path: str, model: Any) -> Optional[str]:
    """Transcribes an audio file using OpenAI's Whisper model."""
    try: # Start try block for error handling
        with st.spinner("Transcribing audio... This may take a few moments."): # Show a spinner during transcription
            result = model.transcribe(audio_file_path) # Perform transcription using the loaded Whisper model
        return result["text"] # Return the transcribed text from the result dictionary
    except Exception as e: # Catch transcription errors
        st.error(f"Audio Transcription Error: {str(e)}") # Display error message
        return None # Return None on error

def extract_audio_from_video(video_file_path: str, audio_output_path: str) -> bool:
    """Extracts audio from a video file and saves it."""
    try: # Start try block for error handling
        with st.spinner("Extracting audio from video..."): # Show a spinner during audio extraction
            video_clip = VideoFileClip(video_file_path) # Load video clip using moviepy
            audio_clip = video_clip.audio # Get audio track from the video
            if audio_clip: # If an audio track exists
                audio_clip.write_audiofile(audio_output_path, codec='pcm_s16le') # Write audio to file (WAV format is often good for Whisper)
                audio_clip.close() # Close the audio clip to free resources
                video_clip.close() # Close the video clip to free resources
                return True # Indicate success
            else: # If no audio track is found
                st.warning("No audio track found in the video.") # Show a warning
                video_clip.close() # Close the video clip
                return False # Indicate no audio track
    except Exception as e: # Catch audio extraction errors
        st.error(f"Audio Extraction Error: {str(e)}") # Display error message
        if 'video_clip' in locals() and video_clip: video_clip.close() # Ensure clip is closed on error
        if 'audio_clip' in locals() and audio_clip: audio_clip.close() # Ensure clip is closed on error
        return False # Return False on error

def fetch_and_parse_url(url: str) -> Optional[str]:
    """Fetches content from a URL (webpage or YouTube) and extracts text."""
    try: # Start try block
        with st.spinner(f"Fetching content from {url}..."): # Show spinner
            # Check if it's a YouTube URL
            if "youtube.com/watch?v=" in url or "youtu.be/" in url: # Common YouTube URL patterns
                video_id = None # Initialize video_id
                if "watch?v=" in url: # Standard YouTube URL
                    video_id = url.split("watch?v=")[1].split("&")[0] # Extract video ID
                elif "youtu.be/" in url: # Shortened YouTube URL
                    video_id = url.split("youtu.be/")[1].split("?")[0] # Extract video ID

                if video_id: # If video ID was found
                    try: # Try to get transcript
                        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id) # Get available transcripts
                        # Try to get a manually created transcript first, then an auto-generated one
                        transcript = transcript_list.find_manually_created_transcript(['en']) \
                                     or transcript_list.find_generated_transcript(['en']) # Prioritize manual English transcript
                        transcript_text = " ".join([item['text'] for item in transcript.fetch()]) # Fetch and join transcript parts
                        return transcript_text # Return transcript
                    except NoTranscriptFound: # If no transcript is found
                        st.error(f"No English transcript found for YouTube video: {url}") # Error message
                        return None # Return None
                    except TranscriptsDisabled: # If transcripts are disabled
                        st.error(f"Transcripts are disabled for YouTube video: {url}") # Error message
                        return None # Return None
                else: # If video ID could not be parsed
                    st.error("Could not parse YouTube video ID.") # Error message
                    return None # Return None
            else: # Regular webpage
                headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'} # Set a user agent
                response = requests.get(url, timeout=10, headers=headers) # Make HTTP GET request with timeout
                response.raise_for_status() # Raise an exception for bad status codes (4XX or 5XX)
                soup = BeautifulSoup(response.content, 'html.parser') # Parse HTML content
                # Remove script and style elements as they don't contain readable content
                for script_or_style in soup(["script", "style"]): # Find all script and style tags
                    script_or_style.decompose() # Remove the element from the parsed tree
                # Get text from common content-holding tags, join paragraphs, and clean up
                text_elements = soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'article', 'div']) # Relevant tags
                text = "\n".join(element.get_text(separator=" ", strip=True) for element in text_elements if element.get_text(strip=True)) # Extract and join text
                return text if text else None # Return extracted text or None if empty
    except requests.exceptions.RequestException as e: # Catch network/HTTP errors
        st.error(f"URL Fetch Error: {str(e)}") # Error message
        return None # Return None
    except Exception as e: # Catch other parsing errors
        st.error(f"Content Parsing Error: {str(e)}") # Error message
        return None # Return None

# --- BERT KEYWORD EXTRACTION ---
def get_bert_embeddings(text: str, tokenizer: Any, model: Any, max_length: int = 512) -> Optional[torch.Tensor]:
    """Generates BERT embeddings for the given text."""
    if not text.strip(): return None # Handle empty text
    try: # Start try block
        inputs = tokenizer(text, return_tensors="pt", truncation=True, max_length=max_length, padding=True) # Tokenize text
        with torch.no_grad(): # Disable gradient calculations for inference
            outputs = model(**inputs) # Get model outputs
        # Use mean pooling of the last hidden state for sentence/document embedding
        embeddings = outputs.last_hidden_state.mean(dim=1) # Pool embeddings
        return embeddings # Return the embedding tensor
    except Exception as e: # Catch errors during embedding generation
        st.warning(f"BERT embedding generation failed: {e}") # Show warning
        return None # Return None

def extract_keywords_with_bert(text: str, models: dict, top_n: int = 10) -> List[str]:
    """Extracts keywords using BERT embeddings by finding terms most similar to the document embedding."""
    if not text or not text.strip(): return [] # Handle empty text
    
    bert_tokenizer = models['bert_tokenizer'] # Get BERT tokenizer from loaded models
    bert_model = models['bert_model'] # Get BERT model from loaded models
    nlp = models['nlp'] # Get spaCy model

    try: # Start try block
        # Get embedding for the whole document
        doc_embedding = get_bert_embeddings(text, bert_tokenizer, bert_model) # Get document embedding
        if doc_embedding is None: return [] # Return empty if document embedding failed

        # Use spaCy to get candidate phrases (noun chunks and significant terms)
        doc_spacy = nlp(text[:50000]) # Process with spaCy (limit length for performance)
        candidate_phrases = list(set([chunk.text.lower() for chunk in doc_spacy.noun_chunks if len(chunk.text.split()) <= 3 and len(chunk.text) > 3])) # Get noun chunks
        candidate_phrases.extend(list(set([token.lemma_.lower() for token in doc_spacy if token.pos_ in ['NOUN', 'PROPN', 'ADJ'] and not token.is_stop and not token.is_punct and len(token.lemma_) > 3]))) # Add other significant tokens
        candidate_phrases = list(set(candidate_phrases)) # Ensure uniqueness

        if not candidate_phrases:
            return [] # Return empty if no candidates

        phrase_embeddings = [] # List to store phrase embeddings
        valid_phrases = [] # List to store phrases for which embeddings were successful
        for phrase in candidate_phrases: # Iterate through candidate phrases
            phrase_emb = get_bert_embeddings(phrase, bert_tokenizer, bert_model) # Get embedding for each phrase
            if phrase_emb is not None: # If embedding is successful
                phrase_embeddings.append(phrase_emb) # Add to list
                valid_phrases.append(phrase) # Add phrase to valid list

        if not phrase_embeddings: return [] # Return if no phrase embeddings

        # Concatenate phrase embeddings into a single tensor
        phrase_embeddings_tensor = torch.cat(phrase_embeddings, dim=0) # Combine tensors

        # Calculate cosine similarity between document embedding and phrase embeddings
        similarities = cosine_similarity(doc_embedding.cpu().numpy(), phrase_embeddings_tensor.cpu().numpy())[0] # Calculate similarities

        # Get top N phrases
        sorted_indices = np.argsort(similarities)[::-1] # Sort by similarity
        top_keywords = [valid_phrases[i] for i in sorted_indices[:top_n]] # Get top N keywords

        return top_keywords # Return extracted keywords
    except Exception as e: # Catch errors during keyword extraction
        st.warning(f"BERT keyword extraction failed: {e}") # Show warning
        return [] # Return empty list on error

# --- SUMMARY GENERATION ---
def generate_summary(models: dict, text: str, audience: str, max_input_length_t5: int = 1024, max_input_length_bart: int = 1024) -> str:
    """Generate audience-specific summaries using appropriate models."""
    if not text:
        return "Cannot generate summary from empty text." # Handle empty input

    try: # Start try block for summary generation
        # Child-friendly summary using BART model
        if audience == "child": # Check if audience is child
            # Truncate text for child summarizer if it's too long (BART has token limits, typically 1024)
            processed_text_for_child = preprocess_text(text, "child")[:max_input_length_bart * 4] # Heuristic for character limit based on token limit
            if not processed_text_for_child.strip():
                return "Not enough content for a child-friendly summary after preprocessing." # Check if text remains

            summary_output = models['child_summarizer']( # Use child summarizer pipeline
                processed_text_for_child, # Input processed text
                max_length=150,  # Shorter summary for children
                min_length=30, # Minimum length for summary
                do_sample=True,  # Use sampling for more creative output
                temperature=0.8, # Control randomness (higher is more random)
                top_p=0.95       # Nucleus sampling parameter
            )
            summary = summary_output[0]['summary_text'] # Extract summary text
            # Simple analogies to make it more relatable
            analogies = [ # List of analogies
                "like building with colorful blocks",
                "similar to how a tiny seed grows into a big plant",
                "just like sharing your favorite toys with a friend",
                "comparable to all the different colors in a rainbow",
                "like different animals in a zoo, each special in its own way"
            ]
            # Format the child summary with an analogy
            return f"Imagine this story is about {summary.lower()}. It's a bit {random.choice(analogies)}." # Return formatted summary

        # Summaries for other audiences (student, researcher, expert) using T5 model
        else: # For other audiences
            # Define prompts for T5 based on the target audience
            prompts = { # Dictionary of prompts
                "student": "summarize for a high school student, focusing on main ideas and concepts: ", # Prompt for student
                "researcher": "summarize for a researcher, highlighting key findings, methodology, and implications: ", # Prompt for researcher
                "expert": "provide a detailed technical summary for an expert in the field, including nuances: " # Prompt for expert
            }
            # T5 typically has a 512 or 1024 token limit for models like t5-base. Chunking by characters.
            # Average token is ~4 chars. Chunk size 3000 chars ~ 750 tokens.
            # Model will truncate to its max_length (e.g., 1024 tokens for t5-base) during encoding.
            chunk_size = 3000 # Characters per chunk for processing long texts
            chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)] # Split text into chunks
            summaries = [] # List to store summaries of each chunk

            # Limit to processing a few chunks to manage time and resources, especially with t5-base
            num_chunks_to_process = min(len(chunks), 2) # Process up to 2 chunks for t5-base to manage time
            with st.spinner(f"Summarizing {num_chunks_to_process} chunk(s) for {audience}..."): # Show spinner for chunk processing
                for chunk_idx in range(num_chunks_to_process): # Iterate through chunks
                    chunk = chunks[chunk_idx] # Get current chunk
                    # Encode the prompt + chunk for T5
                    inputs = models['t5_tokenizer'].encode( # Use T5 tokenizer
                        prompts[audience] + chunk, # Prepend audience-specific prompt
                        return_tensors="pt", # Return PyTorch tensors
                        max_length=max_input_length_t5, # Max tokens for the model input
                        truncation=True # Truncate if longer than max_length
                    ).to(models['t5_model'].device) # Move tensors to the same device as the model

                    # Configure generation parameters based on audience
                    gen_kwargs = { # Dictionary for generation arguments
                        "max_length": 300 if audience == "student" else 500, # Max length of the generated summary
                        "min_length": 75 if audience == "student" else 150,  # Min length of the summary
                        "num_beams": 4,       # Beam search for better quality
                        "early_stopping": True, # Stop when beams converge
                        "temperature": 0.7 if audience == "student" else 0.6 # Control randomness
                    }
                    # Specific parameters for expert summaries
                    if audience == "expert": # If audience is expert
                        gen_kwargs.update({ # Update generation arguments
                            "length_penalty": 2.0, # Encourage longer, more detailed summaries
                            "no_repeat_ngram_size": 3 # Avoid repeating 3-grams
                        })

                    # Generate summary using T5 model
                    outputs = models['t5_model'].generate(inputs, **gen_kwargs) # Generate summary
                    # Decode the generated summary
                    summaries.append(models['t5_tokenizer'].decode(outputs[0], skip_special_tokens=True)) # Decode and add to list

            full_summary = "\n\n[CHUNK BREAK]\n\n".join(summaries) # Join chunk summaries if multiple chunks were processed

            # Add key findings (top TF-IDF terms) for researcher summaries
            if audience == "researcher" and text: # If audience is researcher and text exists
                try: # Try to extract TF-IDF terms
                    vectorizer = TfidfVectorizer(stop_words='english', max_features=1000) # Initialize TF-IDF vectorizer
                    tfidf = vectorizer.fit_transform([text]) # Calculate TF-IDF
                    features = vectorizer.get_feature_names_out() # Get feature names (terms)
                    # Get top 5 terms based on TF-IDF scores
                    top_terms_indices = np.argsort(tfidf.toarray().sum(axis=0))[-5:][::-1] # Get indices of top terms
                    top_terms = [features[i] for i in top_terms_indices] # Get the terms themselves
                    full_summary += "\n\n**KEY TERMS (TF-IDF):**\n- " + "\n- ".join(top_terms) # Append to summary
                except Exception as tfidf_e: # Catch TF-IDF errors
                    st.warning(f"Could not extract TF-IDF key terms for researcher summary: {tfidf_e}") # Show warning

            return full_summary # Return the complete summary
    except Exception as e: # Catch any summary generation error
        st.error(f"Summary Generation Failed: {str(e)}") # Display error message
        return "Could not generate summary due to an error." # Return error message

# --- STUDY MATERIAL GENERATION ---
def generate_study_materials(text: str, nlp_model: Any) -> Optional[Dict[str, List[str]]]:
    """Create comprehensive study materials from text using spaCy."""
    if not text or not text.strip(): return None # Return None if text is empty
    doc = nlp_model(text[:50000]) # Process with spaCy (limit length for performance)
    sentences = list(doc.sents) # Get sentences from the spaCy doc
    if not sentences:
        return None # Return None if no sentences found

    # Extract key phrases (noun chunks, prioritizing those with adjectives/numbers)
    key_phrases = [] # Initialize list for key phrases
    for chunk in doc.noun_chunks: # Iterate through noun chunks identified by spaCy
        # Prioritize phrases that are descriptive or quantitative
        if any(token.pos_ in ('ADJ', 'NUM') for token in chunk) or len(chunk.text.split()) > 1: # Check for adjectives, numbers, or multi-word phrases
            # Basic filter for too short or too common phrases
            if len(chunk.text) > 3 and chunk.text.lower() not in ["this study", "the paper", "the author", "the report"]: # Filter out short/generic phrases
                key_phrases.append(chunk.text.strip()) # Add cleaned phrase
    key_phrases = list(dict.fromkeys(key_phrases)) # Remove duplicates while preserving order
    key_phrases = [phrase for phrase in key_phrases if len(phrase.split()) <= 5][:15] # Limit phrase length and total count

    if not key_phrases: return {"flashcards": ["No key phrases found to generate flashcards."], "questions": [], "diagrams": [], "plan": [], "key_phrases": []} # Return default if no phrases

    # Generate Flashcards with context
    flashcards = [] # Initialize list for flashcards
    for phrase in key_phrases[:10]: # Limit to 10 flashcards
        context_sentences = [s.text for s in sentences if phrase.lower() in s.text.lower()] # Find sentences containing the phrase
        context = random.choice(context_sentences) if context_sentences else f"An important concept related to {phrase}." # Choose a random context sentence
        # Ensure context is not overly long for a flashcard
        context = (context[:200] + '...') if len(context) > 200 else context # Truncate long context
        flashcards.append(f"**Q: What is '{phrase}'?**\nA: *Context:* {context}") # Format flashcard

    # Generate Practice Questions
    questions = [] # Initialize list for questions
    if len(key_phrases) > 1: # If multiple key phrases are available
        for i in range(min(5, len(key_phrases) - 1)): # Limit to 5 questions
            # Create varied question types
            q_type = random.choice(["relate", "significance", "example"]) # Randomly select question type
            if q_type == "relate" and i + 1 < len(key_phrases): # For "relate" type
                questions.append( # Add relation question
                    f"**Q: How might '{key_phrases[i]}' be related to '{key_phrases[i+1]}'?**\n"
                    f"A: *Hint:* Consider their roles, definitions, or how one might influence the other."
                )
            elif q_type == "significance": # For "significance" type
                questions.append( # Add significance question
                    f"**Q: What is the potential significance of '{key_phrases[i]}'?**\n"
                    f"A: *Hint:* Think about its impact or importance in the broader context."
                )
            else: # "example" type as default
                questions.append( # Add example/application question
                    f"**Q: Can you think of an example or application of '{key_phrases[i]}'?**\n"
                    f"A: *Hint:* Try to connect it to a real-world scenario or a concept you already know."
                )
    elif key_phrases: # If only one key phrase
        questions.append(f"**Q: Elaborate on '{key_phrases[0]}'.**\nA: *Hint:* Provide more details or explain its meaning.") # Ask for elaboration

    # Suggest Diagrams based on key phrases
    diagrams = [] # Initialize list for diagram suggestions
    if len(key_phrases) >= 3: # If enough phrases for a concept map
        diagrams.append(f"**Concept Map:** Visually connect '{key_phrases[0]}', '{key_phrases[1]}', and '{key_phrases[2]}' to show their relationships.")
    if len(key_phrases) >= 1: # For timeline suggestion
        diagrams.append(f"**Timeline (if applicable):** If '{key_phrases[0]}' involves a process or history, create a timeline.")
    if len(key_phrases) >= 2: # For Venn diagram suggestion
        diagrams.append(f"**Venn Diagram:** Compare and contrast '{key_phrases[0]}' and '{key_phrases[1]}'.")
    diagrams.append("**Flowchart:** If the text describes a process, try to map it out.") # General flowchart suggestion

    # Generate a simple Study Plan
    study_plan = [] # Initialize list for study plan items
    for i, phrase in enumerate(key_phrases[:min(3, len(key_phrases))]): # Plan for first few key phrases
        study_plan.append(f"**Day {i+1}:** Focus on understanding '{phrase}'. Review flashcards and related text sections.")
    if len(key_phrases) > 3: # If more key phrases exist
        study_plan.append(f"**Day {min(4, len(key_phrases))}:** Review all key phrases and attempt practice questions.")

    return { # Return dictionary of generated study materials
        "flashcards": flashcards, # List of flashcards
        "questions": questions, # List of questions
        "diagrams": diagrams, # List of diagram suggestions
        "plan": study_plan, # List of study plan items
        "key_phrases": key_phrases # List of identified key phrases
    }

# --- VISUALIZATION UTILITIES ---
def create_wordcloud(text: str) -> Optional[plt.Figure]:
    """Generate a word cloud visualization."""
    if not text or not text.strip(): return None # Return None if no text
    try: # Start try block for word cloud generation
        wordcloud_generator = WordCloud( # Initialize WordCloud object
            width=800, # Width of the generated image
            height=400, # Height of the image
            background_color='white', # Background color
            colormap='viridis', # A visually appealing colormap
            stopwords=nltk.corpus.stopwords.words('english'), # Use NLTK's English stopwords
            collocations=False # Avoid showing collocations (pairs of words)
        ).generate(text) # Generate word cloud from text
        fig, ax = plt.subplots(figsize=(12, 6)) # Create matplotlib figure and axes
        ax.imshow(wordcloud_generator, interpolation='bilinear') # Display word cloud image
        ax.axis("off") # Hide axes for a cleaner look
        plt.tight_layout(pad=0) # Adjust layout to prevent clipping
        return fig # Return the figure object
    except Exception as e: # Catch errors during word cloud generation
        st.warning(f"Word cloud generation failed: {e}") # Show warning
        return None # Return None on error

def create_entity_network(entities: List[Tuple[str, str]], nlp_doc: Any) -> Optional[go.Figure]:
    """Create an interactive entity relationship network using Plotly and NetworkX."""
    if not entities: return None # Return None if no entities are provided
    G = nx.Graph() # Initialize a NetworkX graph
    # Add nodes with type grouping and size based on entity length
    unique_entities = list(dict.fromkeys(entities))[:50] # Limit for performance, ensure uniqueness
    for entity_text, label in unique_entities: # Iterate through unique entities
        G.add_node(entity_text, type=label, size=8 + len(entity_text) / 2.5, label=label) # Add node with attributes

    # Add edges between entities that co-occur in the same sentence
    entity_map = {text: label for text, label in unique_entities} # Create a dictionary for faster lookup of entity types
    sentences = list(nlp_doc.sents) # Get sentences from spaCy doc
    for sentence in sentences: # Iterate through sentences
        sent_entities = [ent.text for ent in sentence.ents if ent.text in entity_map] # Get entities present in the current sentence
        # Add edges between co-occurring entities in the sentence
        for i in range(len(sent_entities)): # First entity in pair
            for j in range(i + 1, len(sent_entities)): # Second entity in pair
                if G.has_node(sent_entities[i]) and G.has_node(sent_entities[j]): # Check if both nodes exist
                    if G.has_edge(sent_entities[i], sent_entities[j]): # If edge already exists
                        G[sent_entities[i]][sent_entities[j]]['weight'] += 0.2 # Increase weight for co-occurrence
                    else: # If edge does not exist
                        G.add_edge(sent_entities[i], sent_entities[j], weight=0.5) # Add new edge with initial weight

    if not G.nodes: return None # If no nodes were added (e.g., due to filtering)

    pos = nx.spring_layout(G, k=0.8, iterations=50) # Position nodes using spring layout algorithm for better visualization

    # Create Plotly figure components for edges
    edge_x, edge_y = [], [] # Initialize lists for edge coordinates
    edge_weights = [] # Initialize list for edge weights (for line thickness)
    for edge in G.edges(data=True): # Iterate through edges with their data
        x0, y0 = pos[edge[0]] # Coordinates of the source node
        x1, y1 = pos[edge[1]] # Coordinates of the target node
        edge_x.extend([x0, x1, None]) # X coordinates for edges (None separates lines)
        edge_y.extend([y0, y1, None]) # Y coordinates for edges
        edge_weights.append(edge[2].get('weight', 0.5)) # Get edge weight, default to 0.5

    edge_trace = go.Scatter( # Create Scatter trace for edges
        x=edge_x, y=edge_y, # Edge coordinates
        line=dict(width=0.7, color='#888'), # Edge line style (can be varied by weight if desired)
        hoverinfo='none', # No hover info for edges themselves
        mode='lines' # Draw lines for edges
    )

    # Create Plotly figure components for nodes
    node_x, node_y, node_text, node_color, node_size_viz = [], [], [], [], [] # Initialize lists for node properties
    entity_types = list(set(nx.get_node_attributes(G, 'type').values())) # Get unique entity types for coloring
    # Using a qualitative color scale like 'Accent' or 'Paired' from Plotly or Matplotlib
    # For simplicity, using a predefined list or a cyclical map if many types
    try: # Try to get a colormap
        color_map_types = plt.cm.get_cmap('Accent', len(entity_types) if entity_types else 1) # Get Accent colormap
    except Exception: # Fallback if cmap fails
        def color_map_types(x):
            return (random.random(), random.random(), random.random(), 1.0) # Random color fallback

    for node_idx, node in enumerate(G.nodes()): # Iterate through nodes
        x, y = pos[node] # Node position
        node_x.append(x) # X coordinate for node
        node_y.append(y) # Y coordinate for node
        node_info = G.nodes[node] # Get node attributes
        type_info = node_info.get('type', 'Unknown') # Get entity type
        node_text.append(f"{node}<br>Type: {type_info}") # Text for hover info (node name and type)
        # Assign color based on entity type
        try: # Try to assign color from map
            type_index = entity_types.index(type_info) if type_info in entity_types else -1 # Get index of type
            if type_index != -1: # If type is known
                rgb_color = color_map_types(type_index) # Get color from map
                node_color.append(f'rgb({int(rgb_color[0]*255)},{int(rgb_color[1]*255)},{int(rgb_color[2]*255)})') # Convert to RGB string
            else: # Fallback color for unknown types
                node_color.append('rgb(128,128,128)') # Grey
        except ValueError: # Fallback color if type_info not in entity_types (should be handled by above)
            node_color.append('rgb(128,128,128)') # Grey for unknown types
        node_size_viz.append(node_info.get('size', 10)) # Visual size of node

    node_trace = go.Scatter( # Create Scatter trace for nodes
        x=node_x, y=node_y, # Node coordinates
        mode='markers+text', # Show markers and text for nodes
        text=[node for node in G.nodes()], # Display node name as text near the marker
        textposition='top center', # Position of the text label
        hoverinfo='text', # Show text from node_text on hover
        hovertext=node_text, # Custom hover text
        marker=dict( # Marker style
            showscale=False, # No color scale legend
            colorscale='Rainbow', # Default colorscale if not using specific node_color (overridden by color array)
            color=node_color, # Assign colors based on type
            size=node_size_viz, # Node sizes
            line_width=1.5, # Border width of markers
            line_color='black' # Border color of markers
        )
    )

    # Create the final Plotly figure
    fig = go.Figure( # Initialize Figure object
        data=[edge_trace, node_trace], # Combine edge and node traces
        layout=go.Layout( # Figure layout configuration
            title='Interactive Entity Relationship Network', # Figure title
            titlefont_size=16, # Title font size
            showlegend=False, # Hide legend
            hovermode='closest', # Hover mode for interactivity
            margin=dict(b=10, l=5, r=5, t=30), # Figure margins
            height=600, # Figure height
            xaxis=dict(showgrid=False, zeroline=False, showticklabels=False), # X-axis style (hidden)
            yaxis=dict(showgrid=False, zeroline=False, showticklabels=False)  # Y-axis style (hidden)
        )
    )
    return fig # Return the Plotly figure object

# --- MAIN APPLICATION LOGIC ---
# Load models once at the start
# This is a global call, and the result is used throughout the app.
MODELS = load_models() # It's crucial this completes successfully.

def display_analysis_results(text_to_analyze: str, source_name: str = "Uploaded Content"):
    """
    A centralized function to display summaries, visualizations, BERT keywords, and study materials
    for the provided text.
    """
    st.subheader(f"Analysis Results for: {source_name}") # Display subheader for the source
    # Preprocess text for different audiences/models
    child_text_processed = preprocess_text(text_to_analyze, "child") # Preprocess for child audience
    standard_text_processed = preprocess_text(text_to_analyze) # Standard preprocessing for other uses

    # Tabbed interface for different analysis outputs
    analysis_tabs = st.tabs(["📚 Summaries", "📊 Visualizations", "🔑 BERT Keywords", "📝 Study Aids"]) # Define tabs

    with analysis_tabs[0]: # Summaries Tab
        st.markdown("#### Audience-Specific Summaries") # Section title
        st.info("Summaries are generated for different understanding levels. Select an audience.") # Info message
        # Audience selection for summary
        summary_audience = st.selectbox( # Create a select box for audience choice
            "Choose summary audience:", # Label for select box
            ["child", "student", "researcher", "expert"], # Options for audience
            key=f"summary_audience_{source_name.replace(' ','_')}" # Unique key for state management
        )
        if st.button(f"Generate {summary_audience.capitalize()} Summary", key=f"gen_summary_btn_{source_name.replace(' ','_')}"): # Button to generate summary
            with st.spinner(f"Generating {summary_audience} summary... Please wait."): # Show spinner during generation
                text_for_summary = child_text_processed if summary_audience == "child" else standard_text_processed # Select appropriate processed text
                if not text_for_summary.strip(): # Check if text is empty after processing
                    st.warning("The text is too short or empty after preprocessing for this summary type.") # Show warning
                else: # If text is available
                    summary = generate_summary(MODELS, text_for_summary, summary_audience) # Generate summary
                    summary_container_class = "child-mode" if summary_audience == "child" else "summary-box" # Choose CSS class based on audience
                    st.markdown(f'<div class="{summary_container_class}">{summary}</div>', unsafe_allow_html=True) # Display summary in styled div

    with analysis_tabs[1]: # Visualizations Tab
        st.markdown("#### Text Visualizations") # Section title
        st.write("Visual tools to explore the text's key themes and entities.") # Description
        
        # Word Cloud Generation
        st.markdown("##### Word Cloud") # Sub-section title
        if st.button("Generate Word Cloud", key=f"wc_btn_{source_name.replace(' ','_')}"): # Button for word cloud
            with st.spinner("Creating word cloud..."): # Show spinner
                wordcloud_fig = create_wordcloud(standard_text_processed) # Create word cloud
                if wordcloud_fig: # If figure is generated
                    st.pyplot(wordcloud_fig) # Display matplotlib figure
                    st.caption("Most frequent words in the text (excluding common stop words).") # Caption for context
                else: # If word cloud failed
                    st.info("Not enough text to generate a word cloud or an error occurred.") # Info message

        # Entity Network Visualization
        st.markdown("##### Entity Network") # Sub-section title
        if st.button("Generate Entity Network", key=f"en_btn_{source_name.replace(' ','_')}"): # Button for entity network
            with st.spinner("Analyzing entities and building network..."): # Show spinner
                # Limit input to spaCy for performance on very large texts for NER network
                doc_for_network = MODELS['nlp'](standard_text_processed[:50000]) # Process with spaCy, limit length
                entities = [(ent.text, ent.label_) for ent in doc_for_network.ents] # Extract entities
                if entities: # If entities are found
                    network_fig = create_entity_network(entities, doc_for_network) # Create network figure
                    if network_fig: # If figure is generated
                        st.plotly_chart(network_fig, use_container_width=True) # Display Plotly chart
                        st.caption("Interactive network of named entities and their co-occurrences.") # Caption
                    else: # If network generation failed
                        st.info("Could not generate entity network (e.g., too few entities or layout error).") # Info message
                else: # If no entities found
                    st.info("No named entities found to create a network.") # Info message
        
        # Text Statistics
        st.markdown("##### Text Statistics") # Sub-section title
        if standard_text_processed: # If text is available
            # Limit input to spaCy for performance on very large texts for stats
            doc_stats = MODELS['nlp'](standard_text_processed[:100000]) # Process with spaCy, limit length
            num_sentences = len(list(doc_stats.sents)) # Count sentences
            num_unique_tokens = len(set(token.lemma_.lower() for token in doc_stats if not token.is_stop and not token.is_punct)) # Count unique non-stopword lemmas
            num_words = len([token for token in doc_stats if not token.is_punct]) # Count words (excluding punctuation)
            avg_sent_len = num_words / num_sentences if num_sentences > 0 else 0 # Calculate average sentence length
            
            col1, col2, col3 = st.columns(3) # Create columns for metrics
            with col1: st.metric("Sentences", f"{num_sentences:,}") # Display sentence count
            with col2: st.metric("Unique Terms (Lemmas)", f"{num_unique_tokens:,}") # Display unique term count
            with col3: st.metric("Avg. Sentence Length", f"{avg_sent_len:.1f} words") # Display average sentence length
        else: # If no text for statistics
            st.info("Not enough text for statistics.") # Info message
            
    with analysis_tabs[2]: # BERT Keywords Tab
        st.markdown("#### 🔑 BERT-powered Insights") # Section title
        st.write("Extract key terms using BERT embeddings to capture semantic relevance.") # Description

        if st.button("Extract BERT Keywords", key=f"bert_kw_btn_{source_name.replace(' ','_')}"): # Button for BERT keywords
            with st.spinner("Extracting keywords with BERT... This might take a moment."): # Show spinner
                bert_keywords = extract_keywords_with_bert(standard_text_processed, MODELS, top_n=10) # Extract keywords
                if bert_keywords: # If keywords are found
                    st.markdown("##### Top 10 Keywords (via BERT):") # Sub-title
                    st.info(", ".join([f"`{kw}`" for kw in bert_keywords])) # Display keywords
                    st.caption("These keywords are identified based on their semantic similarity to the overall document context using BERT embeddings.") # Caption
                else: # If no keywords found or error
                    st.warning("Could not extract keywords using BERT. The text might be too short or an error occurred.") # Warning

    with analysis_tabs[3]: # Study Aids Tab (was index 2, now 3)
        st.markdown("#### Study Aids Generator") # Section title
        st.write("Create flashcards, practice questions, and more from the text.") # Description
        if st.button("Generate Study Materials", key=f"study_aid_btn_{source_name.replace(' ','_')}"): # Button for study aids
            with st.spinner("Crafting study materials..."): # Show spinner
                materials = generate_study_materials(standard_text_processed, MODELS['nlp']) # Generate materials
                if materials and materials.get("key_phrases"): # Check if materials and key phrases exist
                    st.markdown("##### 🔑 Key Phrases Identified (spaCy):") # Sub-title
                    st.info(", ".join(materials["key_phrases"])) # Display key phrases

                    st.markdown("##### 📇 Flashcards (Q&A)") # Sub-title
                    for i, card in enumerate(materials["flashcards"]): # Iterate through flashcards
                        question, answer_section = card.split("\nA: ", 1) # Split Q and A
                        with st.expander(f"Flashcard {i+1}: {question.replace('**Q:','').replace('**','').strip()}", expanded=False): # Create expander for each card
                            st.markdown(f"A: {answer_section}") # Display answer

                    st.markdown("##### ❓ Practice Questions") # Sub-title
                    for q_text in materials["questions"]: # Iterate through questions
                        st.markdown(f'<div class="flashcard">{q_text}</div>', unsafe_allow_html=True) # Display question in styled div

                    st.markdown("##### 📊 Recommended Diagram Types") # Sub-title
                    for diag_suggestion in materials["diagrams"]: # Iterate through diagram suggestions
                        st.markdown(f"- {diag_suggestion}") # Display suggestion

                    st.markdown("##### 📅 Suggested Study Plan") # Sub-title
                    for plan_item in materials["plan"]: # Iterate through plan items
                        st.markdown(f"- {plan_item}") # Display plan item
                    
                    # Prepare content for download
                    study_content_str = f"KEY PHRASES (spaCy):\n{', '.join(materials['key_phrases'])}\n\n" # Add key phrases
                    study_content_str += "FLASHCARDS:\n" + "\n\n".join(materials["flashcards"]) + "\n\n" # Add flashcards
                    study_content_str += "PRACTICE QUESTIONS:\n" + "\n\n".join(materials["questions"]) + "\n\n" # Add questions
                    study_content_str += "DIAGRAM SUGGESTIONS:\n" + "\n".join(materials["diagrams"]) + "\n\n" # Add diagram suggestions
                    study_content_str += "STUDY PLAN:\n" + "\n".join(materials["plan"]) # Add study plan
                    
                    st.download_button( # Create download button
                        label="📥 Download All Study Materials (TXT)", # Button label
                        data=study_content_str, # Data to download
                        file_name=f"study_materials_{source_name.replace(' ','_')}.txt", # Filename
                        mime="text/plain" # Mime type
                    )
                elif materials and materials.get("flashcards") and "No key phrases" in materials["flashcards"][0]: # If no key phrases were found explicitly
                    st.warning("No key phrases were found in the text by spaCy to generate detailed study materials.") # Show warning
                else: # General failure to generate materials
                    st.warning("Could not generate study materials. The text might be too short or lack distinct concepts for spaCy analysis.") # Show warning


def main_application():
    """Main function to run the Streamlit application."""
    from PIL import Image # Import PIL for image handling, consider moving to top if used more widely
    
    # --- Sidebar Navigation ---
    try: # Try to load an image for the sidebar
        # Replace 'QR3XoLs.jpeg' with a valid path to your logo or remove if not needed.
        # For robustness, check if file exists or use a default.
        # For this example, assuming it might not be present, we'll wrap it.
        # sidebar_logo_path = 'QR3XoLs.jpeg' # Path to logo
        # if os.path.exists(sidebar_logo_path): # Check if logo exists
        # st.sidebar.image(Image.open(sidebar_logo_path), width=100) # Display logo
        # st.sidebar.caption("Your App Name/Logo Caption") # Caption for logo
        pass # Commenting out the logo part as the file isn't provided with the script
    except FileNotFoundError: # Handle if image not found
        st.sidebar.warning("Logo image not found.") # Warning message
    except Exception as img_e: # Handle other image errors
        st.sidebar.warning(f"Could not load logo: {img_e}") # Warning message

        # --- Sidebar Navigation ---
    st.sidebar.image(Image.open('QR3XoLs.jpeg'), width=100)  # Placeholder logo, width=100
    st.sidebar.caption("pokeman")

    st.sidebar.title("🌟 AI Analyzer Suite 🌟") # Title for the sidebar
    st.sidebar.markdown("---") # Divider in the sidebar

    # Define application modes for navigation using a radio button
    app_mode = st.sidebar.radio( # Create radio button for mode selection
        "Choose a Tool:", # Label for radio buttons
        ["🏠 Home", "💬 Text Input Analysis", "📄 Document File Analysis", "🎤 Media File Analysis (Audio/Video)", "🔗 URL Content Analysis", "💡 About"], # Options
        help="Select the type of content you want to analyze." # Help tooltip
    )
    st.sidebar.markdown("---") # Another divider
    st.sidebar.info("Upload content or paste text, then explore summaries, visualizations, and study aids.") # Informational message

    # --- Page Content Based on Navigation ---
    if app_mode == "🏠 Home": # If Home is selected
        st.title("Welcome to the AI Text & Media Analyzer Suite! 👋") # Main title for home page
        st.markdown("""
            <div class="content-card">
                <p>Unlock insights from your text and media with our powerful AI-driven tools. This suite is designed for students, researchers, educators, and anyone looking to understand content deeply and efficiently.</p>
                <h4>🚀 <strong>What You Can Do:</strong></h4>
                <ul>
                    <li>Paste text directly for quick analysis (<b>Text Input Analysis</b>).</li>
                    <li>Upload documents (PDF, DOCX, PPTX, TXT) to extract and analyze text (<b>Document File Analysis</b>).</li>
                    <li>Analyze spoken content from audio (MP3, WAV, M4A) or video (MP4, MOV, AVI) files (<b>Media File Analysis</b>).</li>
                    <li>Fetch and analyze text from web pages or YouTube transcripts (<b>URL Content Analysis</b>).</li>
                </ul>
                <h4>🎯 <strong>For Each Input, Get:</strong></h4>
                <ul>
                    <li>👶 <strong>Child-friendly</strong> explanations.</li>
                    <li>🎓 <strong>Student-focused</strong> summaries.</li>
                    <li>🔬 <strong>Researcher-oriented</strong> analyses with key terms (TF-IDF & BERT).</li>
                    <li>🧠 <strong>Expert-level</strong> detailed summaries.</li>
                    <li>📊 <strong>Visualizations</strong> like word clouds and entity networks.</li>
                    <li>🔑 <strong>BERT Keywords</strong> for semantic term discovery.</li>
                    <li>📝 <strong>Study Aids</strong> including flashcards, practice questions, and study plans (via spaCy).</li>
                </ul>
                <p>Select an analysis tool from the sidebar to begin!</p>
            </div>
        """, unsafe_allow_html=True) # Display introductory content using HTML for styling
        st.balloons() # Fun balloons animation on home page

    elif app_mode == "💬 Text Input Analysis": # If Text Input Analysis is selected
        st.header("💬 Direct Text Input Analysis") # Page header
        st.markdown("<div class='content-card'>Enter or paste your text below for a comprehensive analysis.</div>", unsafe_allow_html=True) # Info card
        # Text area for user input
        text_input = st.text_area("Paste your text here:", height=250, placeholder="Type or paste your content...", key="text_input_main", help="Enter the text you want to analyze.") # Text input field
        if text_input and text_input.strip(): # If text is entered and not just whitespace
            display_analysis_results(text_input, "Pasted Text") # Call function to display results
        elif text_input and not text_input.strip(): # If only whitespace is entered
            st.info("Please paste some actual text, not just spaces.") # Inform user
        else: # If no text is entered (on initial load)
            st.info("Please paste some text in the area above to enable analysis.") # Prompt user

    elif app_mode == "📄 Document File Analysis": # If Document File Analysis is selected
        st.header("📄 Document File Analysis") # Page header
        st.markdown("<div class='content-card'>Upload your documents (PDF, DOCX, PPTX, TXT) to extract text and generate insights.</div>", unsafe_allow_html=True) # Info card
        # File uploader for document files
        uploaded_file = st.file_uploader( # Create file uploader widget
            "Upload your document:", # Label
            type=["pdf", "docx", "pptx", "txt"], # Supported file types
            key="doc_uploader", # Unique key
            help="Supports PDF, Word (DOCX), PowerPoint (PPTX), and Text (TXT) files." # Help tooltip
        )
        if uploaded_file: # If a file is uploaded
            with st.spinner(f"Processing '{uploaded_file.name}'... Please wait."): # Show spinner
                # Save to a temporary file to pass its path to extraction functions
                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file: # Create temp file
                    tmp_file.write(uploaded_file.getvalue()) # Write uploaded file content to temp file
                    tmp_file_path = tmp_file.name # Get path of the temp file
                
                raw_text = extract_text_from_file(tmp_file_path, uploaded_file.type) # Extract text
                
                try: # Try to delete the temporary file
                    os.unlink(tmp_file_path) # Delete the temporary file after use
                except Exception as e_unlink: # Catch error if deletion fails
                    st.warning(f"Could not delete temporary file {tmp_file_path}: {e_unlink}") # Show warning

                if raw_text and raw_text.strip(): # If text extraction is successful and not empty
                    st.success(f"Successfully extracted text from '{uploaded_file.name}'.") # Success message
                    with st.expander("View Extracted Text", expanded=False): # Expander to view text
                        st.text_area("Extracted Text:", value=raw_text, height=300, disabled=True, key=f"extracted_text_doc_{uploaded_file.name}") # Display extracted text
                    display_analysis_results(raw_text, uploaded_file.name) # Analyze and display results
                elif raw_text is None: # If extraction returned None (error handled in function)
                    pass # Error message already shown by extract_text_from_file
                else: # If extracted text is empty
                    st.error(f"Could not extract significant text from '{uploaded_file.name}' or the document is empty.") # Error message
                    
    elif app_mode == "🎤 Media File Analysis (Audio/Video)": # If Media File Analysis is selected
        st.header("🎤 Audio & Video Transcription and Analysis") # Page header
        st.markdown("<div class='content-card'>Upload audio (MP3, WAV, M4A, OGG, FLAC) or video (MP4, MOV, AVI) files. The audio will be transcribed to text for analysis.</div>", unsafe_allow_html=True) # Info card
        # File uploader for audio/video files
        uploaded_media_file = st.file_uploader( # Create file uploader
            "Upload your audio or video file:", # Label
            type=["mp3", "wav", "m4a", "mp4", "mov", "avi", "ogg", "flac"], # Supported media types
            key="media_uploader", # Unique key
            help="Transcription can take time for long files. Supported formats include MP3, WAV, MP4, etc." # Help tooltip
        )
        if uploaded_media_file: # If a media file is uploaded
            with st.spinner(f"Processing media file '{uploaded_media_file.name}'... This may take some time."): # Show spinner
                # Save to a temporary file
                file_extension = os.path.splitext(uploaded_media_file.name)[1] # Get file extension
                with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_media_file: # Create temp file
                    tmp_media_file.write(uploaded_media_file.getvalue()) # Write content
                    tmp_media_file_path = tmp_media_file.name # Get path

                transcribed_text = None # Initialize transcribed_text
                audio_to_transcribe_path = tmp_media_file_path # Path to the audio to be transcribed

                # Determine if it's video or audio
                if file_extension.lower() in ['.mp4', '.mov', '.avi']: # Video file types
                    # Define a path for the extracted audio in the temp directory
                    audio_output_path = os.path.join(tempfile.gettempdir(), f"extracted_audio_{os.path.basename(tmp_media_file_path)}.wav")
                    if extract_audio_from_video(tmp_media_file_path, audio_output_path): # Extract audio
                        audio_to_transcribe_path = audio_output_path # Update path to the extracted audio
                    else: # If audio extraction failed
                        st.error("Failed to extract audio from video.") # Error message
                        audio_to_transcribe_path = None # Set to None so transcription is skipped
                
                if audio_to_transcribe_path: # If there's an audio file to transcribe
                    transcribed_text = transcribe_audio_with_whisper(audio_to_transcribe_path, MODELS['transcription_model']) # Transcribe
                    # Clean up extracted audio if it was created
                    if audio_to_transcribe_path != tmp_media_file_path and os.path.exists(audio_to_transcribe_path):
                        try: os.unlink(audio_to_transcribe_path)
                        except Exception as e_unlink_audio: st.warning(f"Could not delete temp audio {audio_output_path}: {e_unlink_audio}")


                # Clean up uploaded media temp file
                if os.path.exists(tmp_media_file_path): # Check if temp file exists
                    try: os.unlink(tmp_media_file_path) # Delete temp file
                    except Exception as e_unlink_media: st.warning(f"Could not delete temp media {tmp_media_file_path}: {e_unlink_media}")

                if transcribed_text and transcribed_text.strip(): # If transcription successful and not empty
                    st.success(f"Successfully transcribed audio from '{uploaded_media_file.name}'.") # Success message
                    with st.expander("View Transcription", expanded=False): # Expander for transcription
                        st.text_area("Transcribed Text:", value=transcribed_text, height=300, disabled=True, key=f"transcribed_text_{uploaded_media_file.name}") # Display text
                    display_analysis_results(transcribed_text, f"Transcription of {uploaded_media_file.name}") # Analyze
                elif transcribed_text is None and audio_to_transcribe_path : # Explicit check for None if transcription was attempted and failed
                    st.error(f"Transcription failed for '{uploaded_media_file.name}'. See error above if any.") # Error message (Whisper function shows its own error)
                elif not transcribed_text and audio_to_transcribe_path: # Empty transcription but no explicit error
                     st.warning(f"Transcription resulted in empty text for '{uploaded_media_file.name}'. The media might be silent or unclear.") # Warning
    
    elif app_mode == "🔗 URL Content Analysis": # If URL Content Analysis is selected
        st.header("🔗 Web Content & YouTube Transcript Analysis") # Page header
        st.markdown("<div class='content-card'>Enter a URL (webpage or YouTube video) to fetch its text content for analysis.</div>", unsafe_allow_html=True) # Info card
        # Text input for URL
        url_input = st.text_input("Enter URL (e.g., https://example.com or YouTube video link):", key="url_input_main", help="Make sure the URL is accessible.") # URL input field
        if st.button("Fetch and Analyze URL", key="fetch_url_btn"): # Button to fetch and analyze
            if url_input and url_input.strip(): # If URL is provided
                fetched_text = fetch_and_parse_url(url_input) # Fetch and parse URL content
                if fetched_text and fetched_text.strip(): # If successful and not empty
                    st.success("Successfully fetched and parsed content from URL.") # Success message
                    with st.expander("View Fetched Text", expanded=False): # Expander for text
                        st.text_area("Fetched Text:", value=fetched_text, height=300, disabled=True, key=f"fetched_text_url_{url_input}") # Display text
                    # Use domain as source name, or a generic name if parsing fails
                    try: source_display_name = url_input.split('//')[-1].split('/')[0]
                    except: source_display_name = "Web Content"
                    display_analysis_results(fetched_text, f"Content from {source_display_name}") # Analyze
                elif fetched_text is None: # Error handled by fetch_and_parse_url, so just pass
                    pass # Error message already shown by the function
                else: # If fetched text is empty string
                    st.error("Could not fetch significant text content from the URL or the content is empty.") # Error message
            else: # If URL input is empty
                st.warning("Please enter a URL.") # Warning message

    elif app_mode == "💡 About": # If About is selected
        st.title("💡 About the AI Analyzer Suite") # Page title
        st.markdown( # Display information about the application
            """
            <div class="content-card">
                <h4><strong>AI Text & Media Analyzer Suite - Version 3.0 (BERT Enhanced)</strong></h4>
                <p>This application leverages state-of-the-art Natural Language Processing (NLP) models to provide comprehensive text analysis and understanding tools. It's built with Python using the Streamlit framework for the user interface.</p>
                
                <h5><strong>Core Technologies:</strong></h5>
                <ul>
                    <li><strong>Summarization:</strong>
                        <ul>
                            <li>General Purpose: Google's T5 (Text-To-Text Transfer Transformer) - <code>t5-base</code> model.</li>
                            <li>Child-Friendly: Facebook's BART (Bidirectional Auto-Regressive Transformer) - <code>facebook/bart-large-cnn</code> model.</li>
                        </ul>
                    </li>
                    <li><strong>Keyword Extraction & Embeddings:</strong> BERT (Bidirectional Encoder Representations from Transformers) - <code>bert-base-uncased</code> model.</li>
                    <li><strong>NLP Tasks (Entity Recognition, POS Tagging, etc.):</strong> SpaCy - <code>en_core_web_sm</code> model.</li>
                    <li><strong>Audio Transcription:</strong> OpenAI's Whisper - <code>base</code> model.</li>
                    <li><strong>Text Extraction:</strong> Libraries like <code>PyPDF2</code> (PDFs), <code>python-docx</code> (Word), <code>python-pptx</code> (PowerPoint).</li>
                    <li><strong>Web Content:</strong> <code>requests</code> and <code>BeautifulSoup4</code> for webpages, <code>youtube_transcript_api</code> for YouTube.</li>
                    <li><strong>Visualizations:</strong> <code>WordCloud</code>, <code>Matplotlib</code>, <code>Plotly</code>, and <code>NetworkX</code>.</li>
                </ul>

                <h5><strong>Key Features:</strong></h5>
                <ul>
                    <li>Multi-format input: Direct text, various document types, audio/video files, and web URLs.</li>
                    <li>Audience-specific summarization tailored for children, students, researchers, and experts.</li>
                    <li>BERT-powered keyword extraction for semantic understanding.</li>
                    <li>Automated generation of study aids: key phrases (spaCy), flashcards, practice questions, diagram suggestions, and study plans.</li>
                    <li>Interactive visualizations: Word clouds and entity relationship networks.</li>
                    <li>Detailed text statistics.</li>
                </ul>

                <h5><strong>Developed For:</strong></h5>
                <p>Educational purposes, research assistance, content understanding, and general productivity. This tool aims to make complex information more accessible and actionable.</p>
                
                <h5><strong>Disclaimer:</strong></h5>
                <p>AI-generated content can sometimes be inaccurate or incomplete. Always critically evaluate the results. Processing very large files or long audio/video may take time and consume significant resources. Ensure you have the necessary permissions to process any content you upload or link to.</p>
                <hr>
                <p style="text-align:center;"><em>Happy Analyzing!</em></p>
            </div>
            """,
            unsafe_allow_html=True, # Allow HTML for styling
        )

# Entry point of the script: This ensures main_application() runs when the script is executed
if __name__ == "__main__": # Standard Python idiom to check if the script is run directly
    main_application() # Call the main function to start the Streamlit app