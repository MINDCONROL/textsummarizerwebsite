import streamlit as st
import os
import tempfile
import random
from typing import Optional, List, Dict, Tuple
from transformers import T5Tokenizer, T5ForConditionalGeneration, pipeline
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize, RegexpTokenizer
from nltk.probability import FreqDist
import string
import re
from collections import defaultdict
import textwrap
import networkx as nx
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
import matplotlib.pyplot as plt
st.set_page_config(layout="wide", page_title="Ultimate Text Analyzer", page_icon="🧠")
#from streamlit_option_menu import option_menu
choice = st.sidebar.selectbox("Select an option", ["Home","Document","URL","About"])
if choice == "Home":
    st.title("Welcome to the Home Page")
    st.write("This is the home page of the Streamlit app.")
    st.write("You can use this app to summarize text using the T5 model.")
    st.write("Please enter the text you want to summarize in the text area below.")
    text = st.text_area("Enter text to summarize", height=300)
    if st.button("Summarize"):
        from transformers import pipeline
        summarizer = pipeline("summarization", model="t5-base", tokenizer="t5-base", framework="pt")
        summary = summarizer(text, max_length=100, min_length=30, do_sample=False)
        st.write("Summary:")
        st.write(summary[0]['summary_text'])



elif choice == "Document":
   # Configure NLTK and text processing
    nltk.download('punkt')
    nltk.download('stopwords')
    tokenizer = RegexpTokenizer(r'\w+')

# File processing libraries
    from PyPDF2 import PdfReader
    from docx import Document
    from pptx import Presentation

# Initialize AI models with enhanced parameters
@st.cache_resource
def load_models():
    try:
        # Main summarization model
        t5_tokenizer = T5Tokenizer.from_pretrained("t5-small")
        t5_model = T5ForConditionalGeneration.from_pretrained("t5-small")
        
        # Specialized child-friendly summarizer
        child_summarizer = pipeline(
            "summarization", 
            model="facebook/bart-large-cnn",
            tokenizer="facebook/bart-large-cnn",
            device=0 if st.runtime.exists() else -1
        )
        
        return t5_tokenizer, t5_model, child_summarizer
    except Exception as e:
        st.error(f"Model loading failed: {str(e)}")
        st.stop()

t5_tokenizer, t5_model, child_summarizer = load_models()

# ======================
# ENHANCED TEXT PROCESSING
# ======================
def extract_text_from_file(file_path: str, file_type: str) -> Optional[str]:
    """Improved text extraction with format-specific handling"""
    try:
        if file_type == "application/pdf":
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs if para.text)
        
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return "\n\n".join(text)
        
        elif file_type == "text/plain":
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        return None
    except Exception as e:
        st.error(f"Extraction error: {str(e)}")
        return None

def preprocess_text(text: str, mode: str = "default") -> str:
    """Mode-specific text cleaning"""
    text = re.sub(r'\s+', ' ', text).strip()
    
    if mode == "child":
        text = re.sub(r'[^\w\s.,!?]', '', text)  # Keep only simple punctuation
        text = re.sub(r'\b\w{15,}\b', '', text)  # Remove long words
    else:
        text = re.sub(r'\[.*?\]|\(.*?\)', '', text)  # Remove citations
    
    return text

# ======================
# COMPLETE SUMMARY GENERATORS
# ======================
def generate_child_summary(text: str) -> str:
    """Ultra-simple explanation with analogies"""
    try:
        # First generate base summary
        summary = child_summarizer(
            text[:1024],
            max_length=150,
            min_length=30,
            do_sample=True,
            temperature=0.9,
            top_p=0.95
        )[0]['summary_text']
        
        # Add child-friendly analogies
        analogies = [
            "like building with blocks",
            "similar to how plants grow",
            "just like sharing toys",
            "comparable to a rainbow's colors",
            "like different animals in a zoo"
        ]
        
        return f"This is about {summary.lower()}. It works {random.choice(analogies)}."
    
    except Exception as e:
        st.error(f"Child summary failed: {str(e)}")
        return "Couldn't make a simple version. Ask an adult to explain."

def generate_student_summary(text: str) -> str:
    """Balanced explanation with examples"""
    inputs = t5_tokenizer.encode(
        "summarize for high school student: " + text[:3000],
        return_tensors="pt",
        max_length=512,
        truncation=True
    )
    outputs = t5_model.generate(
        inputs,
        max_length=250,
        min_length=100,
        num_beams=4,
        early_stopping=True,
        temperature=0.7
    )
    return t5_tokenizer.decode(outputs[0], skip_special_tokens=True)

def generate_researcher_summary(text: str) -> str:
    """Technical analysis with key findings"""
    inputs = t5_tokenizer.encode(
        "summarize for researcher with key findings: " + text[:3000],
        return_tensors="pt",
        max_length=512,
        truncation=True
    )
    outputs = t5_model.generate(
        inputs,
        max_length=400,
        min_length=200,
        num_beams=4,
        length_penalty=2.0,
        no_repeat_ngram_size=3,
        temperature=0.5
    )
    summary = t5_tokenizer.decode(outputs[0], skip_special_tokens=True)
    
    # Add key findings section
    vectorizer = TfidfVectorizer(stop_words='english')
    tfidf = vectorizer.fit_transform([text])
    features = vectorizer.get_feature_names_out()
    top_terms = features[np.argsort(tfidf.toarray())[0][-5:][::-1]]
    
    return f"{summary}\n\nKEY FINDINGS:\n- " + "\n- ".join(top_terms)

def generate_expert_summary(text: str) -> str:
    """Comprehensive scholarly summary"""
    chunks = [text[i:i+3000] for i in range(0, len(text), 3000)]
    summaries = []
    
    for chunk in chunks[:3]:  # Limit to 3 chunks
        inputs = t5_tokenizer.encode(
            "summarize in detail for expert: " + chunk,
            return_tensors="pt",
            max_length=512,
            truncation=True
        )
        outputs = t5_model.generate(
            inputs,
            max_length=600,
            min_length=300,
            num_beams=4,
            length_penalty=2.5,
            temperature=0.3
        )
        summaries.append(t5_tokenizer.decode(outputs[0], skip_special_tokens=True))
    
    return "\n\n".join(summaries)

# ======================
# ENHANCED STUDY GUIDE
# ======================
def create_study_guide(text: str, key_phrases: List[str]) -> Dict[str, List[str]]:
    """Complete learning package with 5 components"""
    sentences = sent_tokenize(text)
    
    # 1. Flashcards with context
    flashcards = []
    for phrase in key_phrases[:10]:
        context = next((s for s in sentences if phrase in s.lower()), "")
        flashcards.append(
            f"Q: Explain {phrase}?\nA: {context if context else 'Important concept about ' + phrase}"
        )
    
    # 2. Practice questions with answers
    quiz = []
    for i in range(min(5, len(key_phrases)-1)):
        quiz.append(
            f"Q: How does {key_phrases[i]} relate to {key_phrases[i+1]}?\n"
            f"A: They connect through {random.choice(['similar purpose', 'shared context', 'sequential relationship'])}"
        )
    
    # 3. Diagram suggestions with purposes
    diagrams = [
        f"Concept Map: Connect {', '.join(key_phrases[:3])} to show relationships",
        f"Timeline: Show development of {key_phrases[0]}",
        f"Venn Diagram: Compare {key_phrases[0]} and {key_phrases[1]}"
    ]
    
    # 4. Study plan
    plan = [
        "Day 1: Focus on understanding " + key_phrases[0],
        "Day 2: Study connections between " + " and ".join(key_phrases[:2]),
        "Day 3: Review all concepts and practice questions"
    ]
    
    # 5. Key examples
    examples = []
    for phrase in key_phrases[:3]:
        example_sentences = [s for s in sentences if phrase in s.lower()]
        if example_sentences:
            examples.append(f"Example of {phrase}: {example_sentences[0][:150]}...")
    
    return {
        "flashcards": flashcards,
        "quiz": quiz,
        "diagrams": diagrams,
        "plan": plan,
        "examples": examples
    }

# ======================
# MAIN APPLICATION
# ======================
def main():
    # st.set_page_config(
    #     layout="wide", 
    #     page_title="Ultimate Text Analyzer Pro", 
    #     page_icon="🧠",
    #     initial_sidebar_state="expanded"
    # )
    
    # Custom CSS for better visuals
    st.markdown("""
    <style>
    .summary-box {
        border-left: 5px solid #4CAF50;
        padding: 1rem;
        margin: 1rem 0;
        background-color: #f8f9fa;
    }
    .child-mode {
        background-color: #fff8e1;
        padding: 1rem;
        border-radius: 10px;
    }
    </style>
    """, unsafe_allow_html=True)
    
    st.title("🧠 Ultimate Text Analyzer Pro")
    st.markdown("""
    *Four-level analysis for all understanding needs:*
    - **👶 Child Mode** - Simple explanations with analogies
    - **🎓 Student Mode** - Clear summaries with examples
    - **🔬 Researcher Mode** - Technical analysis with findings
    - **🧠 Expert Mode** - Comprehensive scholarly breakdown
    """)
    
    uploaded_file = st.file_uploader(
        "Upload PDF, DOCX, PPTX, or TXT", 
        type=["pdf", "docx", "pptx", "txt"],
        help="Maximum file size: 50MB"
    )
    
    if uploaded_file:
        with st.spinner("Analyzing your document..."):
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                raw_text = extract_text_from_file(tmp_file.name, uploaded_file.type)
                # os.unlink(tmp_file.name)
            
            if raw_text:
                # Process text for different modes
                child_text = preprocess_text(raw_text, "child")
                standard_text = preprocess_text(raw_text)
                
                # Generate all summaries
                child_summary = generate_child_summary(child_text)
                student_summary = generate_student_summary(standard_text)
                researcher_summary = generate_researcher_summary(standard_text)
                expert_summary = generate_expert_summary(standard_text)
                
                # Create study materials
                tokenizer = RegexpTokenizer(r'\w+')  # Ensure tokenizer is initialized
                key_phrases = [w for w, _ in FreqDist(tokenizer.tokenize(standard_text)).most_common(15)]
                study_guide = create_study_guide(standard_text, key_phrases)
                
                # ======================
                # ENHANCED DISPLAY
                # ======================
                tab1, tab2, tab3, tab4 = st.tabs([
                    "👶 Child Mode", 
                    "🎓 Student Mode", 
                    "🔬 Researcher Mode", 
                    "📚 Study Kit"
                ])
                
                with tab1:
                    st.markdown('<div class="child-mode">', unsafe_allow_html=True)
                    st.header("Simple Explanation")
                    from transformers import T5ForConditionalGeneration, T5Tokenizer
                    from PIL import Image
                    import requests
                    from io import BytesIO
                    import textwrap
                    import matplotlib.pyplot as plt
                    import networkx as nx
                    from wordcloud import WordCloud
                    # Initialize T5 model
                    @st.cache_resource
                    def load_model():
                        model = T5ForConditionalGeneration.from_pretrained('t5-small')
                        tokenizer = T5Tokenizer.from_pretrained('t5-small')
                        return model, tokenizer

                    model, tokenizer = load_model()

                    # Function to generate child-friendly summary
                    def kid_friendly_summary(text):
                        # First get standard summary
                        inputs = tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=512, truncation=True)
                        outputs = model.generate(inputs, max_length=150, min_length=40, length_penalty=2.0, num_beams=4, early_stopping=True)
                        summary = tokenizer.decode(outputs[0])
                        
                        # Simplify further for kids
                        simple_summary = summary.replace("the", "").replace("and", "").replace("that", "")
                        simple_summary = " ".join(simple_summary.split()[:20])  # Keep it very short
                        
                        return summary, simple_summary

                    # Function to generate related image (using placeholder service)
                    def get_related_image(keyword):
                        try:
                            # In a real app, you'd use a proper image API
                            placeholder_url = f"https://source.unsplash.com/300x200/?{keyword},child"
                            response = requests.get(placeholder_url)
                            return Image.open(BytesIO(response.content))
                        except Exception:
                            return Image.new('RGB', (300, 200), color=(73, 109, 137))

                    # Function to create word cloud
                    def create_wordcloud(text):
                        wordcloud = WordCloud(width=600, height=400, background_color='white').generate(text)
                        plt.figure(figsize=(10, 5))
                        plt.imshow(wordcloud, interpolation='bilinear')
                        plt.axis("off")
                        return plt

                    # Function to create simple concept map
                    def create_concept_map(text):
                        words = [w for w in text.split() if len(w) > 3][:6]  # Get important words
                        G = nx.Graph()
                        
                        for i, word in enumerate(words):
                            G.add_node(word, image=get_related_image(word))
                            if i > 0:
                                G.add_edge(words[i-1], word)
                        
                        pos = nx.spring_layout(G)
                        plt.figure(figsize=(8, 6))
                        nx.draw(G, pos, with_labels=True, node_size=2000, node_color='skyblue', font_size=10, font_weight='bold')
                        return plt

                    # Streamlit app
                    def main():
                        st.title("📚 Story Simplifier for Kids")
                        st.subheader("Upload a story and I'll make it easy to understand!")
                        
                        uploaded_file = st.file_uploader("Choose a text file", type=["txt"])
                        
                        if uploaded_file is not None:
                            text = uploaded_file.read().decode("utf-8")
                            
                            st.success("Got your story! Let me read it...")
                            
                            with st.expander("Original Story"):
                                st.text(textwrap.fill(text, width=80))
                            
                            # Generate summaries
                            full_summary, kid_summary = kid_friendly_summary(text)
                            
                            st.header("✨ Super Simple Version")
                            st.write(kid_summary)
                            
                            # Visual elements
                            st.header("🎨 Let's See the Story!")
                            
                            col1, col2 = st.columns(2)
                            
                            with col1:
                                st.subheader("Word Cloud")
                                wordcloud = create_wordcloud(full_summary)
                                st.pyplot(wordcloud)
                                st.caption("Bigger words are more important in the story!")
                            
                            with col2:
                                st.subheader("Story Map")
                                concept_map = create_concept_map(full_summary)
                                st.pyplot(concept_map)
                                st.caption("See how the ideas connect!")
                            
                            # Interactive elements
                            st.header("🎭 Let's Play with the Story!")
                            
                            important_word = st.selectbox(
                                "Choose a word to see a picture:",
                                [w for w in full_summary.split() if len(w) > 4][:10]
                            )
                            
                            st.image(get_related_image(important_word), caption=f"This is about: {important_word}", width=300)
                            
                            # Simple quiz
                            st.header("🤔 Story Quiz")
                            answer = st.radio(
                                "What was the story mainly about?",
                                [kid_summary[:50] + "...", 
                                "Something completely different", 
                                "I'm not sure"]
                            )
                            
                            if answer == kid_summary[:50] + "...":
                                st.balloons()
                                st.success("Yes! You understood the story!")
                            else:
                                st.info("Let's read it again together!")

                            if __name__ == "__main__":
                                main()
                            st.markdown('</div>', unsafe_allow_html=True)
                
                with tab2:
                    st.header("Student Summary")
                    st.markdown('<div class="summary-box">', unsafe_allow_html=True)
                    st.write(student_summary)
                    # Load models
                    from transformers import pipeline
                    import pandas as pd
                    import matplotlib.pyplot as plt
                    import networkx as nx
                    from wordcloud import WordCloud
                    from textblob import TextBlob
                    import plotly.express as px
                    import spacy
                    from collections import Counter
                    from io import BytesIO
                    import concurrent.futures
                    @st.cache_resource
                    def load_models():
                        summarizer = pipeline("summarization", model="t5-small", tokenizer="t5-small")
                        nlp = spacy.load("en_core_web_sm")
                        return summarizer, nlp

                    summarizer, nlp = load_models()

                    # Process text
                    def analyze_text(text):
                        doc = nlp(text)
                        
                        # Extract entities
                        entities = [(ent.text, ent.label_) for ent in doc.ents]
                        entity_counts = Counter([label for _, label in entities])
                        
                        # Extract keywords (noun chunks)
                        keywords = [chunk.text for chunk in doc.noun_chunks if len(chunk.text.split()) <= 3]
                        keyword_counts = Counter(keywords).most_common(10)
                        
                        # Sentiment analysis
                        sentiment = TextBlob(text).sentiment
                        
                        return {
                            "entities": entities,
                            "entity_counts": entity_counts,
                            "keywords": keyword_counts,
                            "sentiment": sentiment,
                            "sentences": list(doc.sents)
                        }

                    # Visualization functions
                    def create_wordcloud(text):
                        wordcloud = WordCloud(width=800, height=400, background_color='white').generate(text)
                        plt.figure(figsize=(12, 6))
                        plt.imshow(wordcloud, interpolation='bilinear')
                        plt.axis("off")
                        return plt

                    def create_entity_graph(entities):
                        G = nx.Graph()
                        for entity, label in entities[:20]:  # Limit to top 20
                            G.add_node(entity, type=label)
                            G.add_edge(label, entity)  # Connect entity to its type
                        
                        pos = nx.spring_layout(G, k=0.5)
                        plt.figure(figsize=(12, 8))
                        nx.draw(G, pos, with_labels=True, node_size=2000, 
                            node_color='skyblue', font_size=10, font_weight='bold')
                        return plt

                    def create_timeline(sentences):
                        timeline_data = []
                        for i, sent in enumerate(sentences[:20]):  # Limit to first 20 sentences
                            timeline_data.append({
                                "Sentence": f"Sentence {i+1}",
                                "Text": sent.text[:50] + "..." if len(sent.text) > 50 else sent.text,
                                "Length": len(sent.text),
                                "Sentiment": TextBlob(sent.text).sentiment.polarity
                            })
                        df = pd.DataFrame(timeline_data)
                        fig = px.line(df, x="Sentence", y="Sentiment", 
                                    hover_data=["Text"], title="Sentiment Timeline")
                        return fig

                    # Streamlit app
                    def main():
                        # st.set_page_config(page_title="Academic Text Analyzer", layout="wide")
                        
                        st.title("📚 Academic Text Analyzer")
                        st.markdown("""
                        <style>
                        .big-font { font-size:18px !important; }
                        .summary-box { background-color: #f0f2f6; border-radius: 10px; padding: 20px; margin-bottom: 20px; }
                        .highlight { background-color: #fffacd; padding: 2px 5px; border-radius: 3px; }
                        </style>
                        """, unsafe_allow_html=True)
                        
                        uploaded_file = st.file_uploader("Upload your academic text (PDF, DOCX, or TXT)", 
                                                    type=["txt", "pdf", "docx"])
                        
                        if uploaded_file:
                            if uploaded_file.type == "application/pdf":
                                text = extract_text_from_pdf(uploaded_file)
                            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                                text = extract_text_from_docx(uploaded_file)
                            else:
                                text = uploaded_file.read().decode("utf-8")
                            
                            with st.expander("🔍 View Original Text"):
                                st.text(text[:2000] + "..." if len(text) > 2000 else text)
                            
                            with st.spinner("Analyzing text..."):
                                analysis = analyze_text(text)
                                
                                # Generate summary in chunks for better performance
                                chunk_size = 1000
                                chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
                                
                                with concurrent.futures.ThreadPoolExecutor() as executor:
                                    summaries = list(executor.map(
                                        lambda chunk: summarizer(chunk, max_length=150, min_length=50)[0]['summary_text'],
                                        chunks
                                    ))
                                
                                full_summary = " ".join(summaries)
                                key_points = "\n".join([f"- {sent.text}" for sent in nlp(full_summary).sents][:5])
                            
                            # Main content columns
                            col1, col2 = st.columns([1, 1])
                            
                            with col1:
                                st.markdown('<div class="summary-box"><h3>📝 Comprehensive Summary</h3><p class="big-font">' + 
                                        full_summary + '</p></div>', unsafe_allow_html=True)
                                
                                st.markdown('<div class="summary-box"><h3>🔑 Key Points</h3><p class="big-font">' + 
                                        key_points + '</p></div>', unsafe_allow_html=True)
                                
                                # Sentiment analysis
                                sentiment = analysis["sentiment"]
                                st.markdown(f"""
                                <div class="summary-box">
                                    <h3>🧠 Sentiment Analysis</h3>
                                    <p class="big-font">
                                        <span class="highlight">Polarity:</span> {sentiment.polarity:.2f} 
                                        ({'Positive' if sentiment.polarity > 0 else 'Negative' if sentiment.polarity < 0 else 'Neutral'})<br>
                                        <span class="highlight">Subjectivity:</span> {sentiment.subjectivity:.2f}
                                        ({'Objective' if sentiment.subjectivity < 0.5 else 'Subjective'})
                                    </p>
                                </div>
                                """, unsafe_allow_html=True)
                            
                            with col2:
                                # Visualizations
                                tab1, tab2, tab3 = st.tabs(["Word Cloud", "Entity Graph", "Sentiment Timeline"])
                                
                                with tab1:
                                    st.pyplot(create_wordcloud(full_summary))
                                    st.caption("Most frequent terms in the summarized content")
                                
                                with tab2:
                                    st.pyplot(create_entity_graph(analysis["entities"]))
                                    st.caption("Relationships between key entities in the text")
                                
                                with tab3:
                                    st.plotly_chart(create_timeline(analysis["sentences"]), use_container_width=True)
                                    st.caption("How sentiment evolves through the document")
                            
                            # Bottom section for interactive elements
                            st.header("📊 Deep Dive Analysis")
                            
                            col3, col4 = st.columns([1, 1])
                            
                            with col3:
                                st.subheader("Top Keywords")
                                keywords_df = pd.DataFrame(analysis["keywords"], columns=["Keyword", "Count"])
                                st.dataframe(keywords_df.style.background_gradient(cmap="Blues"), 
                                            use_container_width=True)
                                
                                selected_keyword = st.selectbox("Explore a keyword:", 
                                                            [kw[0] for kw in analysis["keywords"]])
                                if selected_keyword:
                                    examples = [sent.text for sent in analysis["sentences"] 
                                            if selected_keyword.lower() in sent.text.lower()][:3]
                                    st.write("**Example uses:**")
                                    for ex in examples:
                                        st.write(f"- {ex[:150]}{'...' if len(ex) > 150 else ''}")
                            
                            with col4:
                                st.subheader("Entity Types")
                                entities_df = pd.DataFrame.from_dict(analysis["entity_counts"], 
                                                                orient="index").reset_index()
                                entities_df.columns = ["Entity Type", "Count"]
                                fig = px.pie(entities_df, values="Count", names="Entity Type", 
                                            title="Distribution of Named Entities")
                                st.plotly_chart(fig, use_container_width=True)
                                
                                if st.checkbox("Show all named entities"):
                                    entities_df = pd.DataFrame(analysis["entities"], 
                                                            columns=["Entity", "Type"])
                                    st.dataframe(entities_df, use_container_width=True)

                    # Helper functions for file extraction (would need additional libraries)
                    def extract_text_from_pdf(file):
                        # In a real implementation, use PyPDF2 or pdfplumber
                        return "PDF text extraction would go here"

                    def extract_text_from_docx(file):
                        # In a real implementation, use python-docx
                        return "DOCX text extraction would go here"

                    if __name__ == "__main__":
                        main()
                
                with tab3:
                    st.header("Technical Analysis")
                    from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM
                    import pandas as pd
                    import numpy as np
                    import matplotlib.pyplot as plt
                    import plotly.express as px
                    import plotly.graph_objects as go
                    from wordcloud import WordCloud
                    import spacy
                    from spacy import displacy
                    from collections import Counter
                    import networkx as nx
                    import textwrap
                    from sklearn.feature_extraction.text import TfidfVectorizer
                    from concurrent.futures import ThreadPoolExecutor
                    import base64
                    import time

                                    # ------------------------------
                # 1. SETUP & MODEL LOADING
                # ------------------------------
                # st.set_page_config(
                #     page_title="Research Text Analyzer",
                #     layout="wide",
                #     initial_sidebar_state="expanded"
                # )

                @st.cache_resource
                def load_models():
                    """Load NLP models with progress indicators"""
                    with st.spinner("🔬 Loading research-grade models..."):
                        # T5 Summarization
                        t5_tokenizer = AutoTokenizer.from_pretrained("t5-small")
                        t5_model = AutoModelForSeq2SeqLM.from_pretrained("t5-small")
                        summarizer = pipeline("summarization", model=t5_model, tokenizer=t5_tokenizer)
                        
                        # SpaCy for technical analysis
                        nlp = spacy.load("en_core_web_sm")
                        
                        # TF-IDF Vectorizer
                        tfidf = TfidfVectorizer(max_features=100, stop_words='english')
                        
                    return summarizer, nlp, tfidf

                summarizer, nlp, tfidf = load_models()

                # ------------------------------
                # 2. TEXT PROCESSING FUNCTIONS
                # ------------------------------
                def technical_analysis(text):
                    """Perform deep technical analysis on research text"""
                    doc = nlp(text)
                    
                    # Entity Extraction
                    entities = [(ent.text, ent.label_) for ent in doc.ents]
                    entity_freq = Counter([label for _, label in entities])
                    
                    # Technical Term Extraction (noun phrases with adjectives)
                    tech_terms = [chunk.text for chunk in doc.noun_chunks 
                                if any(t.pos_ == 'ADJ' for t in chunk)]
                    
                    # Claim/Assertion Detection (subjective sentences)
                    claims = [sent.text for sent in doc.sents 
                            if any(t.dep_ == 'nsubj' and t.pos_ == 'VERB' for t in sent)]
                    
                    # TF-IDF Keywords
                    tfidf_matrix = tfidf.fit_transform([text])
                    feature_names = tfidf.get_feature_names_out()
                    tfidf_scores = zip(feature_names, tfidf_matrix.toarray()[0])
                    top_keywords = sorted(tfidf_scores, key=lambda x: x[1], reverse=True)[:15]
                    
                    return {
                        "entities": entities,
                        "entity_freq": entity_freq,
                        "tech_terms": tech_terms,
                        "claims": claims,
                        "keywords": top_keywords,
                        "doc": doc
                    }

                def generate_summary(text, technical_focus=True):
                    """Generate summary with technical focus"""
                    if technical_focus:
                        prefix = "summarize technical research paper: "
                    else:
                        prefix = "summarize: "
                    
                    chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
                    
                    with ThreadPoolExecutor() as executor:
                        summaries = list(executor.map(
                            lambda chunk: summarizer(prefix + chunk, 
                                                max_length=150, 
                                                min_length=50,
                                                do_sample=False)[0]['summary_text'],
                            chunks
                        ))
                    
                    return " ".join(summaries)

                # ------------------------------
                # 3. VISUALIZATION FUNCTIONS
                # ------------------------------
                def create_tech_wordcloud(text):
                    """Technical word cloud with coloring"""
                    wordcloud = WordCloud(
                        width=1000, 
                        height=600,
                        background_color='white',
                        colormap='viridis',
                        collocations=False
                    ).generate(text)
                    
                    fig, ax = plt.subplots(figsize=(16, 9))
                    ax.imshow(wordcloud, interpolation='bilinear')
                    ax.axis("off")
                    ax.set_title("Technical Term Cloud", pad=20, fontsize=18)
                    return fig

                def create_entity_network(entities):
                    """Interactive entity relationship network"""
                    G = nx.Graph()
                    
                    # Add nodes with type grouping
                    for entity, label in entities[:50]:  # Limit for performance
                        G.add_node(entity, type=label, size=5 + len(entity)/2)
                    
                    # Add edges between entities of same type
                    for i, (ent1, type1) in enumerate(entities):
                        for j, (ent2, type2) in enumerate(entities[i+1:i+5]):
                            if type1 == type2:
                                G.add_edge(ent1, ent2, weight=0.3)
                    
                    pos = nx.spring_layout(G, k=0.5)
                    
                    # Create Plotly figure
                    edge_x = []
                    edge_y = []
                    for edge in G.edges():
                        x0, y0 = pos[edge[0]]
                        x1, y1 = pos[edge[1]]
                        edge_x.extend([x0, x1, None])
                        edge_y.extend([y0, y1, None])
                    
                    edge_trace = go.Scatter(
                        x=edge_x, y=edge_y,
                        line=dict(width=0.5, color='#888'),
                        hoverinfo='none',
                        mode='lines')
                    
                    node_x = []
                    node_y = []
                    node_text = []
                    node_size = []
                    for node in G.nodes():
                        x, y = pos[node]
                        node_x.append(x)
                        node_y.append(y)
                        node_text.append(f"{node}<br>Type: {G.nodes[node]['type']}")
                        node_size.append(G.nodes[node]['size'])
                    
                    node_trace = go.Scatter(
                        x=node_x, y=node_y,
                        mode='markers+text',
                        textposition='top center',
                        hoverinfo='text',
                        text=node_text,
                        marker=dict(
                            showscale=True,
                            colorscale='Rainbow',
                            size=node_size,
                            color=[hash(node) % 100 for node in G.nodes()],
                            line_width=2))
                    
                    fig = go.Figure(data=[edge_trace, node_trace],
                                layout=go.Layout(
                                    showlegend=False,
                                    hovermode='closest',
                                    margin=dict(b=20,l=5,r=5,t=40),
                                    height=700,
                                    title="Entity Relationship Network"))
                    
                    return fig

                def create_technical_timeline(doc):
                    """Timeline of technical concepts"""
                    timeline_data = []
                    for i, sent in enumerate(doc.sents):
                        tech_terms = [chunk.text for chunk in sent.noun_chunks 
                                    if any(t.pos_ == 'ADJ' for t in chunk)]
                        if tech_terms:
                            timeline_data.append({
                                "Position": i,
                                "Sentence": f"Sentence {i+1}",
                                "Technical Terms": ", ".join(tech_terms[:3]),
                                "Term Count": len(tech_terms)
                            })
                    
                    df = pd.DataFrame(timeline_data[:100])  # Limit for performance
                    
                    fig = px.scatter(df, x="Position", y="Term Count",
                                    size="Term Count", color="Term Count",
                                    hover_data=["Sentence", "Technical Terms"],
                                    title="Technical Concept Density Timeline",
                                    color_continuous_scale='thermal')
                    
                    fig.update_layout(height=500)
                    return fig

                # ------------------------------
                # 4. STREAMLIT APP LAYOUT
                # ------------------------------
                def main():
                    st.title("🔬 Research Text Analyzer")
                    st.markdown("""
                    <style>
                    .main-title { font-size: 2.5rem !important; }
                    .section-title { border-bottom: 2px solid #eee; padding-bottom: 0.3rem; }
                    .technical-term { background-color: #e6f7ff; padding: 2px 5px; border-radius: 3px; }
                    </style>
                    """, unsafe_allow_html=True)
                    
                    # File Upload
                    uploaded_file = st.file_uploader(
                        "Upload research document (PDF/TXT/DOCX)", 
                        type=["pdf", "txt", "docx"],
                        accept_multiple_files=False
                    )
                    
                    if uploaded_file:
                        # Extract text based on file type
                        # Helper functions for file extraction
                        def extract_text_from_pdf(file):
                            # In a real implementation, use PyPDF2 or pdfplumber
                            return "PDF text extraction would go here"

                        def extract_text_from_docx(file):
                            # In a real implementation, use python-docx
                            return "DOCX text extraction would go here"

                        if uploaded_file.type == "application/pdf":
                            text = extract_text_from_pdf(uploaded_file)
                        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                            text = extract_text_from_docx(uploaded_file)
                        else:
                            text = uploaded_file.read().decode("utf-8")
                        
                        # Analysis Tabs
                        tab1, tab2, tab3, tab4 = st.tabs([
                            "📊 Overview", 
                            "🔍 Deep Analysis", 
                            "📈 Technical Patterns", 
                            "💡 Key Insights"
                        ])
                        
                        with st.spinner("Performing advanced text analysis..."):
                            analysis = technical_analysis(text)
                            summary = generate_summary(text)
                        
                        # ------------------------------
                        # TAB 1: OVERVIEW
                        # ------------------------------
                        with tab1:
                            col1, col2 = st.columns([1, 1])
                            
                            with col1:
                                st.subheader("Technical Summary", divider='blue')
                                st.markdown(f'<div style="background-color:#f8f9fa; padding:15px; border-radius:10px;">{summary}</div>', 
                                        unsafe_allow_html=True)
                                
                                st.subheader("Top Technical Terms", divider='blue')
                                terms_df = pd.DataFrame(analysis["keywords"], columns=["Term", "TF-IDF Score"])
                                st.dataframe(
                                    terms_df.style.background_gradient(cmap="Blues"),
                                    use_container_width=True,
                                    height=300
                                )
                            
                            with col2:
                                st.subheader("Entity Frequency", divider='blue')
                                entity_df = pd.DataFrame.from_dict(analysis["entity_freq"], 
                                                                orient='index').reset_index()
                                entity_df.columns = ["Entity Type", "Count"]
                                
                                fig = px.bar(entity_df, 
                                            x="Count", 
                                            y="Entity Type", 
                                            orientation='h',
                                            color="Count",
                                            color_continuous_scale='deep')
                                st.plotly_chart(fig, use_container_width=True)
                        
                        # ------------------------------
                        # TAB 2: DEEP ANALYSIS
                        # ------------------------------
                        with tab2:
                            st.subheader("Entity Relationship Network", divider='orange')
                            st.plotly_chart(create_entity_network(analysis["entities"]), 
                                        use_container_width=True)
                            
                            st.subheader("Technical Term Cloud", divider='orange')
                            st.pyplot(create_tech_wordcloud(" ".join(analysis["tech_terms"])))
                        
                        # ------------------------------
                        # TAB 3: TECHNICAL PATTERNS
                        # ------------------------------
                        with tab3:
                            st.subheader("Technical Concept Timeline", divider='green')
                            st.plotly_chart(create_technical_timeline(analysis["doc"]), 
                                        use_container_width=True)
                            
                            st.subheader("Research Claims/Assertions", divider='green')
                            for i, claim in enumerate(analysis["claims"][:10]):
                                st.markdown(f"{i+1}. {claim}")
                                st.write("---")
                        
                        # ------------------------------
                        # TAB 4: KEY INSIGHTS
                        # ------------------------------
                        with tab4:
                            st.subheader("Research Highlights", divider='purple')
                            
                            insight_col1, insight_col2 = st.columns(2)
                            
                            with insight_col1:
                                st.markdown("#### 🔥 Most Significant Terms")
                                for term, score in analysis["keywords"][:5]:
                                    st.markdown(f'- <span class="technical-term">{term}</span> (Score: {score:.2f})', 
                                            unsafe_allow_html=True)
                                
                                st.markdown("#### 🧠 Key Entities")
                                entity_types = {
                                    "PERSON": "People/Researchers",
                                    "ORG": "Organizations",
                                    "DATE": "Dates/Timeframes",
                                    "GPE": "Geopolitical Entities"
                                }
                                
                                for ent_type, count in analysis["entity_freq"].most_common(4):
                                    st.markdown(f"- **{entity_types.get(ent_type, ent_type)}**: {count} mentions")
                            
                            with insight_col2:
                                st.markdown("#### 📌 Technical Concept Density")
                                st.metric("Technical Terms per Sentence", 
                                        f"{len(analysis['tech_terms'])/len(list(analysis['doc'].sents)):.2f}")
                                
                                st.markdown("#### 🗂️ Document Structure")
                                sent_lengths = [len(sent.text) for sent in analysis["doc"].sents]
                                st.metric("Average Sentence Length", f"{np.mean(sent_lengths):.1f} characters")
                                st.metric("Total Technical Terms", len(analysis["tech_terms"]))
                                
                                if st.button("📥 Export Analysis Report"):
                                                                        # Define generate_report function
                                                                        def generate_report(analysis, summary):
                                                                            # Placeholder implementation for PDF report generation
                                                                            # In a real app, generate a PDF and return its bytes
                                                                            return b"PDF report bytes would go here"

                                                                        report = generate_report(analysis, summary)
                                                                        st.download_button(
                                                                            label="Download Report (PDF)",
                                                                            data=report,
                                                                            file_name="research_analysis_report.pdf",
                                                                            mime="application/pdf"
                                                                        )
                                    
                            # ------------------------------
                            # HELPER FUNCTIONS
                                    # ------------------------------
                            def extract_text_from_pdf(file):
                                        # Implement with PyPDF2 or pdfplumber
                                        return "PDF text extraction would go here"
                                    
                            def extract_text_from_docx(file):
                                        # Implement with python-docx
                                        return "DOCX text extraction would go here"
                                    
                            def generate_report(analysis, summary):
                                        # Placeholder implementation for PDF report generation
                                        # In a real app, generate a PDF and return its bytes
                                        return b"PDF report bytes would go here"
                                    
                            if __name__ == "__main__":
                                        main()
                                
                with tab4:
                    st.header("Complete Study Kit")
                    import spacy
                    from collections import Counter
                    import networkx as nx
                    from wordcloud import WordCloud
                    import textwrap
                    import random
                    from datetime import datetime, timedelta
                    from io import BytesIO
                    from PIL import Image
                    import base64
                    import concurrent.futures
                    
                                    # ------------------------------
                # 1. SETUP & MODEL LOADING
                # ------------------------------
                # st.set_page_config(
                #     page_title="AI Study Guide Generator",
                #     layout="wide",
                #     initial_sidebar_state="expanded"
                # )

                @st.cache_resource
                def load_models():
                    """Load all required NLP models"""
                    with st.spinner("📚 Loading educational models..."):
                        # T5 for summarization
                        t5_model = T5ForConditionalGeneration.from_pretrained('t5-small')
                        t5_tokenizer = T5Tokenizer.from_pretrained('t5-small')
                        summarizer = pipeline("summarization", model=t5_model, tokenizer=t5_tokenizer)
                        
                        # SpaCy for analysis
                        nlp = spacy.load("en_core_web_sm")
                        
                    return summarizer, nlp

                summarizer, nlp = load_models()

                # ------------------------------
                # 2. STUDY MATERIAL GENERATION
                # ------------------------------
                def generate_study_materials(text):
                    """Generate all study materials from text"""
                    doc = nlp(text)
                    
                    # Key Concepts (noun phrases with adjectives/numbers)
                    key_concepts = [chunk.text for chunk in doc.noun_chunks 
                                if any(t.pos_ in ('ADJ', 'NUM') for t in chunk)]
                    concept_counts = Counter(key_concepts).most_common(20)
                    
                    # Flashcards (term-definition pairs)
                    flashcards = []
                    for sent in doc.sents:
                        if len(sent.ents) > 0:  # Sentences with entities
                            term = random.choice([ent.text for ent in sent.ents])
                            definition = sent.text
                            flashcards.append((term, definition))
                    
                    # Practice Questions (transform statements)
                    questions = []
                    question_types = [
                        ("Explain the concept of {} in your own words", "NOUN"),
                        ("What is the significance of {}?", "ENTITY"),
                        ("Compare and contrast {} and {}", "NOUN_NOUN"),
                        ("How does {} affect {}?", "NOUN_NOUN")
                    ]
                    
                    for template, pattern in question_types:
                        if pattern == "NOUN":
                            noun = random.choice([chunk.text for chunk in doc.noun_chunks])
                            questions.append(template.format(noun))
                        elif pattern == "ENTITY":
                            ent = random.choice([ent.text for ent in doc.ents if ent.label_ in ('ORG', 'PERSON', 'GPE')])
                            questions.append(template.format(ent))
                        elif pattern == "NOUN_NOUN":
                            nouns = random.sample([chunk.text for chunk in doc.noun_chunks], 2)
                            questions.append(template.format(nouns[0], nouns[1]))
                    
                    # Study Plan (schedule based on concepts)
                    study_days = min(7, len(key_concepts))
                    study_plan = []
                    start_date = datetime.now()
                    for i, concept in enumerate(key_concepts[:study_days]):
                        study_plan.append({
                            "Day": (start_date + timedelta(days=i)).strftime("%A, %b %d"),
                            "Focus Area": concept,
                            "Tasks": [
                                f"Review {concept} definition",
                                f"Find 2 examples of {concept}",
                                f"Create diagram for {concept}"
                            ]
                        })
                    
                    # Diagrams (extract relationships)
                    relationships = []
                    for sent in doc.sents:
                        if any(t.dep_ in ('nsubj', 'dobj') for t in sent):
                            subj = next((t.text for t in sent if t.dep_ == 'nsubj'), None)
                            obj = next((t.text for t in sent if t.dep_ == 'dobj'), None)
                            if subj and obj:
                                relationships.append((subj, sent.root.text, obj))
                    
                    return {
                        "key_concepts": concept_counts,
                        "flashcards": flashcards[:20],  # Limit to 20
                        "practice_questions": questions[:15],  # Limit to 15
                        "study_plan": study_plan,
                        "relationships": relationships[:15],  # Limit to 15
                        "doc": doc
                    }

                # ------------------------------
                # 3. VISUALIZATION FUNCTIONS
                # ------------------------------
                def create_concept_map(relationships):
                    """Create interactive concept relationship diagram"""
                    G = nx.DiGraph()
                    
                    for src, rel, target in relationships:
                        G.add_edge(src, target, label=rel)
                    
                    pos = nx.spring_layout(G, k=0.8)
                    
                    edge_x = []
                    edge_y = []
                    edge_text = []
                    for edge in G.edges(data=True):
                        x0, y0 = pos[edge[0]]
                        x1, y1 = pos[edge[1]]
                        edge_x.extend([x0, x1, None])
                        edge_y.extend([y0, y1, None])
                        edge_text.append(edge[2]['label'])
                    
                    edge_trace = go.Scatter(
                        x=edge_x, y=edge_y,
                        line=dict(width=1, color='#888'),
                        hoverinfo='none',
                        mode='lines')
                    
                    node_x = []
                    node_y = []
                    node_text = []
                    for node in G.nodes():
                        x, y = pos[node]
                        node_x.append(x)
                        node_y.append(y)
                        node_text.append(node)
                    
                    node_trace = go.Scatter(
                        x=node_x, y=node_y,
                        mode='markers+text',
                        text=node_text,
                        textposition='top center',
                        hoverinfo='text',
                        marker=dict(
                            showscale=True,
                            colorscale='Rainbow',
                            size=20,
                            color=[hash(node) % 100 for node in G.nodes()],
                            line_width=2))
                    
                    # Create edge labels
                    edge_label_x = []
                    edge_label_y = []
                    edge_label_text = []
                    for i, (src, target) in enumerate(G.edges()):
                        x0, y0 = pos[src]
                        x1, y1 = pos[target]
                        edge_label_x.append((x0 + x1) / 2)
                        edge_label_y.append((y0 + y1) / 2)
                        edge_label_text.append(G.edges[src, target]['label'])
                    
                    edge_label_trace = go.Scatter(
                        x=edge_label_x,
                        y=edge_label_y,
                        mode='text',
                        text=edge_label_text,
                        textfont=dict(size=10, color='red'),
                        hoverinfo='none'
                    )
                    
                    fig = go.Figure(data=[edge_trace, node_trace, edge_label_trace],
                                layout=go.Layout(
                                    showlegend=False,
                                    hovermode='closest',
                                    margin=dict(b=20,l=5,r=5,t=40),
                                    height=600,
                                    title="Concept Relationship Map"))
                    
                    return fig

                def create_study_calendar(study_plan):
                    """Visualize study plan as calendar"""
                    df = pd.DataFrame(study_plan)
                    df['Date'] = pd.to_datetime(df['Day'], format='%A, %b %d')
                    df['Tasks'] = df['Tasks'].apply(lambda x: "\n".join(x))
                    
                    fig = px.timeline(df, 
                                    x_start="Date", 
                                    x_end="Date",
                                    y="Focus Area",
                                    color="Focus Area",
                                    hover_data=["Tasks"],
                                    title="Weekly Study Plan")
                    
                    fig.update_yaxes(autorange="reversed")
                    fig.update_layout(height=400)
                    return fig

                # ------------------------------
                # 4. STREAMLIT APP INTERFACE
                # ------------------------------
                def main():
                    st.title("📖 AI Study Guide Generator")
                    st.markdown("""
                    <style>
                    .flashcard {
                        background-color: #f8f9fa;
                        border-radius: 10px;
                        padding: 15px;
                        margin: 10px 0;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                    }
                    .question {
                        background-color: #e6f7ff;
                        border-left: 4px solid #1890ff;
                        padding: 12px;
                        margin: 8px 0;
                    }
                    </style>
                    """, unsafe_allow_html=True)
                    
                    # File Upload
                    uploaded_file = st.file_uploader(
                        "Upload your textbook, research paper, or notes (PDF/TXT/DOCX)", 
                        type=["pdf", "txt", "docx"]
                    )
                    
                    if uploaded_file:
                        # Extract text
                        if uploaded_file.type == "application/pdf":
                            text = extract_text_from_pdf(uploaded_file)
                        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                            text = extract_text_from_docx(uploaded_file)
                        else:
                            text = uploaded_file.read().decode("utf-8")
                        
                        # Generate materials
                        with st.spinner("🧠 Creating your personalized study materials..."):
                            materials = generate_study_materials(text)
                            summary = summarizer(text, max_length=150, min_length=50)[0]['summary_text']
                        
                        # Main Tabs
                        tab1, tab2, tab3, tab4 = st.tabs([
                            "📚 Study Guide", 
                            "🃏 Flashcards", 
                            "❓ Practice Questions", 
                            "📅 Study Plan"
                        ])
                        
                        # ------------------------------
                        # TAB 1: STUDY GUIDE
                        # ------------------------------
                        with tab1:
                            col1, col2 = st.columns([1, 1])
                            
                            with col1:
                                st.subheader("📝 Chapter Summary")
                                st.markdown(f'<div style="background-color:#f8f9fa; padding:15px; border-radius:10px;">{summary}</div>', 
                                        unsafe_allow_html=True)
                                
                                st.subheader("🔑 Key Concepts")
                                concepts_df = pd.DataFrame(materials["key_concepts"], columns=["Concept", "Frequency"])
                                st.dataframe(
                                    concepts_df.style.background_gradient(cmap="Blues"),
                                    use_container_width=True,
                                    height=400
                                )
                            
                            with col2:
                                st.subheader("🧩 Concept Relationships")
                                st.plotly_chart(create_concept_map(materials["relationships"]), 
                                            use_container_width=True)
                        
                        # ------------------------------
                        # TAB 2: FLASHCARDS
                        # ------------------------------
                        with tab2:
                            st.subheader("📇 Interactive Flashcards")
                            st.caption("Click the flip button to see definitions")
                            
                            cols = st.columns(3)
                            for i, (term, definition) in enumerate(materials["flashcards"]):
                                with cols[i % 3]:
                                    with st.expander(f"**{term}**", expanded=False):
                                        st.markdown(f'<div style="font-size:14px;">{definition}</div>', 
                                                unsafe_allow_html=True)
                            
                            # Flashcard download option
                            flashcard_csv = "\n".join([f"{t}\t{d}" for t, d in materials["flashcards"]])
                            st.download_button(
                                label="📥 Download Flashcards (CSV)",
                                data=flashcard_csv,
                                file_name="flashcards.csv",
                                mime="text/csv"
                            )
                        
                        # ------------------------------
                        # TAB 3: PRACTICE QUESTIONS
                        # ------------------------------
                        with tab3:
                            st.subheader("📝 Practice Questions")
                            
                            for i, question in enumerate(materials["practice_questions"]):
                                st.markdown(
                                    f'<div class="question"><b>Q{i+1}.</b> {question}</div>',
                                    unsafe_allow_html=True
                                )
                            
                            # Answer field for active recall
                            st.subheader("✍️ Your Notes")
                            user_notes = st.text_area("Write your answers here (this saves to your browser)", 
                                                    height=200)
                            
                            if st.button("Generate More Questions"):
                                st.experimental_rerun()
                        
                        # ------------------------------
                        # TAB 4: STUDY PLAN
                        # ------------------------------
                        with tab4:
                            st.subheader("🗓️ Personalized Study Plan")
                            st.plotly_chart(create_study_calendar(materials["study_plan"]), 
                                        use_container_width=True)
                            
                            st.subheader("📋 Daily Breakdown")
                            for day in materials["study_plan"]:
                                with st.expander(f"📌 {day['Day']}: {day['Focus Area']}", expanded=False):
                                    st.markdown("**Tasks:**")
                                    for task in day["Tasks"]:
                                        st.markdown(f"- {task}")
                                    
                                    # Pomodoro timer
                                    if st.button(f"⏱️ Start Timer for {day['Focus Area']}"):
                                        st.session_state.timer_active = True
                                        st.session_state.timer_end = datetime.now() + timedelta(minutes=25)
                                    
                                    if 'timer_active' in st.session_state and st.session_state.timer_active:
                                        time_left = st.session_state.timer_end - datetime.now()
                                        if time_left.total_seconds() > 0:
                                            mins, secs = divmod(int(time_left.total_seconds()), 60)
                                            st.write(f"⏳ Time remaining: {mins}m {secs}s")
                                        else:
                                            st.success("🍅 Pomodoro complete! Take a 5 minute break.")
                                            st.session_state.timer_active = False
                            
                            # Export plan
                            plan_csv = pd.DataFrame(materials["study_plan"]).to_csv(index=False)
                            st.download_button(
                                label="📥 Download Study Plan (CSV)",
                                data=plan_csv,
                                file_name="study_plan.csv",
                                mime="text/csv"
                            )

                # Helper functions for file processing
                def extract_text_from_pdf(file):
                    # Implement with PyPDF2 or pdfplumber
                    return "PDF text would be extracted here"

                def extract_text_from_docx(file):
                    # Implement with python-docx
                    return "DOCX text would be extracted here"

                if __name__ == "__main__":
                    main()

if __name__ == "__main__":
    main()

elif choice == "URL":
      import requests
      from bs4 import BeautifulSoup
      from transformers import pipeline

      st.title("URL Summarizer")

      url = st.text_input("Enter URL to summarize:")

      if st.button("Summarize URL"):
          if url:
              try:
                  response = requests.get(url)
                  response.raise_for_status()  # Raise an exception for bad status codes
                  soup = BeautifulSoup(response.content, 'html.parser')
                  text = soup.get_text()

                  # Initialize the summarizer outside the button's conditional block
                  summarizer = pipeline("summarization", model="t5-base", tokenizer="t5-base", framework="pt")

                  summary = summarizer(text, min_length=100, do_sample=False)
                  st.subheader("Summary:")
                  st.write(summary[0]['summary_text'])

              except requests.exceptions.RequestException as e:
                  st.error(f"Error fetching URL: {e}")
              except Exception as e:
                  st.error(f"An error occurred during summarization: {e}")
          else:
              st.warning("Please enter a URL.")
        
          # st.write("The summarization process uses the T5 model, which is a transformer-based model pre-trained on a large corpus of text.")
          # st.write("The model is capable of generating high-quality summaries for a wide range of text inputs.")
          # st.write("The summarization process may take some time depending on the length of the text and the complexity of the content.")
          # st.write("Please ensure that the URL you enter is accessible and contains text content that can be summarized.")
          # st.write("If the URL is not accessible or does not contain text content, the summarization process may fail.")
          # st.write("You can adjust the maximum and minimum length of the summary by modifying the parameters in the code.")


elif choice == "About":
    st.title("About this App")
    st.write("This app uses the T5 model from Hugging Face's Transformers library to summarize text.")
    st.write("The T5 model is a transformer-based model that has been pre-trained on a large corpus of text.")
    st.write("You can enter any text in the text area and click the 'Summarize' button to get a summary.")
    st.write("This app is built using Streamlit, a Python library for creating web apps.")

    st.write("You can find the source code for this app on GitHub.")
    st.write("If you have any questions or feedback, please feel free to reach out.")
    st.write("Thank you for using this app!")


