### CODE STARTS HERE ###
# Import all original libraries plus additional ones needed for new features
import streamlit as st
import os
import tempfile
import random
from typing import Optional, List, Dict, Tuple, Any
from datetime import datetime
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER
import textwrap
import pandas as pd
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import plotly.graph_objects as go
import networkx as nx
import numpy as np
import torch
from collections import defaultdict
import re
import concurrent.futures
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from transformers import T5Tokenizer, T5ForConditionalGeneration, pipeline, BertTokenizer, BertModel
import spacy
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from youtube_transcript_api import NoTranscriptFound, TranscriptsDisabled
# Add the import where it's actually used for audio extraction from video
from moviepy.editor import VideoFileClip # [cite: 2]
# Remove the redundant import if VideoFileClip is used directly: import moviepy.editor as mp
import whisper
from PIL import Image # [cite: 2]
from spacy import displacy

# New import for Text-to-Speech
from gtts import gTTS
import pygame # For playing audio

# Initialize NLTK with enhanced error handling
try:
    nltk.download('punkt', quiet=True) # [cite: 2]
    nltk.download('stopwords', quiet=True) # [cite: 2]
    nltk.download('wordnet', quiet=True) # [cite: 2]
except Exception as e:
    st.warning(f"Could not download NLTK data automatically: {e}. Some features might be limited.") # [cite: 2]

# Initialize Pygame Mixer
try:
    pygame.mixer.init()
except Exception as e:
    st.warning(f"Could not initialize Pygame mixer for audio playback: {e}. Audio features might be limited.")

# Enhanced Streamlit page configuration
st.set_page_config(
    layout="wide",
    page_title="Advanced AI Text & Media Analyzer Suite",
    page_icon="🧠", # [cite: 3]
    menu_items={
        'Get Help': 'https://example.com/help', # [cite: 3]
        'Report a bug': "https://example.com/bug", # [cite: 3]
        'About': "# Advanced AI Content Analysis Suite" # [cite: 3]
    }
)

# Enhanced CSS styling with additional UI improvements
st.markdown("""
<style>
/* General Streamlit App Styling */
.stApp {
    background-color: #004d40; /* Dark teal background for main page */ /* [cite: 3, 4] */
    color: white; /* Default text color to white for overall readability on dark background */ /* [cite: 4] */
    line-height: 1.6; /* Increased line spacing for better readability from first block */ /* [cite: 4] */
    font-family: 'Segoe UI', Roboto, Arial, sans-serif; /* Modern, readable font stack from first block */ /* [cite: 4] */
}
/* Sidebar Styling (White) */
.css-1d391kg { /* Target the main sidebar container - Class name may vary with Streamlit versions */
    background-color: #FFFFFF !important; /* White background for sidebar from second block */ /* [cite: 5] */
    padding: 20px; /* Increased padding from first block */ /* [cite: 5, 6] */
    box-shadow: 2px 2px 10px rgba(0,0,0,0.05); /* Softer shadow from first block */ /* [cite: 6, 7] */
    border-radius: 10px; /* Rounded corners from first block */ /* [cite: 6, 8] */
}
/* Ensure text elements within the sidebar are black */
/* This might require targeting specific Streamlit component classes within the sidebar if default color inheritance is an issue */
.css-1d391kg label, /* Labels in sidebar (e.g., selectbox, radio, checkbox) */ /* [cite: 9] */
.css-1d391kg .st-bv, /* Markdown text in sidebar */ /* [cite: 9] */
.css-1d391kg .st-cc, /* Text elements potentially */ /* [cite: 9] */
.css-1d391kg .st-cq, /* More text elements potentially */ /* [cite: 9] */
.css-1d391kg div,
.css-1d391kg span,
.css-1d391kg p,
.css-1d391kg li,
.css-1d391kg .stRadio > label > div > p /* Specifically target radio button labels */
 {
    color: black !important; /* Set text color to black for these elements */ /* [cite: 9, 10] */
}

/* Main Content Area Headings */
h1, h2, h3, h4, h5, h6 {
    color: white; /* Header text remains white for visibility on dark teal background */ /* [cite: 12] */
}
/* Buttons Styling */
.stButton>button { /* Target Streamlit buttons */
    background-color: #00796b; /* Teal background from second block */ /* [cite: 13] */
    color: white; /* White text for contrast */ /* [cite: 14] */
    border-radius: 8px; /* Rounded corners from first block */ /* [cite: 15] */
    padding: 12px 24px; /* Increased padding from first block */ /* [cite: 16] */
    font-weight: 500; /* Medium font weight from first block */ /* [cite: 17] */
    transition: all 0.3s ease; /* Smooth transition for hover effects from first block */ /* [cite: 18] */
    cursor: pointer; /* Indicate interactivity from first block */ /* [cite: 19] */
    border: none; /* Remove default border from first block */ /* [cite: 20] */
}
.stButton>button:hover { /* Style for button on hover */
    background-color: #004d40; /* Darker shade on hover from second block */ /* [cite: 21] */
    transform: translateY(-2px); /* Slight lift effect from first block */ /* [cite: 22] */
    box-shadow: 0 4px 12px rgba(0,0,0,0.15); /* More pronounced shadow on hover from first block */ /* [cite: 23] */
}
/* File Uploader Styling */
.stFileUploader label { /* Target file uploader label */
    color: white !important; /* Set file uploader label to white for readability on dark background */ /* [cite: 24] */
}
.stFileUploader>div>section { /* Dropzone area */
    border: 2px dashed #00796b; /* Dotted border with teal accent color (from second block) */ /* [cite: 25] */
    border-radius: 8px; /* Rounded corners from first block */ /* [cite: 26] */
    background-color: rgba(0, 121, 107, 0.1); /* Very light tint of teal */ /* [cite: 27] */
    padding: 20px; /* Padding inside the uploader area from first block */ /* [cite: 28] */
    text-align: center; /* Center align text from first block */ /* [cite: 29] */
}
.stFileUploader section p { /* Text within the dropzone */
    color: white !important; /* Make dropzone text white */ /* [cite: 30] */
}
.stFileUploader section button[data-testid="baseButton-secondary"] { /* Browse files button */
    color: black !important; /* Text color for "Browse files" button */ /* [cite: 31] */
    background-color: #e0f2f1; /* Light background for button from second block */ /* [cite: 32] */
    border: 1px solid #00796b; /* Border for button from second block */ /* [cite: 33] */
    border-radius: 8px; /* Match button border-radius */ /* [cite: 34] */
    padding: 8px 15px; /* Adjust padding for this button */ /* [cite: 35] */
}
.stFileUploader section button[data-testid="baseButton-secondary"]:hover {
    background-color: #c0e0df; /* Slightly darker on hover */ /* [cite: 36] */
}
/* Text Input and Text Area Styling */
.stTextInput>div>div>input, .stTextArea>div>textarea {
    border: 1px solid #00796b; /* Teal border */ /* [cite: 37] */
    border-radius: 8px; /* Rounded corners from first block */
    padding: 10px; /* Padding inside input fields from first block */ /* [cite: 38] */
    width: 100%; /* Full width from first block */ /* [cite: 39] */
    box-sizing: border-box; /* Include padding and border in size from first block */ /* [cite: 40] */
    color: white !important; /* Set input text color to white */ /* [cite: 41] */
    background-color: #005548; /* Slightly lighter background than main page */ /* [cite: 42] */
}
.stTextInput label, .stTextArea label { /* Target labels for text input and text area */
    color: white !important; /* Set label text color to white */ /* [cite: 43] */
}
.stTextInput>div>div>input:focus, .stTextArea>div>textarea:focus {
    border-color: #00796b; /* Highlight color on focus (teal) */ /* [cite: 44] */
    outline: none; /* Remove default outline */ /* [cite: 45] */
    box-shadow: 0 0 0 0.2rem rgba(0, 121, 107, 0.25); /* Glow effect on focus (teal) */ /* [cite: 46] */
}
/* Content Cards and Similar Display Boxes */
.content-card {
    background-color: #e0f2f1; /* Light teal background from second block */ /* [cite: 47] */
    border-radius: 10px; /* Rounded corners from first block */ /* [cite: 48] */
    padding: 25px; /* Increased padding inside the card from first block */ /* [cite: 49] */
    margin-bottom: 20px; /* Space below the card from first block */ /* [cite: 50] */
    box-shadow: 0 2px 5px rgba(0,0,0,0.08); /* Subtle shadow from first block */ /* [cite: 51] */
    border-left: 5px solid #00796b; /* Accent color left border (teal) from second block */ /* [cite: 52] */
    line-height: 1.7; /* Ensure good line spacing within cards from first block */ /* [cite: 53] */
    color: black; /* Set text color within these boxes to black from second block */ /* [cite: 54] */
}
.summary-box { /* Styling for summary boxes from second block */
    border-left: 5px solid #00796b; /* Teal left border for emphasis */ /* [cite: 55] */
    padding: 1rem; /* Padding inside the box */ /* [cite: 56] */
    margin: 1rem 0; /* Margin around the box */ /* [cite: 57] */
    background-color: #e0f2f1; /* Light teal background */ /* [cite: 58] */
    border-radius: 5px; /* Rounded corners */ /* [cite: 59] */
    box-shadow: 2px 2px 5px rgba(0,0,0,0.1); /* Subtle shadow for depth */ /* [cite: 60] */
    color: black; /* Set text color within these boxes to black */ /* [cite: 61] */
}
/* Styling for Cornell Notes and Mind Map outputs */
.cornell-notes-output, .mind-map-output {
    background-color: #424242; /* Gray background */
    color: white; /* White text */
    padding: 15px;
    border-radius: 8px;
    margin-bottom: 15px;
    border-left: 5px solid #00acc1; /* Accent color (cyan) */
}
.cornell-notes-output h4, .mind-map-output h4 {
    color: #80deea; /* Lighter cyan for headings inside */
    margin-top: 0;
}
.cornell-notes-output p, .mind-map-output p,
.cornell-notes-output li, .mind-map-output li {
    color: white;
}
.child-mode { /* Specific styling for child-mode content from second block */
    background-color: #fff8e1; /* Light yellow for child mode */ /* [cite: 62] */
    padding: 1rem; /* Padding */ /* [cite: 62] */
    border-radius: 10px; /* Rounded corners */ /* [cite: 63] */
    border-left: 5px solid #ffc107; /* Amber border */ /* [cite: 63] */
    color: black; /* Set text color in child mode to black */ /* [cite: 64] */
}
.flashcard { /* Styling for flashcard elements from second block */
    background-color: #f1f8e9; /* Light green background */ /* [cite: 65] */
    border-radius: 8px; /* Rounded corners */ /* [cite: 65] */
    padding: 15px; /* Padding */ /* [cite: 66] */
    margin: 10px 0; /* Margin */ /* [cite: 66] */
    box-shadow: 0 2px 4px rgba(0,0,0,0.1); /* Subtle shadow */ /* [cite: 67] */
    border-left: 4px solid #8bc34a; /* Green border */ /* [cite: 67] */
    color: black; /* Set text color in flashcards to black */ /* [cite: 68] */
}
/* Tabs Styling */
.stTabs [data-baseweb="tab-list"] { /* Target the list of tabs */
    gap: 24px; /* Space between tabs from second block */ /* [cite: 69] */
}
.stTabs [data-baseweb="tab"] { /* Target individual tab elements */
    height: 50px; /* Height of tabs from second block */ /* [cite: 70] */
    white-space: pre-wrap; /* Allow text wrapping in tabs from second block */ /* [cite: 71] */
    background-color: #e0f2f1; /* Light teal for inactive tabs from second block */ /* [cite: 72] */
    border-radius: 8px 8px 0px 0px; /* Rounded top corners from first block */ /* [cite: 73] */
    padding: 10px 15px; /* Padding for tabs (combined) */ /* [cite: 74] */
    color: black !important; /* Set text color for inactive tabs to black from second block */ /* [cite: 75] */
    transition: background-color 0.3s ease, color 0.3s ease; /* Smooth transition */ /* [cite: 76] */
}
.stTabs [aria-selected="true"] { /* Target the currently selected (active) tab */
    background-color: #00796b !important; /* Teal for active tab from second block */ /* [cite: 77] */
    color: white !important; /* Set text color for active tab to white from second block */ /* [cite: 78] */
    font-weight: bold; /* Bold text for active tab from first block */ /* [cite: 79] */
    border-radius: 8px 8px 0 0; /* Ensure rounded top corners */ /* [cite: 80] */
    padding: 10px 15px; /* Ensure consistent padding */ /* [cite: 81] */
}
/* Expander Styling */
.stExpander {
    background-color: #f8f9fa; /* Light background for expanders (from first block - readable on dark main) */ /* [cite: 82] */
    border-radius: 8px; /* Rounded corners from first block */ /* [cite: 83] */
    margin-bottom: 15px; /* Space below expanders from first block */ /* [cite: 84] */
    border: 1px solid #e0e0e0; /* Subtle border from first block */ /* [cite: 85] */
}
.stExpander p, .stExpander div, .stExpander summary {
     color: black !important; /* Ensure text inside expander is black */ /* [cite: 86, 89] */
}
.stExpander div[role="button"] p { /* Style for the expander header text */
    color: #1a374d !important; /* Slightly darker color for header text */ /* [cite: 88] */
}
/* Metrics Styling */
.stMetric { /* Target metric display elements */
    background-color: #e0f2f1; /* Light teal background (using color from second block's palette) */ /* [cite: 90] */
    padding: 15px; /* Padding inside metrics from first block */ /* [cite: 91] */
    border-radius: 8px; /* Rounded corners from first block */ /* [cite: 92] */
    border-left: 4px solid #00796b; /* Accent color left border (teal) */ /* [cite: 93] */
    margin-bottom: 15px; /* Space below metrics from first block */ /* [cite: 94] */
}
.stMetric label, .stMetric div[data-testid="stMetricValue"], .stMetric div[data-testid="stMetricDelta"] {
    color: black !important; /* Set text color in metrics to black from second block */ /* [cite: 95, 96, 97, 98] */
}
/* Responsive Layout Improvements for Smaller Screens */
@media (max-width: 768px) {
    .stTabs [data-baseweb="tab-list"] {
        flex-direction: column; /* Stack tabs vertically on small screens from first block */ /* [cite: 99] */
    }
    .stButton>button {
        width: 100%; /* Full width buttons on small screens from first block */ /* [cite: 100] */
        text-align: center; /* [cite: 101] */
    }
}
/* Additional minor improvements */
.stMarkdown h3 {
    margin-top: 1.5em; /* Space above h3 in markdown from first block */ /* [cite: 102] */
}
.stMarkdown ul, .stMarkdown ol {
    margin-bottom: 1em; /* Space below lists in markdown from first block */ /* [cite: 103] */
}
/* Ensure general markdown text is readable on dark background */
.stMarkdown { /* [cite: 104] */
    color: white;
}
.stMarkdown strong { /* [cite: 104] */
    color: white; /* Ensure bold text is white */
}
.stMarkdown em {
     color: #b2dfdb; /* A slightly lighter teal for emphasis */ /* [cite: 105] */
}
/* Streamlit default text input/area labels - sometimes these need explicit color */
.st-emotion-cache-1x8cf18 e1i5pmtj0 { /* Example class, replace with actual if needed for selectbox label */
   color: white !important; /* [cite: 106] */
}
/* Make sure selectbox label is white */
.stSelectbox label {
    color: white !important;
}
</style>
""", unsafe_allow_html=True)


# Enhanced model loading with progress tracking and fallback options
@st.cache_resource
def load_models() -> Dict[str, Any]:
    """Load and cache all required NLP models with enhanced error handling and progress tracking."""
    try:
        with st.spinner("Loading AI models... This may take a few minutes"):
            t5_tokenizer = T5Tokenizer.from_pretrained("t5-base") # [cite: 107]
            t5_model = T5ForConditionalGeneration.from_pretrained("t5-base") # [cite: 107]
            bert_tokenizer = BertTokenizer.from_pretrained("bert-base-uncased") # [cite: 107]
            bert_model = BertModel.from_pretrained("bert-base-uncased") # [cite: 107]
            child_summarizer = pipeline( # [cite: 108]
                "summarization",
                model="facebook/bart-large-cnn", # [cite: 108]
                tokenizer="facebook/bart-large-cnn" # [cite: 108]
            )
            nlp = spacy.load("en_core_web_sm") # [cite: 109]
            nlp.add_pipe('sentencizer') # [cite: 109]
            transcription_model = whisper.load_model("base") # [cite: 109]
            st.success("All models loaded successfully!") # [cite: 110]
            return {
                't5_tokenizer': t5_tokenizer, # [cite: 110]
                't5_model': t5_model, # [cite: 110]
                'bert_tokenizer': bert_tokenizer, # [cite: 110]
                'bert_model': bert_model, # [cite: 110]
                'child_summarizer': child_summarizer, # [cite: 111]
                'nlp': nlp, # [cite: 111]
                'transcription_model': transcription_model # [cite: 111]
            }
    except Exception as e:
        st.error(f"Model loading failed: {str(e)}") # [cite: 111]
        st.stop()


# Enhanced text extraction with better error handling and file type support
def extract_text_from_file(file_path: str, file_type: str) -> Optional[str]:
    """Universal text extraction with improved error handling and support for more file types.""" # [cite: 112]
    try:
        if file_type == "application/pdf": # [cite: 112]
            reader = PdfReader(file_path) # [cite: 112]
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text()) # [cite: 112]
            return text if text.strip() else None # [cite: 112]
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document": # [cite: 112]
            doc = DocxDocument(file_path) # [cite: 113]
            return "\n".join(para.text for para in doc.paragraphs if para.text) # [cite: 113]
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation": # [cite: 113]
            prs = Presentation(file_path) # [cite: 113]
            text_runs = []
            for slide in prs.slides: # [cite: 113]
                for shape in slide.shapes: # [cite: 113]
                    if hasattr(shape, "text") and shape.text.strip(): # [cite: 114]
                        text_runs.append(shape.text) # [cite: 114]
            return "\n\n".join(text_runs) if text_runs else None # [cite: 114]
        elif file_type == "text/plain": # [cite: 114]
            with open(file_path, 'r', encoding='utf-8') as f: # [cite: 114]
                return f.read() # [cite: 115]
        return None
    except Exception as e:
        st.error(f"File extraction error: {str(e)}") # [cite: 115]
        return None
    finally:
        if 'file_path' in locals() and os.path.exists(file_path): # [cite: 115]
            try:
                os.unlink(file_path) # [cite: 115]
            except Exception as e:
                st.warning(f"Could not delete temporary file {file_path}: {e}") # [cite: 116]


# Advanced text preprocessing with more cleaning options
def preprocess_text(text: str, mode: str = "default") -> str:
    """Enhanced text cleaning with mode-specific rules and better handling of special cases.""" # [cite: 116]
    if not text:
        return "" # [cite: 116]
    text = re.sub(r'\s+', ' ', text).strip() # [cite: 116]
    text = re.sub(r'\[.*?\]|\(.*?\)', '', text) # Remove citations # [cite: 117]
    if mode == "child": # [cite: 117]
        text = re.sub(r'[^\w\s.,!?\'\"\\]', '', text) # Added escaping for backslash # [cite: 117]
        text = re.sub(r'\b\w{15,}\b', '', text) # [cite: 117]
        text = text.lower() # Make it more child-friendly # [cite: 117]
    elif mode == "technical": # [cite: 117]
        text = re.sub(r'[^\w\s.,;:\-+\'"()\[\]{}]', '', text) # [cite: 117]
    return text # [cite: 118]


# Utility function to get BERT embeddings
def get_bert_embeddings(text: str, tokenizer, model):
    """Get the mean-pooled BERT embedding for the input text.""" # [cite: 118]
    if tokenizer is None or model is None: # [cite: 118]
        st.warning("BERT tokenizer or model not loaded.") # [cite: 118]
        return None
    try:
        inputs = tokenizer(text, return_tensors="pt", truncation=True, max_length=512) # Increased max_length # [cite: 118]
        with torch.no_grad(): # [cite: 119]
            outputs = model(**inputs) # [cite: 119]
            last_hidden_state = outputs.last_hidden_state # [cite: 119]
            attention_mask = inputs['attention_mask'] # [cite: 119]
            mask_expanded = attention_mask.unsqueeze(-1).expand(last_hidden_state.size()).float() # [cite: 119]
            sum_embeddings = torch.sum(last_hidden_state * mask_expanded, 1) # [cite: 119]
            sum_mask = torch.clamp(mask_expanded.sum(1), min=1e-9) # [cite: 120]
            mean_pooled = sum_embeddings / sum_mask # [cite: 120]
            return mean_pooled # [cite: 120]
    except Exception as e:
        st.warning(f"BERT embedding generation failed: {e}") # [cite: 120]
        return None


# Enhanced BERT keyword extraction
def extract_keywords_with_bert(text: str, models: dict) -> Dict[str, Dict[str, Any]]:
    """Advanced keyword extraction with semantic grouping, meaning generation, and no fixed limit.""" # [cite: 121]
    if not text or not text.strip(): return {} # [cite: 121]
    if 'bert_tokenizer' not in models or models['bert_tokenizer'] is None or \
       'bert_model' not in models or models['bert_model'] is None or \
       'nlp' not in models or models['nlp'] is None: # [cite: 121]
        st.warning("Required models for keyword extraction not loaded.") # [cite: 122]
        return {}
    try:
        doc_embedding = get_bert_embeddings(text, models['bert_tokenizer'], models['bert_model']) # [cite: 122]
        if doc_embedding is None: return {} # [cite: 122]

        doc_spacy = models['nlp'](text[:100000]) # Increased limit # [cite: 122]
        candidates = list(set( # [cite: 123]
            [chunk.text.lower() for chunk in doc_spacy.noun_chunks if isinstance(chunk, spacy.tokens.span.Span) # [cite: 123]
             and len(chunk.text.split()) <= 5 and len(chunk.text) > 3] # [cite: 123]
        ))
        candidates.extend(list(set( # [cite: 123]
            [token.lemma_.lower() for token in doc_spacy if isinstance(token, spacy.tokens.token.Token) # [cite: 123]
             and token.pos_ in ['NOUN', 'PROPN', 'ADJ', 'VERB'] and not token.is_stop # [cite: 124]
             and not token.is_punct and len(token.lemma_) > 3] # [cite: 124]
        )))
        candidates = list(set(c for c in candidates if len(c.split()) < 6 and len(c) > 2)) # Filter out very long candidates
        if not candidates: return {} # [cite: 124]

        phrase_embeddings = [] # [cite: 125]
        valid_phrases = [] # [cite: 125]
        for phrase in candidates: # [cite: 125]
            phrase_emb = get_bert_embeddings(phrase, models['bert_tokenizer'], models['bert_model']) # [cite: 125]
            if phrase_emb is not None: # [cite: 125]
                phrase_embeddings.append(phrase_emb) # [cite: 125]
                valid_phrases.append(phrase) # [cite: 125]
        if not phrase_embeddings: return {} # [cite: 126]

        phrase_embeddings_tensor = torch.cat(phrase_embeddings, dim=0) # [cite: 126]
        similarities = cosine_similarity( # [cite: 126]
            doc_embedding.cpu().numpy(), # [cite: 126]
            phrase_embeddings_tensor.cpu().numpy() # [cite: 126]
        )[0]
        keywords_with_meanings = {} # [cite: 127]
        sorted_indices = np.argsort(similarities)[::-1] # [cite: 127]

        for i in sorted_indices: # [cite: 127]
            phrase = valid_phrases[i] # [cite: 127]
            score = float(similarities[i]) # [cite: 127]
            if score < 0.3: continue # Threshold to filter less relevant keywords

            context_sentences = []
            if hasattr(doc_spacy, 'sents') and doc_spacy.sents is not None: # [cite: 127]
                 try:
                     context_sentences = [sent.text for sent in doc_spacy.sents if isinstance(sent, spacy.tokens.span.Span) and phrase in sent.text.lower()][:2] # Get up to 2 context sentences # [cite: 128]
                 except Exception as e:
                     st.warning(f"Error processing sentences for keyword context: {e}") # [cite: 129]
            context = " ".join(context_sentences) if context_sentences else "" # [cite: 129]

            pos_tags = []
            if models.get('nlp'): # [cite: 129]
                try:
                    doc_phrase = models['nlp'](phrase) # [cite: 130]
                    pos_tags = [token.pos_ for token in doc_phrase if hasattr(token, 'pos_')] # [cite: 130]
                except Exception as e:
                    st.warning(f"Error processing phrase for POS tags: {e}") # [cite: 130]
            category = "Concept" # [cite: 130]
            if 'VERB' in pos_tags: category = "Action" # [cite: 131]
            elif 'ADJ' in pos_tags: category = "Attribute" # [cite: 131]

            keywords_with_meanings[phrase] = { # [cite: 131]
                'score': score, # [cite: 131]
                'meaning': f"{category} related to: {context[:150]}..." if context else f"Important {category}", # [cite: 132]
                'category': category, # [cite: 132]
                'context': context[:300] + "..." if context else "" # [cite: 132]
            }
            if len(keywords_with_meanings) >= 30: break # Limit to top 30 keywords
        return keywords_with_meanings
    except Exception as e:
        st.warning(f"Advanced keyword extraction failed: {e}") # [cite: 133]
        return {}


# Enhanced summary generation
def generate_advanced_summary(models: dict, text: str, audience: str) -> Dict[str, Any]:
    """Generate comprehensive summaries with multiple note-taking methods and structured formats.""" # [cite: 133]
    if not text or not text.strip(): # [cite: 133]
        return {"error": "Cannot generate summary from empty text."} # [cite: 133]
    try:
        processed_text = preprocess_text(text, "child" if audience == "child" else "default") # [cite: 134]
        base_summary = "Summary could not be generated." # [cite: 134]

        summary_max_length = min(max(750, len(processed_text.split()) // 2), 1024) # Dynamic max length for summary
        summary_min_length = max(150, summary_max_length // 4)


        if processed_text: # [cite: 134]
            if audience == "child": # [cite: 134]
                if models.get('child_summarizer') and callable(models['child_summarizer']): # [cite: 135]
                    try:
                        summary_output = models['child_summarizer']( # [cite: 135]
                            processed_text[:4096], # [cite: 135]
                            max_length=200, # [cite: 136]
                            min_length=50, # [cite: 136]
                            do_sample=True, # [cite: 136]
                            temperature=0.8, # [cite: 137]
                            top_p=0.95 # [cite: 137]
                        )
                        if isinstance(summary_output, list) and len(summary_output) > 0 and isinstance(summary_output[0], dict) and 'summary_text' in summary_output[0]: # [cite: 138]
                            base_summary = summary_output[0]['summary_text'] # [cite: 138]
                        else:
                            st.warning("Child summarizer returned unexpected output format or is empty.") # [cite: 139]
                            base_summary = "Child summary generation failed due to unexpected output." # [cite: 139]
                    except Exception as e:
                        st.warning(f"Child summary generation failed: {e}") # [cite: 140]
                        base_summary = "Child summary generation failed." # [cite: 140]
                else:
                    st.warning("Child summarizer model not loaded or not callable.") # [cite: 140]
                    base_summary = "Child summarizer not available." # [cite: 141]
            else: # student, researcher, expert
                if models.get('t5_tokenizer') and models.get('t5_model'): # [cite: 141]
                    try:
                        prompts = { # [cite: 142]
                            "student": "Summarize the following text for a high school student. Focus on key concepts, provide bullet points for main ideas, and explain any complex terms simply: ", # [cite: 142]
                            "researcher": "Create a detailed research summary. Include key findings, methodology (if apparent), limitations (if any), and implications. Structure with headings and bullet points where appropriate: ", # [cite: 142]
                            "expert": "Provide a concise yet comprehensive technical summary. Focus on novel contributions, technical analysis, and potential future research directions or applications. Use precise language: " # [cite: 143]
                        }
                        prompt_text = prompts.get(audience, "Summarize in detail with bullet points and insights: ") + processed_text[:10000] # Increased input limit # [cite: 143]
                        inputs = models['t5_tokenizer'].encode( # [cite: 144]
                            prompt_text,
                            return_tensors="pt",
                            max_length=2048, # Increased model max input length # [cite: 144]
                            truncation=True # [cite: 145]
                        )
                        outputs = models['t5_model'].generate( # [cite: 145]
                            inputs, # [cite: 146]
                            max_length=summary_max_length, # Dynamic and increased # [cite: 146]
                            min_length=summary_min_length, # Dynamic and increased # [cite: 146]
                            num_beams=5, # Increased beams # [cite: 146]
                            early_stopping=True, # [cite: 147]
                            temperature=0.7 if audience == "student" else 0.65, # [cite: 147]
                            no_repeat_ngram_size=3 # Add to reduce repetition
                        )
                        if isinstance(outputs, torch.Tensor) and outputs.shape[0] > 0: # [cite: 148]
                            base_summary = models['t5_tokenizer'].decode(outputs[0], skip_special_tokens=True) # [cite: 148]
                        else:
                            st.warning("T5 model generated unexpected output format or is empty.") # [cite: 149]
                            base_summary = "Advanced summary generation failed due to unexpected output." # [cite: 149]
                    except Exception as e:
                        st.warning(f"T5 summary generation failed: {e}") # [cite: 150]
                        base_summary = "Advanced summary generation failed." # [cite: 150]
                else:
                    st.warning("T5 model or tokenizer not loaded.") # [cite: 150]
                    base_summary = "Advanced summarization model not available." # [cite: 151]
        else:
            st.warning("Processed text is empty, skipping summary generation.") # [cite: 151]
            base_summary = "No text to summarize." # [cite: 151]

        doc = None # [cite: 151]
        sentences = [] # [cite: 151]
        entities = [] # [cite: 152]
        try:
            if processed_text and models.get('nlp'): # [cite: 152]
                 doc = models['nlp'](processed_text[:100000]) # Increased limit # [cite: 152]
                 sentences = [sent for sent in doc.sents if isinstance(sent, spacy.tokens.span.Span) and len(sent.text.split()) > 5] # [cite: 153]
                 entities = [(ent.text, ent.label_) for ent in doc.ents if isinstance(ent, spacy.tokens.span.Span)] # [cite: 153]
            elif processed_text and not models.get('nlp'): # [cite: 154]
                 st.warning("spaCy model not loaded for detailed analysis.") # [cite: 154]
            else:
                 st.info("No processed text for detailed analysis (notes, etc.).") # [cite: 154]
        except Exception as e:
            st.warning(f"spaCy processing failed: {e}") # [cite: 154]

        cornell_notes = {"Cue": [], "Notes": [], "Summary": []} # [cite: 155]
        if sentences: # [cite: 155]
             for i, sent in enumerate(sentences[:20]): # Increased limit # [cite: 155]
                 cornell_notes["Cue"].append(f"Key Point {i+1}") # [cite: 156]
                 cornell_notes["Notes"].append(sent.text) # [cite: 156]
                 summary_snippet = ' '.join(sent.text.split()[:7]) + "..." if sent.text else "..." # [cite: 156]
                 cornell_notes["Summary"].append(f"Central theme: {summary_snippet}") # [cite: 156]

        mind_map = { # [cite: 157]
            "Central Idea": base_summary.split('.')[0] if base_summary and '.' in base_summary else base_summary[:min(70, len(base_summary))] + "..." if base_summary else "Document Analysis", # [cite: 157]
            "Main Branches": [], # [cite: 158]
            "Sub-branches": defaultdict(list) # [cite: 158]
        }
        if entities: # [cite: 158]
             for i, (ent_text, label) in enumerate(entities[:7]): # Increased branches # [cite: 158]
                 mind_map["Main Branches"].append(f"{label}: {ent_text}") # [cite: 159]
                 if sentences: # [cite: 159]
                     related_sents = [sent.text for sent in sentences if ent_text.lower() in sent.text.lower()][:3] # Increased sub-branches # [cite: 159]
                     for sent_text in related_sents: # [cite: 159]
                         mind_map["Sub-branches"][f"Branch {i+1}"].append(sent_text) # [cite: 160]

        linear_notes = [] # [cite: 160]
        if sentences: # [cite: 160]
            for i, sent in enumerate(sentences[:25]): # Increased limit # [cite: 160]
                linear_notes.append(f"{i+1}. {sent.text}") # [cite: 160]

        concept_map_suggestions = [] # [cite: 161]
        concepts = [chunk.text for chunk in (doc.noun_chunks if doc and hasattr(doc, 'noun_chunks') else []) if isinstance(chunk, spacy.tokens.span.Span) and len(chunk.text.split()) > 1][:7] # [cite: 161]
        if len(concepts) >= 2: # [cite: 161]
            concept_map_suggestions.append(f"Consider a Concept Map showing relationships between: {', '.join(concepts)}") # [cite: 162]
            if len(concepts) >= 3: # [cite: 162]
                 concept_map_suggestions.append(f"Explore how '{concepts[0]}' influences '{concepts[1]}' and '{concepts[2]}'.") # [cite: 162]
            if len(concepts) >= 4:
                concept_map_suggestions.append(f"Illustrate connections: '{concepts[0]}' -> '{concepts[1]}' -> '{concepts[2]}' -> '{concepts[3]}'. (e.g., as a flowchart or sequence)")


        return { # [cite: 162]
            "base_summary": base_summary, # [cite: 162]
            "cornell_notes": cornell_notes, # [cite: 163]
            "mind_map": mind_map, # [cite: 163]
            "linear_notes": linear_notes, # [cite: 163]
            "concept_map_suggestions": concept_map_suggestions, # [cite: 163]
            "key_phrases": [chunk.text for chunk in (doc.noun_chunks if doc and hasattr(doc, 'noun_chunks') else []) if isinstance(chunk, spacy.tokens.span.Span)][:20], # [cite: 163]
            "entities": entities[:20] # [cite: 163]
        }
    except Exception as e:
        st.error(f"An unexpected error occurred during advanced summary generation: {str(e)}") # [cite: 164]
        return {"error": f"Summary generation failed: {str(e)}"} # [cite: 164]


# Enhanced study materials generation
def generate_advanced_study_materials(text: str, nlp_model: Any) -> Dict[str, Any]:
    """Generate comprehensive study materials with abstractive content and contextual understanding.""" # [cite: 164]
    if not text or not text.strip(): # [cite: 164]
        return {} # [cite: 165]
    if nlp_model is None: # [cite: 165]
        st.warning("spaCy model not loaded for study materials generation.") # [cite: 165]
        return {}
    try:
        doc = nlp_model(text[:100000]) # Increased limit # [cite: 165]
        sentences = [sent for sent in doc.sents if isinstance(sent, spacy.tokens.span.Span)] # [cite: 165]
        flashcards = [] # [cite: 166]
        important_concepts = [chunk.text for chunk in (doc.noun_chunks if hasattr(doc, 'noun_chunks') else []) if isinstance(chunk, spacy.tokens.span.Span) and # [cite: 166]
                              len(chunk.text.split()) <= 3 and # [cite: 166]
                              any(t.pos_ in ['NOUN', 'PROPN'] for t in chunk)][:15] # Increased flashcards # [cite: 167]
        for concept in important_concepts: # [cite: 167]
            context_sents = [sent.text for sent in sentences if concept.lower() in sent.text.lower()] # [cite: 167]
            context = random.choice(context_sents) if context_sents else "" # [cite: 167]
            question = f"What is '{concept}' and its significance in this context?" # [cite: 168]
            answer = f"'{concept}' is relevant as it pertains to: {context[:200]}..." if context else f"'{concept}' is an important concept mentioned in the text." # [cite: 168]
            flashcards.append(f"Q: {question}\nA: {answer}") # [cite: 168]

        questions = [] # [cite: 169]
        concepts = [chunk.text for chunk in (doc.noun_chunks if hasattr(doc, 'noun_chunks') else []) if isinstance(chunk, spacy.tokens.span.Span) and len(chunk.text.split()) > 1][:12] # [cite: 169]
        if len(concepts) >= 2: # [cite: 169]
            for i in range(min(8, len(concepts) - 1)): # Increased questions # [cite: 169]
                questions.append( # [cite: 170]
                    f"Explain the relationship or interaction between '{concepts[i]}' and '{concepts[i+1]}' as presented in the text." # [cite: 170]
                )
        if len(concepts) >=3:
            questions.append(f"How might '{concepts[0]}' influence or be influenced by '{concepts[2]}', based on the provided information?")

        diagrams = [] # [cite: 170]
        if len(concepts) >= 3: # [cite: 170]
            diagrams.append(f"Concept Map: Create a visual map connecting key concepts such as '{concepts[0]}', '{concepts[1]}', and '{concepts[2]}', showing their primary relationships (e.g., hierarchical, causal, associative).") # [cite: 171]
        diagrams.append("Timeline: If the text describes events or a process over time, construct a timeline highlighting major milestones or steps.") # [cite: 171]
        diagrams.append("Venn Diagram: If there are two or more central themes/entities (e.g., '{concepts[0]}' vs '{concepts[1]}' if applicable), use a Venn diagram to illustrate their similarities and differences.") # [cite: 171]
        diagrams.append("Flowchart: For any process, workflow, or decision-making sequence described, design a flowchart to map out the steps and connections.") # [cite: 171]
        diagrams.append("SWOT Analysis: If the text discusses a strategy, project, or entity, consider a SWOT analysis (Strengths, Weaknesses, Opportunities, Threats).")


        study_plan = [ # [cite: 171]
            f"Day 1: Initial Read-Through & Highlighting. Focus on understanding the main arguments and identify core concepts like '{important_concepts[0] if important_concepts else 'main topics'}', '{important_concepts[1] if len(important_concepts)>1 else 'secondary topics'}'.", # [cite: 172]
            "Day 2: Deep Dive with Summary & Notes. Review the generated 'Comprehensive Summary', 'Cornell Notes', and 'Linear Notes'. Actively engage by adding your own annotations.", # [cite: 172]
            "Day 3: Concept Reinforcement with Flashcards. Use the generated 'Flashcards' to test your recall and understanding of key terms and their context.", # [cite: 172]
            "Day 4: Visualize Relationships. Sketch one of the 'Diagram Suggestions' (e.g., a concept map or flowchart) to solidify your understanding of connections between ideas.", # [cite: 172]
            "Day 5: Active Recall & Questioning. Attempt to answer the 'Practice Questions' without looking at the notes. Identify areas needing more review.",
            "Day 7: Spaced Repetition Review. Briefly review the summary, your sketched diagrams, and the flashcards for challenging concepts.", # [cite: 173]
            "Weekly: Quick Re-engagement. Spend 15-20 minutes reviewing the most critical concepts or a section you found difficult to reinforce long-term memory." # [cite: 173]
        ]
        return { # [cite: 173]
            "flashcards": flashcards, # [cite: 173]
            "questions": questions, # [cite: 173]
            "diagrams": diagrams, # [cite: 173]
            "plan": study_plan, # [cite: 173]
            "key_concepts": important_concepts, # [cite: 174]
            "summary": " ".join(sent.text for sent in sentences[:7]) if sentences else "" # [cite: 174]
        }
    except Exception as e:
        st.warning(f"Study materials generation failed: {e}") # [cite: 174]
        return {}

# PDF generation function
def create_pdf_report(content: Dict[str, Any], title: str) -> BytesIO: # [cite: 175]
    """Generate a professional PDF report from analysis content.""" # [cite: 175]
    buffer = BytesIO() # [cite: 175]
    doc = SimpleDocTemplate(buffer, pagesize=letter) # [cite: 175]
    styles = getSampleStyleSheet() # [cite: 175]

    # Modify existing styles
    if 'Heading1' in styles: styles['Heading1'].fontSize = 18; styles['Heading1'].leading = 22; styles['Heading1'].textColor = colors.HexColor("#004D40"); styles['Heading1'].spaceAfter = 12 # [cite: 175]
    if 'Heading2' in styles: styles['Heading2'].fontSize = 14; styles['Heading2'].leading = 18; styles['Heading2'].textColor = colors.HexColor("#00796B"); styles['Heading2'].spaceBefore = 12; styles['Heading2'].spaceAfter = 6 # [cite: 176]
    if 'Heading3' in styles: styles['Heading3'].fontSize = 12; styles['Heading3'].leading = 16; styles['Heading3'].textColor = colors.HexColor("#00796B"); styles['Heading3'].spaceBefore = 8; styles['Heading3'].spaceAfter = 4
    if 'Bullet' in styles: styles['Bullet'].fontSize = 10; styles['Bullet'].leading = 14; styles['Bullet'].leftIndent = 20 # [cite: 176]
    if 'BodyText' in styles: styles['BodyText'].fontSize = 10; styles['BodyText'].leading = 14; styles['BodyText'].alignment = TA_JUSTIFY # [cite: 177]

    # Custom styles
    styles.add(ParagraphStyle(name='JustifyBody', parent=styles['BodyText'], alignment=TA_JUSTIFY)) # [cite: 178]
    styles.add(ParagraphStyle(name='CenteredSmall', parent=styles['BodyText'], alignment=TA_CENTER, fontSize=9, textColor=colors.grey))
    styles.add(ParagraphStyle(name='FlashcardQ', parent=styles['BodyText'], textColor=colors.HexColor("#00796B"), spaceBefore=6))
    styles.add(ParagraphStyle(name='FlashcardA', parent=styles['JustifyBody'], leftIndent=10, spaceAfter=6))
    def create_pdf_report(analysis, source_name):
        styles = getSampleStyleSheet()

        # Check if 'Code' style already exists before adding it
        if 'Code' not in styles:
            styles.add(ParagraphStyle(
                name='Code',
                parent=styles['BodyText'],
                fontName='Courier',
                fontSize=9,
                textColor=colors.darkslategray,
                backColor=colors.whitesmoke,
                firstLineIndent=0,
                leftIndent=6,
                rightIndent=6,
                borderPadding=5,
                borderRadius=2,
                borderColor=colors.lightgrey,
                borderWidth=1,
                padding=5 ))
        # ))styles.add(ParagraphStyle(name='Code', parent=styles['BodyText'], fontName='Courier', fontSize=9, textColor=colors.darkslategray, backColor=colors.whitesmoke, firstLineIndent=0, leftIndent=6, rightIndent=6, borderPadding=5, borderRadius=2, borderColor=colors.lightgrey, borderWidth=1, padding=5))


    elements = [] # [cite: 178]
    elements.append(Paragraph(title, styles['Heading1'])) # [cite: 178]
    elements.append(Paragraph(f"Report Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['CenteredSmall'])) # [cite: 179]
    elements.append(Spacer(1, 0.25 * inch))

    def add_section(heading_text, content_data, style_key='JustifyBody', is_list=False, list_prefix="• "): # [cite: 179]
        if content_data:
            elements.append(Paragraph(heading_text, styles['Heading2'])) # [cite: 179]
            if is_list:
                for item in content_data:
                    if isinstance(item, str): elements.append(Paragraph(f"{list_prefix}{item}", styles[style_key])) # [cite: 180]
            elif isinstance(content_data, str):
                elements.append(Paragraph(content_data, styles[style_key])) # [cite: 180]
            elements.append(Spacer(1, 0.15 * inch))

    add_section("Comprehensive Summary", content.get("base_summary")) # [cite: 179]

    if content.get('cornell_notes') and content['cornell_notes'].get('Cue'): # [cite: 180]
        elements.append(Paragraph("Cornell Method Notes", styles['Heading2'])) # [cite: 180]
        data = [["<b>Cue / Keywords</b>", "<b>Notes / Details</b>", "<b>Summary / Main Idea</b>"]]
        for cue, note, summ in zip(content['cornell_notes'].get('Cue', []), content['cornell_notes'].get('Notes', []), content['cornell_notes'].get('Summary', [])): # [cite: 182]
            if isinstance(cue, str) and isinstance(note, str) and isinstance(summ, str): # [cite: 183]
                 data.append([Paragraph(cue, styles['BodyText']), Paragraph(note, styles['JustifyBody']), Paragraph(summ, styles['JustifyBody'])])
        if len(data) > 1:
            table = Table(data, colWidths=[1.5*inch, 3.5*inch, 2*inch])
            table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#00796B")),
                ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0,0), (-1,0), 12),
                ('BACKGROUND', (0,1), (-1,-1), colors.HexColor("#E0F2F1")),
                ('GRID', (0,0), (-1,-1), 1, colors.darkgrey),
                ('LEFTPADDING', (0,0), (-1,-1), 6),
                ('RIGHTPADDING', (0,0), (-1,-1), 6),
                ('TOPPADDING', (0,0), (-1,-1), 6),
                ('BOTTOMPADDING', (0,0), (-1,-1), 6),
            ]))
            elements.append(table)
            elements.append(Spacer(1, 0.15 * inch))


    add_section("Key Concepts / Phrases", content.get("key_phrases") or content.get("key_concepts"), 'Bullet', is_list=True) # [cite: 183]
    add_section("Named Entities", [f"<b>{ent_tuple[0]}:</b> {ent_tuple[1]}" for ent_tuple in content.get("entities", []) if isinstance(ent_tuple, tuple) and len(ent_tuple) == 2 and isinstance(ent_tuple[0], str) and isinstance(ent_tuple[1], str)], 'BodyText', is_list=True, list_prefix="") # [cite: 200]


    if content.get('flashcards'): # [cite: 184]
        elements.append(Paragraph("Flashcards", styles['Heading2'])) # [cite: 185]
        for card in content['flashcards']: # [cite: 185]
            if isinstance(card, str) and "\nA: " in card: # [cite: 186]
                q, a = card.split("\nA: ", 1) # [cite: 186]
                elements.append(Paragraph(q.replace("Q: ", "<b>Q: </b>", 1), styles['FlashcardQ'])) # [cite: 186]
                elements.append(Paragraph(a, styles['FlashcardA'])) # [cite: 186]
        elements.append(Spacer(1, 0.15 * inch)) # [cite: 186]

    if content.get('mind_map') and content['mind_map'].get('Central Idea'): # [cite: 187]
        elements.append(Paragraph("Mind Map Structure Suggestion", styles['Heading2'])) # [cite: 187]
        mm = content['mind_map']
        if isinstance(mm['Central Idea'], str): elements.append(Paragraph(f"<b>Central Idea:</b> {mm['Central Idea']}", styles['BodyText'])) # [cite: 188]
        if mm.get('Main Branches'): # [cite: 188]
            elements.append(Paragraph("<b>Main Branches:</b>", styles['BodyText'])) # [cite: 188]
            for branch in mm['Main Branches']: # [cite: 188]
                if isinstance(branch, str): elements.append(Paragraph(f"- {branch}", styles['Bullet'])) # [cite: 189]
            for i, branch_name in enumerate(mm.get('Main Branches', [])): # [cite: 189]
                sub_branches = mm['Sub-branches'].get(f"Branch {i+1}") # [cite: 189]
                if sub_branches and isinstance(sub_branches, list): # [cite: 189]
                    elements.append(Paragraph(f"  <u>Sub-branches for {branch_name}:</u>", styles['BodyText'])) # [cite: 190]
                    for sub_branch in sub_branches: # [cite: 190]
                        if isinstance(sub_branch, str): elements.append(Paragraph(f"    • {sub_branch}", styles['Bullet'])) # [cite: 190]
        elements.append(Spacer(1, 0.15 * inch)) # [cite: 191]

    add_section("Linear Notes", content.get("linear_notes"), 'BodyText', is_list=True, list_prefix="") # [cite: 191]
    add_section("Concept Map Suggestions", content.get("concept_map_suggestions"), 'Bullet', is_list=True) # [cite: 192]
    add_section("Practice Questions", content.get("questions"), 'BodyText', is_list=True, list_prefix=f"{i+1}. ") # [cite: 194]
    add_section("Diagram Suggestions", content.get("diagrams"), 'Bullet', is_list=True) # [cite: 196]
    add_section("Study Plan", content.get("plan"), 'BodyText', is_list=True, list_prefix=f"{i+1}. ") # [cite: 197]


    try:
        doc.build(elements) # [cite: 201]
    except Exception as e:
        st.error(f"Error building PDF: {e}")
        # Fallback: build with only title if complex elements fail
        elements = [Paragraph(title, styles['h1']), Paragraph(f"Report generation encountered an error with content: {e}", styles['BodyText'])]
        doc.build(elements)

    buffer.seek(0) # [cite: 201]
    return buffer # [cite: 201]


# Text-to-speech function
def play_audio_summary(text: str):
    if not text.strip():
        st.toast("No summary text to play.", icon="🔈")
        return
    if not pygame.mixer.get_init():
        st.warning("Audio player not initialized. Cannot play summary.")
        return

    try:
        with st.spinner("Generating audio summary..."):
            tts = gTTS(text=text, lang='en', slow=False)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
                tts.save(fp.name)
                audio_file_path = fp.name

        pygame.mixer.music.load(audio_file_path)
        pygame.mixer.music.play()
        st.toast("Playing summary...", icon="🔊")

        # Keep app alive while playing, then cleanup
        # This part can be tricky with Streamlit's execution model.
        # A simple approach:
        while pygame.mixer.music.get_busy():
            pygame.time.Clock().tick(10) # Check every 100ms

    except ConnectionError:
        st.error("Failed to connect to Google Text-to-Speech service. Please check your internet connection.")
    except Exception as e:
        st.error(f"Error generating or playing audio: {e}")
    finally:
        if 'audio_file_path' in locals() and os.path.exists(audio_file_path):
            try:
                # Ensure pygame mixer is stopped before deleting
                if pygame.mixer.music.get_busy():
                    pygame.mixer.music.stop()
                # Short delay to ensure file lock is released
                pygame.time.wait(500)
                os.unlink(audio_file_path)
            except Exception as e_del:
                st.warning(f"Could not delete temporary audio file {audio_file_path}: {e_del}")


# Enhanced display function
def display_advanced_analysis_results(text_to_analyze: str, models: dict, source_name: str = "Uploaded Content"): # [cite: 201]
    """Display comprehensive analysis results with all new features.""" # [cite: 201]
    if not text_to_analyze or not text_to_analyze.strip():
        st.warning(f"No text content provided from {source_name} to analyze.")
        return

    st.subheader(f"Advanced Analysis Results for: {source_name}") # [cite: 201]
    processed_text = preprocess_text(text_to_analyze) # [cite: 201]

    tab1, tab2, tab3, tab4 = st.tabs(["📚 Summaries & Notes", "🔑 Keywords", "📝 Study Aids", "📊 Visualizations"]) # [cite: 202]

    with tab1: # [cite: 202]
        st.markdown("#### Advanced Summarization & Note-Taking") # [cite: 202]
        audience_options = ["student", "researcher", "expert", "child"] # [cite: 202]
        # Check if session state has a pre-selected audience, else default
        default_audience_index = audience_options.index(st.session_state.get('audience_choice', "student"))
        audience = st.selectbox("Select target audience for summary:", audience_options, index=default_audience_index, key=f"audience_sb_{source_name}") # [cite: 202]
        st.session_state.audience_choice = audience # Store choice for persistence across reruns

        if st.button(f"Generate {audience.capitalize()} Summary & Notes", key=f"gen_summary_btn_{source_name}"): # [cite: 202]
            if processed_text: # [cite: 202]
                with st.spinner(f"Generating advanced {audience} summary and notes... This may take a moment."): # [cite: 203]
                    analysis_results = generate_advanced_summary(models, text_to_analyze, audience) # [cite: 203]
                    st.session_state[f'analysis_results_{source_name}'] = analysis_results # Cache results in session state
            else:
                 st.warning("No text to summarize.") # [cite: 218]

        # Display results if they exist in session state
        if f'analysis_results_{source_name}' in st.session_state:
            analysis = st.session_state[f'analysis_results_{source_name}']
            if 'error' in analysis: # [cite: 203]
                st.error(analysis['error']) # [cite: 203]
            else:
                st.markdown("### Comprehensive Summary") # [cite: 204]
                base_summary_text = analysis.get("base_summary", "Summary not available.") # [cite: 204]
                st.markdown(f'<div class="content-card">{base_summary_text}</div>', unsafe_allow_html=True) # [cite: 204]

                if base_summary_text != "Summary not available." and base_summary_text.strip():
                    if st.button("🔈 Listen to Summary", key=f"tts_btn_{source_name}"):
                        play_audio_summary(base_summary_text)

                if analysis.get('cornell_notes') and analysis['cornell_notes'].get('Cue'): # [cite: 205]
                     with st.expander("Cornell Method Notes (View)", expanded=False): # [cite: 205]
                         cornell_html = '<div class="cornell-notes-output"><h4>Cornell Notes</h4>' # [cite: 205]
                         cornell_html += "<table><tr><th>Cue</th><th>Notes</th><th>Summary Snippet</th></tr>"
                         for cue, note, summ in zip(analysis['cornell_notes'].get('Cue', []), analysis['cornell_notes'].get('Notes', []), analysis['cornell_notes'].get('Summary',[])): # [cite: 208]
                             if isinstance(cue, str) and isinstance(note, str) and isinstance(summ, str): # [cite: 208]
                                 cornell_html += f"<tr><td><strong>{cue}</strong></td><td>{note}</td><td><em>{summ}</em></td></tr>" # [cite: 209]
                         cornell_html += "</table></div>"
                         st.markdown(cornell_html, unsafe_allow_html=True)

                if analysis.get('mind_map') and analysis['mind_map'].get('Central Idea'): # [cite: 209]
                    with st.expander("Mind Map Structure Suggestion (View)", expanded=False): # [cite: 210]
                        mm = analysis['mind_map'] # [cite: 210]
                        mind_map_html = '<div class="mind-map-output"><h4>Mind Map Outline</h4>' # [cite: 210]
                        mind_map_html += f"<p><strong>Central Idea:</strong> {mm['Central Idea']}</p>" # [cite: 210]
                        if mm.get('Main Branches'): # [cite: 210]
                            mind_map_html += "<p><strong>Main Branches:</strong></p><ul>" # [cite: 211]
                            for i, branch in enumerate(mm['Main Branches']): # [cite: 211]
                                if isinstance(branch, str): mind_map_html += f"<li>{branch}</li>" # [cite: 212]
                                sub_branches = mm['Sub-branches'].get(f"Branch {i+1}") # [cite: 213]
                                if sub_branches and isinstance(sub_branches, list): # [cite: 213]
                                    mind_map_html += "<ul>" # [cite: 214]
                                    for sub_branch in sub_branches: # [cite: 214]
                                        if isinstance(sub_branch, str): mind_map_html += f"<li><em>{sub_branch[:150]}...</em></li>" # [cite: 216]
                                    mind_map_html += "</ul>"
                            mind_map_html += "</ul>"
                        mind_map_html += "</div>"
                        st.markdown(mind_map_html, unsafe_allow_html=True)

                if analysis.get("linear_notes"):
                    with st.expander("Linear Notes (View)", expanded=False):
                        st.markdown(f'<div class="content-card"><h4>Linear Notes</h4>{"<br>".join(analysis["linear_notes"])}</div>', unsafe_allow_html=True)


                pdf_buffer = create_pdf_report(analysis, f"Analysis Report - {source_name}") # [cite: 216]
                st.download_button( # [cite: 216]
                     label="📥 Download Full Analysis Report (PDF)", # [cite: 217]
                     data=pdf_buffer, # [cite: 217]
                     file_name=f"analysis_report_{source_name.replace(' ','_').split('.')[0]}.pdf", # [cite: 217]
                     mime="application/pdf", # [cite: 218]
                     key=f"dl_report_btn_{source_name}"
                )

    with tab2: # [cite: 218]
        st.markdown("#### Semantic Keywords") # [cite: 218]
        if st.button("Extract Semantic Keywords", key=f"extract_kw_btn_{source_name}"): # [cite: 218]
            if processed_text: # [cite: 219]
                with st.spinner("Analyzing text for semantic keywords... This may take a moment."): # [cite: 219]
                    keywords = extract_keywords_with_bert(processed_text, models) # [cite: 219]
                    st.session_state[f'keywords_{source_name}'] = keywords # Cache
            else:
                 st.warning("No text to extract keywords from.") # [cite: 226]

        if f'keywords_{source_name}' in st.session_state:
            keywords = st.session_state[f'keywords_{source_name}']
            if keywords: # [cite: 219]
                st.markdown("### Key Concepts with Contextual Meaning") # [cite: 220]
                for kw, data in keywords.items(): # [cite: 220]
                    if isinstance(kw, str) and isinstance(data, dict): # [cite: 220]
                         with st.expander(f"{kw} ({data.get('category', 'N/A')}) - Score: {data.get('score', 0.0):.2f}"): # [cite: 221]
                            st.markdown(f"**Meaning Suggestion:** {data.get('meaning', 'N/A')}") # [cite: 221]
                            if data.get('context'): st.markdown(f"**Context:** _{data['context']}_") # [cite: 222]

                if keywords: # [cite: 222]
                    keywords_df = pd.DataFrame.from_dict(keywords, orient='index') # [cite: 223]
                    csv = keywords_df.to_csv(index=True).encode('utf-8') # [cite: 223]
                    st.download_button( # [cite: 223]
                        label="📥 Download Keywords (CSV)", # [cite: 224]
                        data=csv, # [cite: 224]
                        file_name=f"keywords_{source_name.replace(' ','_').split('.')[0]}.csv", # [cite: 224]
                        mime="text/csv", # [cite: 225]
                        key=f"dl_kw_btn_{source_name}"
                    )
            else:
                st.info("No significant keywords could be extracted or the process failed.") # [cite: 226]


    with tab3: # [cite: 226]
        st.markdown("#### Study Aids Generator") # [cite: 227]
        if st.button("Generate Study Materials", key=f"gen_study_btn_{source_name}"): # [cite: 227]
            if processed_text: # [cite: 227]
                with st.spinner("Creating advanced study materials... This may take a moment."): # [cite: 227]
                    materials = generate_advanced_study_materials(processed_text, models.get('nlp')) # [cite: 227]
                    st.session_state[f'study_materials_{source_name}'] = materials # Cache
            else:
              st.warning("No text to generate study materials from.") # [cite: 243]

        if f'study_materials_{source_name}' in st.session_state:
            materials = st.session_state[f'study_materials_{source_name}']
            if materials: # [cite: 228]
                if materials.get('flashcards'): # [cite: 228]
                    st.markdown("### Flashcards") # [cite: 228]
                    for card in materials['flashcards']: # [cite: 229]
                         if isinstance(card, str) and "\nA: " in card: # [cite: 229]
                             q, a = card.split("\nA: ", 1) # [cite: 230]
                             with st.expander(q): st.write(a) # [cite: 230]
                else: st.info("No flashcards generated.") # [cite: 231]

                if materials.get('questions'): # [cite: 232]
                    st.markdown("### Practice Questions") # [cite: 232]
                    for q_idx, q_text in enumerate(materials['questions']): # [cite: 232]
                         if isinstance(q_text, str): st.markdown(f'<div class="flashcard">{q_idx+1}. {q_text}</div>', unsafe_allow_html=True) # [cite: 233]
                else: st.info("No practice questions generated.") # [cite: 234]

                if materials.get('diagrams'): # [cite: 234]
                    st.markdown("### Diagram Suggestions") # [cite: 235]
                    for diagram in materials['diagrams']: # [cite: 235]
                         if isinstance(diagram, str): st.markdown(f"- {diagram}") # [cite: 236]
                else: st.info("No diagram suggestions generated.") # [cite: 237]

                if materials.get('plan'): # [cite: 237]
                    st.markdown("### Suggested Study Plan") # [cite: 237]
                    for item_idx, item_text in enumerate(materials['plan']): # [cite: 238]
                         if isinstance(item_text, str): st.markdown(f"{item_idx+1}. {item_text}") # [cite: 239]
                else: st.info("No study plan generated.") # [cite: 239]

                pdf_buffer_study = create_pdf_report(materials, f"Study Materials - {source_name}") # [cite: 240]
                st.download_button( # [cite: 240]
                   label="📥 Download Study Materials (PDF)", # [cite: 241]
                   data=pdf_buffer_study, # [cite: 241]
                   file_name=f"study_materials_{source_name.replace(' ','_').split('.')[0]}.pdf", # [cite: 241]
                   mime="application/pdf", # [cite: 241]
                   key=f"dl_study_btn_{source_name}"
                )
            else:
                st.warning("Could not generate study materials from this content.") # [cite: 242]


    with tab4: # [cite: 243]
        st.markdown("#### Text Visualizations") # [cite: 243]
        if st.button("Generate Word Cloud", key=f"wc_btn_{source_name}"): # [cite: 243]
            if processed_text: # [cite: 243]
                with st.spinner("Creating word cloud..."): # [cite: 243]
                    fig_wc = create_wordcloud(processed_text) # [cite: 243]
                    if fig_wc: 
                        st.pyplot(fig_wc) # [cite: 244]
                    else: 
                        st.warning("Could not generate word cloud.") # [cite: 244]
            else: 
                st.warning("No text to generate word cloud from.") # [cite: 245]

        if st.button("Generate Entity Network", key=f"en_btn_{source_name}"): # [cite: 245]
            if processed_text: # [cite: 245]
                with st.spinner("Analyzing entities for network... This might take a moment."): # [cite: 245]
                    if models.get('nlp'): # [cite: 246]
                        doc = models['nlp'](processed_text[:100000]) # Limit for performance # [cite: 246]
                        entities = [(ent.text, ent.label_) for ent in doc.ents if isinstance(ent, spacy.tokens.span.Span)] # [cite: 246]
                        if entities: # [cite: 247]
                            fig_en = create_entity_network(entities, doc) # [cite: 247]
                            if fig_en: 
                                st.plotly_chart(fig_en, use_container_width=True) # [cite: 248]
                            else: 
                                st.warning("Could not generate entity network, not enough linked entities found.") # [cite: 248]
                        else: 
                            st.info("No named entities found to visualize.") # [cite: 249]
                    else: 
                        st.warning("spaCy model not loaded for entity network visualization.") # [cite: 249]
            else:
                st.warning("No text to generate entity network from.") # [cite: 250]


# Utility function to create a word cloud
def create_wordcloud(text: str): # [cite: 250]
    """Generate a word cloud from the input text and return a matplotlib figure.""" # [cite: 250]
    if not text or not text.strip(): return None # [cite: 250]
    try:
        stop_words = set(stopwords.words('english')) # [cite: 250]
        wordcloud_obj = WordCloud( # [cite: 250]
            width=800, height=400, background_color='white', stopwords=stop_words, # [cite: 251]
            colormap='viridis', min_font_size=10, prefer_horizontal=0.9
        ).generate(text) # [cite: 251]
        fig, ax = plt.subplots(figsize=(12, 6)) # [cite: 251]
        ax.imshow(wordcloud_obj, interpolation='bilinear') # [cite: 251]
        ax.axis('off') # [cite: 251]
        plt.tight_layout(pad=0) # [cite: 252]
        return fig # [cite: 252]
    except Exception as e:
        st.warning(f"Could not generate word cloud: {e}") # [cite: 252]
        return None


# Utility function to create an entity network
def create_entity_network(entities, doc): # [cite: 252]
    """Create a Plotly network graph of named entities and their co-occurrences.""" # [cite: 252]
    if not entities or doc is None: return None # [cite: 252]
    try:
        G = nx.Graph() # [cite: 252]
        unique_entities_map = {ent_text: label for ent_text, label in list(set(entities))} # [cite: 253]
        if not unique_entities_map: return None # [cite: 253]

        for ent, label in unique_entities_map.items(): G.add_node(ent, label=label, size=10) # [cite: 253]

        # Add edges based on co-occurrence in the same sentence, weighted by frequency
        if hasattr(doc, 'sents') and doc.sents is not None: # [cite: 254]
            try:
                for sent in doc.sents: # [cite: 254]
                     if isinstance(sent, spacy.tokens.span.Span): # [cite: 254]
                        sent_ents_texts = list(set([ent.text for ent in sent.ents if isinstance(ent, spacy.tokens.span.Span) and ent.text in G.nodes()])) # [cite: 255]
                        for i in range(len(sent_ents_texts)): # [cite: 255]
                            for j in range(i + 1, len(sent_ents_texts)): # [cite: 255]
                                entity1, entity2 = sent_ents_texts[i], sent_ents_texts[j] # [cite: 256]
                                if G.has_node(entity1) and G.has_node(entity2) and entity1 != entity2: # [cite: 257]
                                    if G.has_edge(entity1, entity2): # [cite: 257]
                                        G[entity1][entity2]['weight'] += 1 # [cite: 258]
                                    else:
                                        G.add_edge(entity1, entity2, weight=1) # [cite: 259]
            except Exception as e:
                 st.warning(f"Error processing sentences for entity network: {e}") # [cite: 259]

        if len(G.nodes) == 0 or len(G.edges) == 0: return None # [cite: 259]

        # Prune graph: remove nodes with degree 0 (isolates) after edge creation
        isolates = list(nx.isolates(G))
        G.remove_nodes_from(isolates)
        if len(G.nodes) < 2 or len(G.edges) == 0 : return None


        pos = nx.spring_layout(G, k=0.85, iterations=70, seed=42, weight='weight') # Adjusted for weights # [cite: 259]
        edge_x, edge_y = [], [] # [cite: 260]
        edge_weights = []
        for edge in G.edges(data=True): # [cite: 260]
            x0, y0 = pos[edge[0]]; x1, y1 = pos[edge[1]] # [cite: 260]
            edge_x.extend([x0, x1, None]); edge_y.extend([y0, y1, None]) # [cite: 260]
            edge_weights.append(edge[2].get('weight', 1))


        edge_trace = go.Scatter(x=edge_x, y=edge_y, line=dict(width=0.7, color='#888'), hoverinfo='none', mode='lines') # [cite: 261]
        node_x, node_y, node_text, node_sizes, node_colors_list = [], [], [], [], [] # [cite: 261]

        label_colors = { # [cite: 262]
            'PERSON': '#FF6347', 'ORG': '#4682B4', 'GPE': '#32CD32', 'LOC': '#FFD700', # [cite: 262]
            'DATE': '#6A5ACD', 'CARDINAL': '#FF69B4', 'NORP': '#00CED1', 'FAC': '#F4A460', # [cite: 263]
            'PRODUCT': '#98FB98', 'EVENT': '#87CEEB', 'LAW': '#DC143C', 'LANGUAGE': '#AFEEEE', # [cite: 263]
            'MONEY': '#ADFF2F', 'QUANTITY': '#FFB6C1', 'ORDINAL': '#E6E6FA', 'TIME': '#DDA0DD',
            'WORK_OF_ART': '#BDB76B', 'PERCENT': '#FFA07A',
            'OTHER': '#A9A9A9' # [cite: 264]
        }
        # Calculate node sizes based on degree (number of connections)
        degrees = dict(G.degree())
        min_size, max_size = 10, 30

        for node in G.nodes(): # [cite: 264]
            x, y = pos[node]; node_x.append(x); node_y.append(y) # [cite: 265]
            label = G.nodes[node].get('label', 'OTHER') # [cite: 265]
            degree = degrees.get(node, 1)
            size = min_size + (max_size - min_size) * (degree / max(1,max(degrees.values()))) # Normalize size

            node_text.append(f"{node} ({label}) - Connections: {degree}") # [cite: 265]
            node_colors_list.append(label_colors.get(label, '#A9A9A9')) # [cite: 265]
            node_sizes.append(size)


        node_trace = go.Scatter( # [cite: 265]
            x=node_x, y=node_y, mode='markers+text', textfont=dict(size=10, color='black'), # [cite: 266]
            texttemplate='%{customdata[0]}', customdata=[[t.split('(')[0].strip()] for t in node_text],
            textposition="top center", hoverinfo='text', hovertext=node_text, # [cite: 266]
            marker=dict(showscale=False, color=node_colors_list, size=node_sizes, line_width=1.5, line_color='black') # [cite: 267]
        )

        fig = go.Figure(data=[edge_trace, node_trace], # [cite: 267]
                        layout=go.Layout( # [cite: 267]
                            showlegend=False, hovermode='closest', margin=dict(b=10, l=5, r=5, t=30), # [cite: 268]
                            xaxis=dict(showgrid=False, zeroline=False, showticklabels=False), # [cite: 268]
                            yaxis=dict(showgrid=False, zeroline=False, showticklabels=False), # [cite: 269]
                            title=dict(text="Entity Co-occurrence Network", x=0.5, font=dict(size=16, color='white')), # [cite: 269]
                            plot_bgcolor='#004d40', paper_bgcolor='#004d40' # [cite: 270]
                        ))
        return fig # [cite: 270]
    except Exception as e:
        st.warning(f"Could not generate entity network: {e}") # [cite: 270]
        return None


# Utility function to extract audio from video
def extract_audio_from_video(video_path: str, audio_output_path: str) -> bool: # [cite: 270]
    """Extract audio from a video file and save it as a WAV file.""" # [cite: 270]
    try:
        clip = VideoFileClip(video_path) # [cite: 271]
        clip.audio.write_audiofile(audio_output_path, codec='pcm_s16le', fps=16000) # [cite: 271]
        clip.close() # [cite: 271]
        return True # [cite: 271]
    except Exception as e:
        st.warning(f"Audio extraction failed: {e}") # [cite: 271]
        return False


# Utility function to transcribe audio
def transcribe_audio_with_whisper(audio_path: str, transcription_model: Any) -> str: # [cite: 271]
    """Transcribe audio file using Whisper model.""" # [cite: 271]
    if transcription_model is None: # [cite: 271]
        st.warning("Transcription model not loaded.") # [cite: 272]
        return ""
    try:
        result = transcription_model.transcribe(audio_path, fp16=False) # Disable fp16 if causing issues # [cite: 272]
        return result["text"] # [cite: 272]
    except Exception as e:
        st.warning(f"Transcription failed: {e}") # [cite: 272]
        return ""
    finally:
        if 'audio_path' in locals() and os.path.exists(audio_path): # [cite: 272]
            try: os.unlink(audio_path) # [cite: 273]
            except Exception as e: st.warning(f"Could not delete temporary audio file {audio_path}: {e}") # [cite: 273]


# Utility function to fetch and parse URL
def fetch_and_parse_url(url: str) -> str: # [cite: 273]
    """Fetches the content from a URL and extracts visible text.""" # [cite: 273]
    try:
        headers = { # [cite: 273]
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36 AdvancedAIAnalyzer/1.0" # [cite: 274]
        }
        response = requests.get(url, headers=headers, timeout=20) # Increased timeout # [cite: 275]
        response.raise_for_status() # [cite: 275]
        soup = BeautifulSoup(response.content, "html.parser") # [cite: 275]
        for tag in soup(["script", "style", "header", "footer", "nav", "aside", "form", "noscript", "img", "meta", "link", "figure", "iframe"]): # Added more tags # [cite: 276]
            tag.decompose() # [cite: 276]
        text_chunks = [chunk.get_text(" ", strip=True) for chunk in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'article', 'div'])]
        text = "\n".join(chunk for chunk in text_chunks if chunk and len(chunk.split()) > 5) # Join with newline, filter short chunks
        text = re.sub(r'\s*\n\s*', '\n', text).strip() # Normalize multiple newlines
        text = re.sub(r'\s+', ' ', text).strip() # Normalize whitespace # [cite: 276]
        return text
    except requests.exceptions.Timeout:
        st.error(f"Timeout error fetching URL: {url}. The server did not respond in time.")
        return ""
    except requests.exceptions.HTTPError as e:
        st.error(f"HTTP error fetching URL: {e}. Status code: {e.response.status_code}")
        return ""
    except requests.exceptions.RequestException as e: # [cite: 277]
        st.error(f"Error fetching URL: {e}") # [cite: 277]
        return ""
    except Exception as e:
        st.error(f"Failed to parse URL content: {e}") # [cite: 277]
        return ""

# Main application function
def main_application(): # [cite: 277]
    """Main application function with enhanced features integrated.""" # [cite: 277]
    MODELS = load_models() # [cite: 277]

    # Sidebar Navigation
    try: # [cite: 278]
        # Replace with your actual logo path or remove if not needed.
        # sidebar_logo_path = 'QR3XoLs.jpeg' # Path to logo [cite: 279, 281]
        # if os.path.exists(sidebar_logo_path): # Check if logo exists [cite: 281]
        #     st.sidebar.image(Image.open(sidebar_logo_path), width=100) # Display logo [cite: 281]
        # else:
        #     st.sidebar.warning("Logo image 'QR3XoLs.jpeg' not found. Using default.")
        # For this example, let's assume the image might not be found or we use a placeholder emoji
        st.sidebar.markdown("<h1 style='text-align: center; color: black;'>🧠</h1>", unsafe_allow_html=True) # Using emoji as placeholder
    except Exception as img_e: # [cite: 282]
        st.sidebar.warning(f"Could not load logo: {img_e}") # [cite: 282]
        # --- Sidebar Navigation ---
    st.sidebar.image(Image.open('QR3XoLs.jpeg'), width=100)  # Placeholder logo, width=100
    st.sidebar.caption("pokeman")


    st.sidebar.title("AI Analyzer Suite") # [cite: 282]
    st.sidebar.markdown("---") # [cite: 283]
    app_mode = st.sidebar.radio( # [cite: 283]
        "Navigation:", # [cite: 283]
        ["🏠 Home", "💬 Text Analysis", "📄 Document Analysis", "🎤 Media Analysis", "🌐 Web Analysis", "ℹ️ About"], # [cite: 283]
        help="Select the type of content you want to analyze" # [cite: 283]
    )
    st.sidebar.markdown("---") # [cite: 283]
    st.sidebar.info("Upload content, paste text, or enter a URL to generate advanced analysis, summaries, and study materials.") # [cite: 283]

    if "current_analysis_content" not in st.session_state:
        st.session_state.current_analysis_content = None
    if "current_source_name" not in st.session_state:
        st.session_state.current_source_name = "None"


    if app_mode == "🏠 Home": # [cite: 284]
        st.title("🧠 Advanced AI Content Analysis Suite") # [cite: 284]
        st.markdown( # [cite: 284]
            """
            <div class="content-card">
                <h3>📚 Comprehensive Content Understanding Platform</h3>
                <p>🔍 This advanced tool provides deep analysis of text, documents, media, and web content using state-of-the-art NLP techniques.</p>
                <h4>✨ Key Features Include:</h4>
                <ul>
                    <li>🏷️ Semantic keyword extraction with contextual meanings & scores</li>
                    <li>📝 Multiple note-taking formats (Cornell, Mind Map Outline, Linear)</li>
                    <li>💡 Abstractive study materials generation (Flashcards, Questions, Diagrams)</li>
                    <li>📄 Professional PDF report generation for all analyses</li>
                    <li>📊 Enhanced visualizations: Word Clouds & Interactive Entity Networks</li>
                    <li>🗣️ Audio playback for generated summaries (Text-to-Speech)</li>
                    <li>🌐 Full analysis pipeline for web URL content</li>
                </ul>
                <h4>❓ How to Use:</h4>
                <ol>
                    <li>🖱️ Select an analysis mode from the sidebar (Text, Document, Media, Web).</li>
                    <li>⬆️ Upload your file, paste text, or enter a URL.</li>
                    <li>⚙️ Click the respective 'Analyze' or 'Generate' buttons within the selected mode.</li>
                    <li>🔎 Explore different analysis tabs (Summaries, Keywords, Study Aids, Visualizations).</li>
                    <li>⬇️ Download comprehensive reports or specific outputs (PDF, CSV).</li>
                </ol>
                <p><em>Navigate using the sidebar to begin your analysis.</em></p>
            </div>
            """, unsafe_allow_html=True # [cite: 289]
        )

    elif app_mode == "💬 Text Analysis": # [cite: 289]
        st.header("📝 Direct Text Input Analysis") # [cite: 289]
        text_input = st.text_area("Paste your text here for analysis:", height=250, # [cite: 290]
                                 placeholder="Enter any text content you wish to analyze in detail...", # [cite: 290]
                                 help="Paste text for summarization, keyword extraction, study aid generation, and visualizations.") # [cite: 290]
        if st.button("Analyze Pasted Text", key="analyze_direct_text"): # [cite: 290]
            if text_input and text_input.strip(): # [cite: 291]
                st.session_state.current_analysis_content = text_input # [cite: 291]
                st.session_state.current_source_name = "Pasted_Text_Input" # [cite: 291]
                # Clear previous analysis results for this specific source to avoid stale data
                if f'analysis_results_{st.session_state.current_source_name}' in st.session_state:
                    del st.session_state[f'analysis_results_{st.session_state.current_source_name}']
                if f'keywords_{st.session_state.current_source_name}' in st.session_state:
                    del st.session_state[f'keywords_{st.session_state.current_source_name}']
                if f'study_materials_{st.session_state.current_source_name}' in st.session_state:
                    del st.session_state[f'study_materials_{st.session_state.current_source_name}']
            else:
                st.warning("Please enter some text to analyze.") # [cite: 291]
                st.session_state.current_analysis_content = None

        if st.session_state.current_analysis_content and st.session_state.current_source_name == "Pasted_Text_Input":
            display_advanced_analysis_results(st.session_state.current_analysis_content, MODELS, st.session_state.current_source_name)

    elif app_mode == "📄 Document Analysis": # [cite: 291]
        st.header("📂 Document File Analysis") # [cite: 291]
        uploaded_file = st.file_uploader( # [cite: 291]
            "Upload document (PDF, DOCX, PPTX, TXT):", # [cite: 292]
            type=["pdf", "docx", "pptx", "txt"], # [cite: 292]
            help="Supported formats: PDF, Word (DOCX), PowerPoint (PPTX), Plain Text (TXT)" # [cite: 292]
        )
        if uploaded_file: # [cite: 292]
            # Generate a unique key for the button based on file name to handle re-uploads correctly
            analyze_doc_button_key = f"analyze_doc_{uploaded_file.name}"
            if st.button(f"Analyze Document: {uploaded_file.name}", key=analyze_doc_button_key):
                with st.spinner(f"Processing {uploaded_file.name}... This might take a few moments."): # [cite: 292]
                    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file: # [cite: 293]
                         tmp_file.write(uploaded_file.getvalue()) # [cite: 293]
                         tmp_file_path = tmp_file.name # [cite: 293]
                    raw_text = extract_text_from_file(tmp_file_path, uploaded_file.type) # [cite: 293]

                    if raw_text and raw_text.strip(): # [cite: 294]
                        st.session_state.current_analysis_content = raw_text # [cite: 294]
                        st.session_state.current_source_name = uploaded_file.name # [cite: 295]
                        if f'analysis_results_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'analysis_results_{st.session_state.current_source_name}']
                        if f'keywords_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'keywords_{st.session_state.current_source_name}']
                        if f'study_materials_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'study_materials_{st.session_state.current_source_name}']

                        with st.expander("View Extracted Text from Document", expanded=False): # [cite: 294]
                            st.text_area("Extracted Text Preview:", value=raw_text[:5000]+"...", height=200, disabled=True) # [cite: 294]
                    else:
                        st.error("Could not extract meaningful text from the document. The document might be empty, image-based (requiring OCR not yet implemented), or corrupted.") # [cite: 295]
                        st.session_state.current_analysis_content = None
        # Display results if content is loaded for the current document source
        if st.session_state.current_analysis_content and uploaded_file and st.session_state.current_source_name == uploaded_file.name:
             display_advanced_analysis_results(st.session_state.current_analysis_content, MODELS, st.session_state.current_source_name)


    elif app_mode == "🎤 Media Analysis": # [cite: 295]
        st.header("🎧 Audio/Video Transcription & Analysis") # [cite: 295]
        uploaded_media = st.file_uploader( # [cite: 295]
            "Upload media file (MP3, WAV, MP4, MOV, AVI, etc.):", # [cite: 296]
            type=["mp3", "wav", "mp4", "m4a", "ogg", "flac", "mov", "avi", "mpeg"], # Added mpeg # [cite: 296]
            help="Supported audio & video formats for transcription." # [cite: 296]
        )
        if uploaded_media: # [cite: 296]
            analyze_media_button_key = f"analyze_media_{uploaded_media.name}"
            if st.button(f"Transcribe & Analyze Media: {uploaded_media.name}", key=analyze_media_button_key):
                with st.spinner(f"Processing and transcribing {uploaded_media.name}... This can take some time depending on the file size."): # [cite: 296]
                    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_media.name)[1]) as tmp_media: # [cite: 297]
                        tmp_media.write(uploaded_media.getvalue()) # [cite: 297]
                        tmp_media_path = tmp_media.name # [cite: 297]

                    transcribed_text = None # [cite: 297]
                    audio_path_for_transcription = tmp_media_path # [cite: 298]
                    original_media_temp_path_to_delete = None # Keep track if original video is converted

                    if uploaded_media.name.lower().endswith(('.mp4', '.mov', '.avi', '.mpeg')): # [cite: 298]
                        audio_output_filename = f"extracted_audio_{os.path.basename(tmp_media_path)}.wav" # [cite: 298]
                        audio_output_path = os.path.join(tempfile.gettempdir(), audio_output_filename) # [cite: 298]
                        if extract_audio_from_video(tmp_media_path, audio_output_path): # [cite: 298]
                            audio_path_for_transcription = audio_output_path # [cite: 299]
                            original_media_temp_path_to_delete = tmp_media_path # Mark original video temp for deletion
                        else: # [cite: 299]
                            st.error("Failed to extract audio from video file.")
                            audio_path_for_transcription = None # Ensure no transcription if audio extraction fails
                    # Cleanup original video temp file if audio was extracted
                    if original_media_temp_path_to_delete and os.path.exists(original_media_temp_path_to_delete):
                        try: os.unlink(original_media_temp_path_to_delete) # [cite: 300]
                        except Exception as e_del: st.warning(f"Could not delete temporary video file {original_media_temp_path_to_delete}: {e_del}") # [cite: 300]


                    if audio_path_for_transcription and os.path.exists(audio_path_for_transcription): # [cite: 301]
                         transcribed_text = transcribe_audio_with_whisper(audio_path_for_transcription, MODELS.get('transcription_model')) # [cite: 301]
                         # transcribe_audio_with_whisper now handles deletion of its input audio_path_for_transcription
                    elif not uploaded_media.name.lower().endswith(('.mp4', '.mov', '.avi', '.mpeg')) and tmp_media_path and os.path.exists(tmp_media_path): # Original was audio # [cite: 301]
                        # This case means it was an audio file to begin with, and tmp_media_path is the one to transcribe and delete
                        transcribed_text = transcribe_audio_with_whisper(tmp_media_path, MODELS.get('transcription_model')) # [cite: 302]
                         # transcribe_audio_with_whisper will delete tmp_media_path
                    else:
                         if not audio_path_for_transcription : st.error("Could not find a valid audio file to transcribe after processing media.") # [cite: 302]


                    if transcribed_text and transcribed_text.strip(): # [cite: 303]
                        st.session_state.current_analysis_content = transcribed_text # [cite: 303]
                        st.session_state.current_source_name = f"Transcript_{uploaded_media.name}" # [cite: 303]
                        if f'analysis_results_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'analysis_results_{st.session_state.current_source_name}']
                        if f'keywords_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'keywords_{st.session_state.current_source_name}']
                        if f'study_materials_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'study_materials_{st.session_state.current_source_name}']
                        with st.expander("View Full Transcription", expanded=False): # [cite: 303]
                            st.text_area("Transcribed Text:", value=transcribed_text, height=200, disabled=True) # [cite: 303]
                    else:
                        st.error("Could not transcribe meaningful text from the media file. The media might be silent or the transcription failed.") # [cite: 304]
                        st.session_state.current_analysis_content = None
        # Display results if content is loaded for the current media source
        if st.session_state.current_analysis_content and uploaded_media and st.session_state.current_source_name == f"Transcript_{uploaded_media.name}":
            display_advanced_analysis_results(st.session_state.current_analysis_content, MODELS, st.session_state.current_source_name)


    elif app_mode == "🌐 Web Analysis": # [cite: 304]
        st.header("🔗 Web URL Content Analysis") # [cite: 304]
        url_input = st.text_input("Enter URL to analyze:", placeholder="https://example.com/article", key="url_input_field") # [cite: 304]
        if st.button("Fetch & Analyze URL Content", key="analyze_url_content"): # [cite: 304]
            if url_input and url_input.strip(): # [cite: 304]
                with st.spinner("Fetching and analyzing URL content... This may take a moment."): # [cite: 305]
                    fetched_text = fetch_and_parse_url(url_input) # [cite: 305]
                    if fetched_text and fetched_text.strip(): # [cite: 305]
                        try: # [cite: 306]
                            parsed_url = requests.utils.urlparse(url_input) # [cite: 306]
                            source_name_url = parsed_url.netloc.replace("www.","") + (parsed_url.path.replace("/","_")[:30] if parsed_url.path else "") or "Web_Content" # [cite: 306]
                        except Exception: source_name_url = "Web_Content_Unknown_URL" # [cite: 306]

                        st.session_state.current_analysis_content = fetched_text # [cite: 307]
                        st.session_state.current_source_name = source_name_url # [cite: 307]
                        if f'analysis_results_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'analysis_results_{st.session_state.current_source_name}']
                        if f'keywords_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'keywords_{st.session_state.current_source_name}']
                        if f'study_materials_{st.session_state.current_source_name}' in st.session_state: del st.session_state[f'study_materials_{st.session_state.current_source_name}']

                        with st.expander("View Fetched Web Content Preview", expanded=False): # [cite: 307]
                            st.text_area("Fetched Content:", value=fetched_text[:5000]+"...", height=200, disabled=True) # [cite: 307]
                    else:
                        st.error("Could not fetch meaningful content from the URL. Please ensure the URL is valid, publicly accessible, and contains substantial text content.") # [cite: 308]
                        st.session_state.current_analysis_content = None
            else:
                st.warning("Please enter a valid URL to analyze.") # [cite: 309]
        # Display results if content is loaded for the current URL source
        if st.session_state.current_analysis_content and url_input and st.session_state.current_source_name != "Pasted_Text_Input" and st.session_state.current_source_name != "None" and not st.session_state.current_source_name.startswith("Transcript_") and not any(st.session_state.current_source_name.lower().endswith(ext) for ext in ['.pdf','.docx','.pptx','.txt']): # Heuristic to check if it's likely a URL source
            display_advanced_analysis_results(st.session_state.current_analysis_content, MODELS, st.session_state.current_source_name)


    elif app_mode == "ℹ️ About": # [cite: 309]
        st.title("ℹ️ About the Advanced AI Analyzer") # [cite: 310]
        st.markdown("""
            <div class="content-card">
                <h3>🤖 Advanced AI Content Analysis Suite</h3>
                <p>✨ Version 5.0 (Enhanced) - Your comprehensive partner for deep content understanding.</p>
                <h4>⚙️ Core Technologies Utilized:</h4>
                <ul>
                    <li>🧠 Transformers (T5 for summarization, BERT for embeddings/keywords) from Hugging Face</li>
                    <li>📝 spaCy for advanced Natural Language Processing (NLP) tasks like entity recognition, noun chunking, and sentence segmentation</li>
                    <li>🎙️ OpenAI Whisper for accurate audio transcription (video and audio files)</li>
                    <li>📄 ReportLab for generating professional, structured PDF reports</li>
                    <li>📊 Plotly & Matplotlib for creating interactive and static data visualizations</li>
                    <li>🔊 gTTS (Google Text-to-Speech) for audio playback of summaries</li>
                    <li>🌐 BeautifulSoup & Requests for fetching and parsing web content</li>
                </ul>
                <h4>🔑 Key Features at a Glance:</h4>
                <ul>
                    <li>📦 Multi-format Content Analysis: Handles direct text input, documents (PDF, DOCX, PPTX, TXT), media files (MP3, WAV, MP4, MOV, etc.), and web URLs.</li>
                    <li>👥 Audience-Specific Summarization: Tailors summaries for different understanding levels (student, researcher, expert, child).</li>
                    <li>🎯 Contextual Keyword Extraction: Identifies key terms with semantic meanings, categories, and contextual relevance scores.</li>
                    <li>🗒️ Diverse Note-Taking Aids: Generates Cornell notes, mind map outlines, and linear notes.</li>
                    <li>💡 Abstractive Study Materials: Creates flashcards, conceptual questions, diagram suggestions, and a structured study plan.</li>
                    <li>🕸️ Interactive Visualizations: Offers dynamic entity co-occurrence networks and insightful word clouds.</li>
                    <li>📄 Comprehensive PDF Reports: Allows download of detailed analysis and study materials in a portable format.</li>
                    <li>🗣️ Smart Audio Assistance: Provides text-to-speech functionality for listening to generated summaries.</li>
                </ul>
                <p>🧑‍🔬 This tool is designed to empower students, researchers, educators, content creators, and professionals by providing powerful AI-driven insights from various forms of content, making learning and analysis more efficient and profound.</p>
                <p><em>Developed with ❤️ and a lot of Python!</em></p>
            </div>
        """, unsafe_allow_html=True) # [cite: 316]

if __name__ == "__main__":
    main_application() # [cite: 316]
### CODE ENDS HERE ###